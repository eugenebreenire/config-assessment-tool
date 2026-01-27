# compare_tool/powerpoint.py

import logging
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def generate_powerpoint_from_analysis(comparison_result_path, powerpoint_output_path, current_file_path, previous_file_path):
    logging.debug("Generating PowerPoint presentation...")

    try:
        # prs = Presentation(template_path)  # Open the template

        # Define the relative path for the template using the TEMPLATE_FOLDER
        template_folder = config.get('TEMPLATE_FOLDER', 'templates')  # 'templates' is the default folder name
        template_path = os.path.join(template_folder, 'template.pptx')

        # Load the 'Analysis' sheet from the current workbook (uploaded by the user)
        df_current_analysis = pd.read_excel(current_file_path, sheet_name='Analysis')

        # Count the number of valid applications by counting rows where 'name' column is not NaN or empty
        number_of_apps = df_current_analysis['name'].dropna().str.strip().ne('').sum()

        # Log the number of valid applications
        logging.info(f"Number of applications in the current 'Analysis' sheet: {number_of_apps}")

        # Check if the template exists, otherwise, ask the user for input or use environment variables
        if not os.path.exists(template_path):
            template_path = os.getenv('TEMPLATE_PATH', template_path)  # Allow user to set this via an environment variable
            if not os.path.exists(template_path):
                template_path = input("Template not found! Please provide the full path to the template: ")

        # Load the template
        prs = Presentation(template_path)
        logging.debug(f"Template loaded from: {template_path}")

        # Load Summary sheets (current and previous) to drive Key Callouts on Slide 2.
        current_summary_df = pd.read_excel(current_file_path, sheet_name='Summary')
        previous_summary_df = pd.read_excel(previous_file_path, sheet_name='Summary')

        # Load the Summary sheet
        summary_df = pd.read_excel(comparison_result_path, sheet_name='Summary')
        logging.debug("Loaded Summary sheet successfully.")
        logging.debug(f"Summary DataFrame head:\n{summary_df.head()}")

        # Load Summary sheets to drive the Key Callouts slide.
        current_summary_df = pd.read_excel(current_file_path, sheet_name='Summary')
        previous_summary_df = pd.read_excel(previous_file_path, sheet_name='Summary')

        # Load the Analysis sheet
        df_analysis = pd.read_excel(comparison_result_path, sheet_name='Analysis')
        # Load the 'AppAgentsAPM' sheet from the Excel file
        df_app_agents = pd.read_excel(comparison_result_path, sheet_name='AppAgentsAPM')
        # Load the 'MachineAgentsAPM' sheet from the Excel file
        df_machine_agents = pd.read_excel(comparison_result_path, sheet_name='MachineAgentsAPM')
        # Load the 'BusinessTransactionsAPM' sheet from the Excel file
        df_BTs = pd.read_excel(comparison_result_path, sheet_name='BusinessTransactionsAPM')
        # Load the 'BackendsAPM' sheet from the Excel file
        df_Backends = pd.read_excel(comparison_result_path, sheet_name='BackendsAPM')
        # Load the 'OverheadAPM' sheet from the Excel file
        df_Overhead = pd.read_excel(comparison_result_path, sheet_name='OverheadAPM')
        # Load the 'ServiceEndpointsAPM' sheet from the Excel file
        df_ServiceEndpoints = pd.read_excel(comparison_result_path, sheet_name='ServiceEndpointsAPM')
        # Load the 'ErrorConfigurationAPM' sheet from the Excel file
        df_ErrorConfiguration = pd.read_excel(comparison_result_path, sheet_name='ErrorConfigurationAPM')
        # Load the 'HealthRulesAndAlertingAPM' sheet from the Excel file
        df_HealthRulesAndAlerting = pd.read_excel(comparison_result_path, sheet_name='HealthRulesAndAlertingAPM')
        # Load the 'DataCollectorsAPM' sheet from the Excel file
        df_DataCollectors = pd.read_excel(comparison_result_path, sheet_name='DataCollectorsAPM')
        # Load the 'DashboardsAPM' sheet from the Excel file
        df_Dashboards = pd.read_excel(comparison_result_path, sheet_name='DashboardsAPM')

        # Function to find table placeholders by name
        def find_table_placeholder_by_name(slide, name):
            for shape in slide.shapes:
                if shape.is_placeholder and shape.name == name:
                    return shape
            return None  # Return None if not found

        def insert_table_at_placeholder(slide, placeholder_name, rows, cols):
            """Insert a table at the position of a placeholder."""
            placeholder = find_table_placeholder_by_name(slide, placeholder_name)
            
            if not placeholder:
                logging.error(f"Placeholder '{placeholder_name}' not found on the slide.")
                return None

            # Get placeholder dimensions
            left = placeholder.left
            top = placeholder.top
            width = placeholder.width
            height = placeholder.height

            logging.debug(f"Inserting table at placeholder position: left={left}, top={top}, width={width}, height={height}")

            # Insert table at the placeholder's position
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            return table_shape.table  # Return the inserted table

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Slide 2 (index 1) — Assessment Result - Key Callouts **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 

        slide = prs.slides[1]  # Slide 2

        # Existing local helpers to parse percentages and arrows.
        def _parse_percent_to_float(val):
            if pd.isna(val):
                return None
            if isinstance(val, (int, float)):
                return float(val)
            s = str(val).strip()
            if '→' in s:
                s = s.split('→')[-1].strip()
            s = s.replace('%', '')
            try:
                return float(s)
            except Exception:
                return None

        def _get_tier_percent(df, tier):
            name_map = {c.lower(): c for c in df.columns}
            candidates = [f"{tier.lower()} %", f"{tier.lower()}%", f"percentage{tier.lower()}", f"{tier.lower()}percentage"]
            for cand in candidates:
                if cand in name_map:
                    return _parse_percent_to_float(df[name_map[cand]].iloc[0])
            needed = ['bronze', 'silver', 'gold', 'platinum']
            if all(k in name_map for k in needed):
                try:
                    total = 0.0
                    counts = {}
                    for k in needed:
                        val = pd.to_numeric(df[name_map[k]].iloc[0], errors='coerce')
                        counts[k] = 0.0 if pd.isna(val) else float(val)
                        total += counts[k]
                    if total > 0:
                        return (counts[tier.lower()] / total) * 100.0
                except Exception:
                    return None
            return None

        def _arrow(curr, prev):
            if curr is None or prev is None:
                return '→'
            if curr > prev:
                return '↑'
            if curr < prev:
                return '↓'
            return '→'

        def _trend_word(curr, prev):
            if curr is None or prev is None:
                return "held steady"
            if curr > prev:
                return "increased"
            if curr < prev:
                return "decreased"
            return "held steady"

        def _delta_pp(curr, prev):
            if curr is None or prev is None:
                return None
            return round(curr - prev, 1)

        # NEW: Coverage helpers for APM adoption (local grade_token here; module-level functions consume it).
        def _grade_token(s: str):
            if not s:
                return None
            m = re.search(r'(platinum|gold|silver|bronze)', str(s), re.I)
            return m.group(1).lower() if m else None

        def _apps_coverage(path):
            try:
                df = pd.read_excel(path, sheet_name='Analysis')
                total = int(df['name'].dropna().astype(str).str.strip().ne('').sum())
                if total == 0:
                    return (0, 0, 0.0)
                rated = int(df['OverallAssessment'].apply(_grade_token).notna().sum())
                pct = (rated / total) * 100.0
                return (total, rated, pct)
            except Exception:
                return (0, 0, None)

        def _arrow_threshold(curr, prev, threshold_pp=0.5):
            if curr is None or prev is None:
                return '→'
            delta = curr - prev
            if delta >= threshold_pp:
                return '↑'
            if delta <= -threshold_pp:
                return '↓'
            return '→'

        def _fmt_pp_delta(prev, curr):
            if prev is None or curr is None:
                return None
            d = curr - prev
            sign = '+' if d > 0 else '−' if d < 0 else '±'
            return f"{prev:.1f}%→{curr:.1f}% ({sign}{abs(d):.1f} pp)."

        # Derive Gold/Platinum percentages from Summary sheets.
        curr_gold = _get_tier_percent(current_summary_df, 'Gold')
        prev_gold = _get_tier_percent(previous_summary_df, 'Gold')
        curr_plat = _get_tier_percent(current_summary_df, 'Platinum')
        prev_plat = _get_tier_percent(previous_summary_df, 'Platinum')

        # NEW: Coverage from the original APM workbooks (for the first row).
        total_prev, rated_prev, cov_prev = _apps_coverage(previous_file_path)
        total_curr, rated_curr, cov_curr = _apps_coverage(current_file_path)
        cov_arrow = _arrow_threshold(cov_curr, cov_prev)
        cov_outcome = (
            f"Coverage: {cov_curr:.1f}% of apps rated ({rated_curr}/{total_curr})."
            if cov_curr is not None and total_curr > 0
            else "Coverage data not available."
        )
        cov_prev_curr = _fmt_pp_delta(cov_prev, cov_curr)
        if cov_prev_curr:
            cov_outcome = f"{cov_outcome} {cov_prev_curr}"

        # Overall result and next focus (unchanged from your prior logic).
        try:
            df_cmp = pd.read_excel(comparison_result_path, sheet_name='Analysis')
        except Exception:
            df_cmp = None

        def _count_changes(df, col):
            if df is None or col not in df.columns:
                return 0, 0
            s = df[col].astype(str)
            up = s.str.contains('Upgraded', case=False, na=False).sum()
            down = s.str.contains('Downgraded', case=False, na=False).sum()
            return int(up), int(down)

        up_overall, down_overall = _count_changes(df_cmp, 'OverallAssessment')
        overall_result_text = "Increase" if up_overall > down_overall else "Decrease" if down_overall > up_overall else "Even"

        area_cols = [
            'AppAgentsAPM','MachineAgentsAPM','BusinessTransactionsAPM',
            'BackendsAPM','OverheadAPM','ServiceEndpointsAPM',
            'ErrorConfigurationAPM','HealthRulesAndAlertingAPM',
            'DataCollectorsAPM','DashboardsAPM'
        ]
        pretty = {
            'AppAgentsAPM': 'App Agents',
            'MachineAgentsAPM': 'Machine Agents',
            'BusinessTransactionsAPM': 'Business Transactions',
            'BackendsAPM': 'Backends',
            'OverheadAPM': 'Overhead',
            'ServiceEndpointsAPM': 'Service Endpoints',
            'ErrorConfigurationAPM': 'Error Configuration',
            'HealthRulesAndAlertingAPM': 'Health Rules & Alerting',
            'DataCollectorsAPM': 'Data Collectors',
            'DashboardsAPM': 'Dashboards',
        }
        downgraded_counts = []
        if df_cmp is not None:
            for col in area_cols:
                if col in df_cmp.columns:
                    s = df_cmp[col].astype(str)
                    cnt = s.str.contains('Downgraded', case=False, na=False).sum()
                    downgraded_counts.append((col, int(cnt)))
        downgraded_counts.sort(key=lambda x: x[1], reverse=True)
        focus_list = [pretty[c] for c, n in downgraded_counts if n > 0][:2]
        next_focus_text = ", ".join(focus_list) if focus_list else "Maintain current progress"

        delta_gold = _delta_pp(curr_gold, prev_gold)
        delta_plat = _delta_pp(curr_plat, prev_plat)

        def _fmt_outcome(prev, curr, delta):
            if prev is None or curr is None:
                return "Data not available."
            sign = "+" if delta is not None and delta > 0 else "−" if delta is not None and delta < 0 else "±"
            if delta is None:
                return f"{prev:.1f}%→{curr:.1f}%"
            return f"{prev:.1f}%→{curr:.1f}% ({sign}{abs(delta):.1f} pp)."

        # Table headers and rows.
        headers = [
            "AppD Maturity Progression & Engagement",
            "Commentary",
            "Outcomes",
            "Change/Status Since Last",
        ]

        rows = [
            # UPDATED first row: APM-only coverage.
            [
                "B/S/G/P Model Adoption & Maturity Status",
                f"B/S/G/P model applied to APM; assessment covered {int(total_curr)} apps.",
                cov_outcome,
                cov_arrow,
            ],
            [
                "Gold Status Apps",
                f"Gold-or-better coverage {_trend_word(curr_gold, prev_gold)} across the portfolio.",
                _fmt_outcome(prev_gold, curr_gold, delta_gold),
                _arrow(curr_gold, prev_gold),
            ],
            [
                "Platinum Status Apps",
                f"Platinum presence {_trend_word(curr_plat, prev_plat)}; teams progressing on prerequisites.",
                _fmt_outcome(prev_plat, curr_plat, delta_plat),
                _arrow(curr_plat, prev_plat),
            ],
            [
                "Maturity Partnership",
                "Working cadence in place; recommendations implemented during this period.",
                f"Overall result: {overall_result_text}. Next focus: {next_focus_text}.",
                "↑" if overall_result_text == "Increase" else "↓" if overall_result_text == "Decrease" else "→",
            ],
        ]

        # Insert table, set header styles, autosize arrow column, and render rows.
        key_callouts_ph = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if key_callouts_ph:
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(rows) + 1, len(headers))
        else:
            table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.6), Inches(2.1), Inches(9.0), Inches(4.0)).table

        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(12)
            cell.text_frame.word_wrap = False

        autosize_col_to_header(table, 3, header_pt=12, padding_in=0.6, avg_char_em=0.55)

        for r_idx, row in enumerate(rows, start=1):
            for c_idx, value in enumerate(row):
                cell = table.cell(r_idx, c_idx)
                if c_idx == 3 and value in ("↑", "↓", "→"):
                    set_arrow_cell(cell, value, color=PINK, size_pt=36)
                else:
                    cell.text = str(value)
                    p = cell.text_frame.paragraphs[0]
                    p.font.size = Pt(12)

        # NEW: color the oval named "Oval 10" to reflect overall maturity (APM current Analysis).
        overall_tier = overall_maturity_from_df(df_current_analysis, grade_func=_grade_token)
        if overall_tier:
            color_oval_for_maturity(slide, shape_name="Oval 10", tier=overall_tier, update_text=False)

    # --- Add notes explaining why the status is Silver (or other tier) ---
        def _tier_counts(df):
            counts = {'bronze': 0, 'silver': 0, 'gold': 0, 'platinum': 0}
            col = 'OverallAssessment'
            if df is None or col not in df.columns:
                return counts, 0
            for v in df[col]:
                t = _grade_token(v)
                if t in counts:
                    counts[t] += 1
            total = sum(counts.values())
            return counts, total

        def _pct(n, d):
            return (n / d) * 100.0 if d else 0.0

        tier_counts, tier_total = _tier_counts(df_current_analysis)
        b, s, g, p = (tier_counts['bronze'], tier_counts['silver'], tier_counts['gold'], tier_counts['platinum'])
        pb, ps, pg, pp_ = (_pct(b, tier_total), _pct(s, tier_total), _pct(g, tier_total), _pct(p, tier_total))

        # Build a clear, single-paragraph rationale line.
        rationale = (
            f"Status is {overall_tier} because it has the largest share of rated apps this run. "
            f"Distribution — Platinum {pp_:.1f}% ({p}), Gold {pg:.1f}% ({g}), "
            f"Silver {ps:.1f}% ({s}), Bronze {pb:.1f}% ({b})."
        )

        coverage_note = (
            f"Rated coverage: {cov_curr:.1f}% ({rated_curr}/{total_curr})."
            if cov_curr is not None and total_curr > 0 else "Rated coverage: n/a."
        )

        next_focus_note = f"Next focus: {next_focus_text}."

        # Write to slide notes.
        notes = slide.notes_slide  # creates one if it doesn't exist
        tf = notes.notes_text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = "Overall tier selection: majority of app ratings in Analysis; ties prefer the higher tier."
        p2 = tf.add_paragraph()
        p2.text = rationale
        p3 = tf.add_paragraph()
        p3.text = f"{coverage_note} {next_focus_note}"
        
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Now handle Slide 4 table with "Upgraded" applications **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        slide = prs.slides[3]  # Slide 4 (index 3)
        upgraded_apps = df_analysis[df_analysis['OverallAssessment'].str.contains('upgraded', case=False, na=False)]['name'].tolist()

        # Count the number of applications in the current "Analysis" sheet
        current_analysis_df = pd.read_excel(current_file_path, sheet_name='Analysis')  # Load the current "Analysis" sheet
        number_of_apps = len(current_analysis_df)

        # Insert the count into TextBox 7
        textbox_7 = None
        for shape in slide.shapes:
            if shape.name == "TextBox 7":
                textbox_7 = shape
                break

        if textbox_7:
            textbox_7.text = f"{number_of_apps}"  # Set the text with the count
        else:
            logging.warning("TextBox 8 not found on Slide 3.")

        # Insert Upgraded Applications Table onto Slide 3 (Slide index 2) - using Table Placeholder 1
        upgraded_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")  # We are now using the same placeholder
        if upgraded_placeholder:
            logging.debug("Found Upgraded Applications table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(upgraded_apps) + 1, 1)
        else:
            logging.warning("Upgraded Applications table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(upgraded_apps) + 1, 1, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  

        # Add header for the new table
        table.cell(0, 0).text = "Applications with Upgraded Metrics"
        table.cell(0, 0).text_frame.paragraphs[0].font.size = Pt(12)

        # Populate the table with upgraded applications
        for idx, app in enumerate(upgraded_apps):
            table.cell(idx + 1, 0).text = app
            table.cell(idx + 1, 0).text_frame.paragraphs[0].font.size = Pt(12)

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert Summary Table onto Slide 5 (Slide index 4) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*
        slide = prs.slides[4]  # Slide 5 (index 4)
        summary_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")  # Placeholder for Summary Table

        if summary_placeholder:
            logging.debug("Found Summary table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(summary_df) + 1, len(summary_df.columns))
        else:
            logging.warning("Summary table placeholder not found. Adding manually.")
            # Explicitly add a new table with defined dimensions for Slide 5
            table = slide.shapes.add_table(len(summary_df) + 1, len(summary_df.columns), Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  

        # Set column headers for the Summary table
        for col_idx, column in enumerate(summary_df.columns):
            table.cell(0, col_idx).text = str(column)
            table.cell(0, col_idx).text_frame.paragraphs[0].font.size = Pt(12)

        # Populate table with Summary data
        for row_idx, row in summary_df.iterrows():
            for col_idx, value in enumerate(row):
                table.cell(row_idx + 1, col_idx).text = str(value)
                table.cell(row_idx + 1, col_idx).text_frame.paragraphs[0].font.size = Pt(12)

        # Add the title for Slide 4 (Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Comparison Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white


        # Load the Analysis sheet
        df = pd.read_excel(comparison_result_path, sheet_name='Analysis')

        columns = [
            'AppAgentsAPM', 'MachineAgentsAPM', 'BusinessTransactionsAPM',
            'BackendsAPM', 'OverheadAPM', 'ServiceEndpointsAPM',
            'ErrorConfigurationAPM', 'HealthRulesAndAlertingAPM',
            'DataCollectorsAPM', 'DashboardsAPM', 'OverallAssessment'
        ]

        results = {}
        total_applications = len(df)
        
        for col in columns:
            df[col] = df[col].astype(str)
            upgraded_count = df[col].str.contains('upgraded', case=False, na=False).sum()
            downgraded_count = df[col].str.contains('downgraded', case=False, na=False).sum()

            # Total applications is the length of the column
            total_applications = len(df[col])

            overall_result = "Increase" if upgraded_count > downgraded_count else "Decrease" if downgraded_count > upgraded_count else "Even"
            percentage_value = 0 if overall_result == "Even" else round((upgraded_count / total_applications) * 100)

            # Log the results for each column
            # logging.debug(f"Column: {col}")
            # logging.debug(f"Upgraded Count: {upgraded_count}")
            # logging.debug(f"Total Applications: {total_applications}")
            # logging.debug(f"Overall Result: {overall_result}")
            # logging.debug(f"Percentage: {percentage_value}%")

            results[col] = {
                'upgraded': upgraded_count,
                'downgraded': downgraded_count,
                'overall_result': overall_result,
                'percentage': percentage_value
            }

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert Summary Table from Previous Workbook onto Slide 5 (Table Placeholder 4) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        slide = prs.slides[4]  # Slide 5 (index 4)

        # Load the previous summary data
        previous_summary_df = pd.read_excel(previous_file_path, sheet_name='Summary')

        # Add to Table Placeholder 4 (for previous summary)
        summary_placeholder_previous = find_table_placeholder_by_name(slide, "Table Placeholder 4")  # Placeholder for Previous Summary Table
        if summary_placeholder_previous:
            logging.debug("Found Table Placeholder 4. Inserting table for previous summary.")
            table_previous = insert_table_at_placeholder(slide, "Table Placeholder 4", len(previous_summary_df) + 1, len(previous_summary_df.columns))
        else:
            logging.warning("Table Placeholder 4 not found. Adding manually.")
            table_previous = slide.shapes.add_table(len(previous_summary_df) + 1, len(previous_summary_df.columns), Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  

        # Set column headers for the previous summary table
        for col_idx, column in enumerate(previous_summary_df.columns):
            table_previous.cell(0, col_idx).text = str(column)
            table_previous.cell(0, col_idx).text_frame.paragraphs[0].font.size = Pt(12)

        # Populate the table with previous summary data
        for row_idx, row in previous_summary_df.iterrows():
            for col_idx, value in enumerate(row):
                table_previous.cell(row_idx + 1, col_idx).text = str(value)
                table_previous.cell(row_idx + 1, col_idx).text_frame.paragraphs[0].font.size = Pt(12)

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert Summary Table from Current Workbook onto Slide 5 (Table Placeholder 4)
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # Load the current summary data
        current_summary_df = pd.read_excel(current_file_path, sheet_name='Summary')

        # Add to Table Placeholder 3 (for current summary)
        summary_placeholder_current = find_table_placeholder_by_name(slide, "Table Placeholder 3")  # Placeholder for Current Summary Table
        if summary_placeholder_current:
            logging.debug("Found Table Placeholder 3. Inserting table for current summary.")
            table_current = insert_table_at_placeholder(slide, "Table Placeholder 3", len(current_summary_df) + 1, len(current_summary_df.columns))
        else:
            logging.warning("Table Placeholder 3 not found. Adding manually.")
            table_current = slide.shapes.add_table(len(current_summary_df) + 1, len(current_summary_df.columns), Inches(0.5), Inches(6), Inches(9), Inches(4)).table  

        # Set column headers for the current summary table
        for col_idx, column in enumerate(current_summary_df.columns):
            table_current.cell(0, col_idx).text = str(column)
            table_current.cell(0, col_idx).text_frame.paragraphs[0].font.size = Pt(12)

        # Populate the table with current summary data
        for row_idx, row in current_summary_df.iterrows():
            for col_idx, value in enumerate(row):
                table_current.cell(row_idx + 1, col_idx).text = str(value)
                table_current.cell(row_idx + 1, col_idx).text_frame.paragraphs[0].font.size = Pt(12)

        # Add the title for Slide 4 (Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Comparison Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white


        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*
        # Insert Overall Assessment Table onto Slide 7 (Slide index 6)
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*
        slide = prs.slides[6]  # Slide 7
        overall_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")  # Placeholder name

        if overall_placeholder:
            # logging.debug("Found Overall Assessment table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", 2, 5)
        else:
            # logging.warning("Overall Assessment table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(2, 5, Inches(0.5), Inches(1.5), Inches(9), Inches(1.5)).table  

        headers = ['Metric', '# of Apps Improved', '# Apps Degraded', 'Overall Result', 'Percentage Value']
        for col_idx, header in enumerate(headers):
            table.cell(0, col_idx).text = header
            table.cell(0, col_idx).text_frame.paragraphs[0].font.size = Pt(14)

        overall_assessment = results['OverallAssessment']
        table.cell(1, 0).text = 'OverallAssessment'
        table.cell(1, 1).text = str(overall_assessment['upgraded'])
        table.cell(1, 2).text = str(overall_assessment['downgraded'])
        table.cell(1, 3).text = overall_assessment['overall_result']
        table.cell(1, 4).text = f"{overall_assessment['percentage']}%"

        if overall_assessment['overall_result'] == "Increase":
            table.cell(1, 4).fill.solid()
            table.cell(1, 4).fill.fore_color.rgb = RGBColor(0, 255, 0)  # Green

        # Add the title for Slide 7
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Overall Assessment Result"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*
        # Insert Status Table onto Slide 8 (Slide index 7)
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*
        slide = prs.slides[7]  # Slide 8 (index 7)
        status_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")  # Placeholder name

        num_rows = len(columns)  # Should match the expected row count
        num_cols = 5  

        if status_placeholder:
            # logging.debug("Found Status table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", num_rows, num_cols)
        else:
            # logging.warning("Status table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  

        headers = ['Metric', '# of Apps Improved', '# Apps Degraded', 'Overall Result', 'Percentage Value']
        for col_idx, header in enumerate(headers):
            table.cell(0, col_idx).text = header
            table.cell(0, col_idx).text_frame.paragraphs[0].font.size = Pt(14)

        for i, col in enumerate(columns[:-1]):  
            table.cell(i + 1, 0).text = col
            table.cell(i + 1, 1).text = str(results[col]['upgraded'])
            table.cell(i + 1, 2).text = str(results[col]['downgraded'])
            table.cell(i + 1, 3).text = results[col]['overall_result']
            table.cell(i + 1, 4).text = f"{results[col]['percentage']}%"

            if results[col]['overall_result'] == "Increase":
                table.cell(i + 1, 4).fill.solid()
                table.cell(i + 1, 4).fill.fore_color.rgb = RGBColor(0, 255, 0)  # Green

        # Add the title for Slide 8
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "APM Maturity Assessment Result"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert APM AGENT Downgrade Table onto Slide 12 (Slide index 11) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[11]  # Slide 12 (index 11)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the AppAgentsAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['AppAgentsAPM']
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 12 (Slide index 11)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            logging.debug(f"Grade: {grade}, Applications: {', '.join(downgrade_data[grade]['applications'])}")

            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = ', '.join(downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 12 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "APM Agent - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'percentAgentsLessThan1YearOld': 'Rectangle 11',
            'metricLimitNotHit': 'Rectangle 10',
            'percentAgentsLessThan2YearsOld': 'Rectangle 12',
            'percentAgentsReportingData': 'Rectangle 13',
            'percentAgentsRunningSameVersion': 'Rectangle 14'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_app_agents.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'declined' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        # for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 12
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 15.") 


        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert MACHINE AGENT Downgrade Table onto Slide 13 (Slide index 12) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[12]  # Slide 13 (index 12)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the MachineAgentsAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['MachineAgentsAPM']  # Use MachineAgentsAPM column for comparison
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 13 (Slide index 12)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            logging.debug(f"Grade: {grade}, Applications: {', '.join(downgrade_data[grade]['applications'])}")
            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = ', '.join(downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 11 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Machine Agent - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'percentAgentsLessThan1YearOld': 'Rectangle 8',
            'percentAgentsLessThan2YearsOld': 'Rectangle 9',
            'percentAgentsReportingData': 'Rectangle 10',
            'percentAgentsRunningSameVersion': 'Rectangle 11',
            'percentAgentsInstalledAlongsideAppAgents': 'Rectangle 12'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_machine_agents.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'declined' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        # for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 13
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 13.") 

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert BT Downgrade Table onto Slide 14 (Slide index 13) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[13]  # Slide 14 (index 13)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the BusinessTransactionsAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['BusinessTransactionsAPM']  # Use BusinessTransactionsAPM column for comparison
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 14 (Slide index 13)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            logging.debug(f"Grade: {grade}, Applications: {', '.join(downgrade_data[grade]['applications'])}")
            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = ', '.join(downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 14 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Business Transactions - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'numberOfBTs': 'Rectangle 17',
            'percentBTsWithLoad': 'Rectangle 18',
            'btLockdownEnabled': 'Rectangle 19',
            'numberCustomMatchRules': 'Rectangle 20'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_BTs.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'decreased' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        # for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 14
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 14.") 

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert Backend Downgrade Table onto Slide 15 (Slide index 14) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[14]  # Slide 15 (index 14)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the BackendsAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['BackendsAPM']  # Use BackendsAPM column for comparison
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 15 (Slide index 14)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            # Convert all items in the applications list to strings
            applications_str = ', '.join(str(app) for app in downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            
            # Log the grade and applications
            logging.debug(f"Grade: {grade}, Applications: {applications_str}")
            
            # Populate the table with the data
            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = applications_str  # Display the application names
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 15 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Backends - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'percentBackendsWithLoad': 'Rectangle 10',
            'backendLimitNotHit': 'Rectangle 11',
            'numberOfCustomBackendRules': 'Rectangle 12'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_Backends.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'decreased' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        # for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 15
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 15.") 

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert SEP Downgrade Table onto Slide 16 (Slide index 15) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[15]  # Slide 16 (index 15)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the ServiceEndpointsAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['ServiceEndpointsAPM']  # Use ServiceEndpointsAPM column for comparison
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 16 (Slide index 15)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            logging.debug(f"Grade: {grade}, Applications: {', '.join(downgrade_data[grade]['applications'])}")
            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = ', '.join(downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 14 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Service Endpoints - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'numberOfCustomServiceEndpointRules': 'Rectangle 10',
            'serviceEndpointLimitNotHit': 'Rectangle 11',
            'percentServiceEndpointsWithLoadOrDisabled': 'Rectangle 12'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_ServiceEndpoints.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'decreased' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        #for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 16
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 14.") 

        
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert ERROR CONFIG Downgrade Table onto Slide 17 (Slide index 16) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[16]  # Slide 17 (index 16)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the ErrorConfigurationAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['ErrorConfigurationAPM']  # Use ErrorConfigurationAPM column for comparison
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 15 (Slide index 14)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            logging.debug(f"Grade: {grade}, Applications: {', '.join(downgrade_data[grade]['applications'])}")
            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = ', '.join(downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 17 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Error Configuration - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'successPercentageOfWorstTransaction': 'Rectangle 10',
            'numberOfCustomRules': 'Rectangle 11'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_ErrorConfiguration.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'decreased' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        #for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 17
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 15.") 

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert HR & ALERTS Downgrade Table onto Slide 18 (Slide index 17) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[17]  # Slide 18 (index 17)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the HealthRulesAndAlertingAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['HealthRulesAndAlertingAPM']  # Use HealthRulesAndAlertingAPM column for comparison
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 18 (Slide index 17)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            logging.debug(f"Grade: {grade}, Applications: {', '.join(downgrade_data[grade]['applications'])}")
            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = ', '.join(downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 18 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Health Rules & Alerting - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'numberOfHealthRuleViolations': 'Rectangle 10',
            'numberOfDefaultHealthRulesModified': 'Rectangle 11',
            'numberOfActionsBoundToEnabledPolicies': 'Rectangle 12',
            'numberOfCustomHealthRules': 'Rectangle 13'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_HealthRulesAndAlerting.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'decreased' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        #for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 18
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 16.") 

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert DATA COLLECTORS Downgrade Table onto Slide 19 (Slide index 18) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[18]  # Slide 19 (index 18)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the DataCollectorsAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['DataCollectorsAPM']  # Use DataCollectorsAPM column for comparison
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 19 (Slide index 18)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            logging.debug(f"Grade: {grade}, Applications: {', '.join(downgrade_data[grade]['applications'])}")
            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = ', '.join(downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 19 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Data Collectors - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'numberOfDataCollectorFieldsConfigured': 'Rectangle 10',
            'numberOfDataCollectorFieldsCollectedInSnapshots': 'Rectangle 11',
            'numberOfDataCollectorFieldsCollectedInAnalytics': 'Rectangle 12',
            'biqEnabled': 'Rectangle 13'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_DataCollectors.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'decreased' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        #for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 19
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 17.") 

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert DASHBOARDS Downgrade Table onto Slide 20 (Slide index 19) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[19]  # Slide 20 (index 19)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the DashboardsAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['DashboardsAPM']  # Use DashboardsAPM column for comparison
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 20 (Slide index 19)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            logging.debug(f"Grade: {grade}, Applications: {', '.join(downgrade_data[grade]['applications'])}")
            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = ', '.join(downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 20 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Dashboards - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

                # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'numberOfDashboards': 'Rectangle 10',
            'percentageOfDashboardsModifiedLast6Months': 'Rectangle 11',
            'numberOfDashboardsUsingBiQ': 'Rectangle 12'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_Dashboards.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'decreased' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        #for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 20
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 20.") 

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert OVERHEAD Downgrade Table onto Slide 21 (Slide index 20) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[20]  # Slide 21 (index 20)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the OverheadAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['OverheadAPM']  # Use OverheadAPM column for comparison
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 21 (Slide index 20)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            logging.debug(f"Grade: {grade}, Applications: {', '.join(downgrade_data[grade]['applications'])}")
            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = ', '.join(downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 20 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Overhead - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'developerModeNotEnabledForAnyBT': 'Rectangle 10',
            'findEntryPointsNotEnabled': 'Rectangle 11',
            'aggressiveSnapshottingNotEnabled': 'Rectangle 12',
            'developerModeNotEnabledForApplication': 'Rectangle 13'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_Overhead.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'changed' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        #for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 20
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 20.")

        # Save the PowerPoint
        prs.save(powerpoint_output_path)
        logging.debug(f"PowerPoint saved to {powerpoint_output_path}.")

    except Exception as e:
        logging.error(f"Error generating PowerPoint: {e}", exc_info=True)
        raise