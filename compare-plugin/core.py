import os
import json
import sys
import logging
import webbrowser
from flask import Flask, request, send_file, render_template
import openpyxl
from openpyxl.styles import PatternFill, Font
from openpyxl.utils.dataframe import dataframe_to_rows
from openpyxl import load_workbook
import pandas as pd
from copy import copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.shapes.placeholder import TablePlaceholder
import xlwings as xw
import datetime as dt
import re
from flask import jsonify
from typing import Optional, Dict, Any, Tuple, List
from pathlib import Path

# Set up logging configuration
def setup_logging():
    # Set log level to INFO directly
    log_level = logging.INFO  # INFO, WARNING, ERROR, CRITICAL
    logging.basicConfig(
        level=log_level,
        format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
    )
    logging.info("Logging is set up!")

# Call the setup_logging function to initialize logging
setup_logging()

app = Flask(__name__, template_folder='templates', static_folder='static')

@app.route('/', methods=['GET'], endpoint='home')
def home():
    return render_template('index.html')

@app.route('/insights', methods=['GET'], endpoint='insights_page')
def insights_page():
    return render_template('insights.html')

# --- Insights API helpers and routes ---

def _last_json_path_for_domain(domain):
    return app.config.get(f"LAST_JSON_{domain.upper()}")

def _last_result_path_for_domain(domain):
    return app.config.get(f"LAST_RESULT_{domain.upper()}")

def ts_now():
    return dt.datetime.utcnow().strftime("%Y%m%d_%H%M%S")

def ymd(dt_obj):
    try:
        return dt_obj.strftime("%Y%m%d")
    except Exception:
        return None

def file_mtime_ymd(path):
    try:
        return dt.datetime.utcfromtimestamp(os.path.getmtime(path)).strftime("%Y%m%d")
    except Exception:
        return None

def workbook_created_or_modified_ymd(path):
    try:
        wb = load_workbook(path, read_only=True, data_only=True)
        props = wb.properties
        created = ymd(props.created)
        modified = ymd(props.modified)
        wb.close()
        # Prefer created; otherwise modified; otherwise filesystem mtime
        return created or modified or file_mtime_ymd(path)
    except Exception as e:
        logging.warning("Could not read workbook properties: %s", e)
        return file_mtime_ymd(path)

def normalize_header(v):
    return str(v).strip().lower() if v is not None else ""

def _digits(s: Optional[str]) -> str:
    """Keep only digits; safe for strings like '2024-11-06'."""
    return "".join(re.findall(r"\d", (s or "")))

def _norm_day(s: Optional[str]) -> str:
    """Return YYYYMMDD (8 digits) or '' if not available."""
    d = _digits(s)
    return d[:8] if len(d) >= 8 else ""

def _norm_compare(s: Optional[str]) -> str:
    """Return YYYYMMDDHHMMSS (14 digits) or best-effort fallback."""
    d = _digits(s)
    if len(d) >= 14:
        return d[:14]
    elif len(d) >= 8:
        # If we only have a day, pad with zeros for a stable tie-breaker.
        return d[:8] + "000000"
    return ""

def safe_slug(s, max_len=60):
    s = (s or "controller").lower()
    s = re.sub(r"[^a-z0-9]+", "-", s)
    s = re.sub(r"-+", "-", s).strip("-")
    return s[:max_len] or "controller"

def find_controller_name_xlsx(path):
    try:
        wb = load_workbook(path, read_only=True, data_only=True)
        for ws in wb.worksheets:
            # Inspect first 10 rows for a header called "controller"
            for r in range(1, min(ws.max_row, 10) + 1):
                row_vals = [ws.cell(row=r, column=c).value for c in range(1, ws.max_column + 1)]
                headers = [normalize_header(v) for v in row_vals]
                if "controller" in headers:
                    idx = headers.index("controller") + 1  # 1-based column index
                    # Find first non-empty value under that column
                    for rr in range(r + 1, ws.max_row + 1):
                        val = ws.cell(row=rr, column=idx).value
                        if val and str(val).strip():
                            wb.close()
                            return str(val).strip()
        wb.close()
    except Exception as e:
        logging.warning("Controller name not found in %s: %s", path, e)
    return None

def build_compare_token(previous_path, current_path, controller_hint=None):
    prev_date = workbook_created_or_modified_ymd(previous_path) or "unknownprev"
    curr_date = workbook_created_or_modified_ymd(current_path) or "unknowncurr"
    compare_date = ts_now()
    # Prefer hint, else try to read from current, else previous
    controller = controller_hint or find_controller_name_xlsx(current_path) or find_controller_name_xlsx(previous_path) or "controller"
    controller_slug = safe_slug(controller)
    return controller_slug, prev_date, curr_date, compare_date

def _find_app_name_column(df):
    candidates = ["name", "application", "applicationName", "Application Name"]
    for c in candidates:
        if c in df.columns:
            return c
    for c in df.columns:
        if df[c].dtype == object:
            return c
    return df.columns[0] if len(df.columns) else None

#####################################################################################
############## Utility for Index on Read (compare multiple output) ##################
#####################################################################################

def _slug(s: Optional[str]) -> str:
    import re
    s = (s or "").strip().lower()
    s = re.sub(r"[^a-z0-9]+", "-", s)
    return re.sub(r"-+", "-", s).strip("-")

def _infer_compare_from_name(name: str, prefix: str) -> Optional[str]:
    base = name[len(prefix):-5]  # strip 'summary_<domain>_' and '.json'
    if "." in base:
        parts = base.split(".")
        return parts[-1] if len(parts) >= 4 else None
    return base or None

def scan_runs(
    result_folder: str,
    domain: str,
    controller_filter: Optional[str] = None,
    limit: int = 50,
) -> List[Dict[str, Any]]:
    prefix = f"summary_{domain.lower()}_"
    want_slug = _slug(controller_filter) if controller_filter else None

    try:
        names = [n for n in os.listdir(result_folder) if n.startswith(prefix) and n.endswith(".json")]
    except FileNotFoundError:
        names = []

    items: List[Dict[str, Any]] = []
    for name in names:
        path = os.path.join(result_folder, name)
        try:
            with open(path, "r", encoding="utf-8") as f:
                data = json.load(f)
        except Exception:
            continue

        meta = data.get("meta") or {}
        overall = data.get("overall") or {}
        tiers = data.get("tiers") or {}

        controller = (meta.get("controller") or "Unknown").strip()
        prev = (meta.get("previousDate") or "Unknown").strip()
        curr = (meta.get("currentDate") or "Unknown").strip()
        compare = (meta.get("compareDate") or "").strip()
        if not compare:
            compare = _infer_compare_from_name(name, prefix) or ""

        if want_slug:
            if _slug(controller) != want_slug and want_slug not in name.lower():
                continue

        sort_prev = _norm_day(prev)
        sort_curr = _norm_day(curr)
        sort_cmp  = _norm_compare(compare)

        items.append({
            "file": name,
            "path": path,
            "controller": controller,
            "previousDate": prev,
            "currentDate": curr,
            "compareDate": compare,
            "improved": int(overall.get("improved") or 0),
            "degraded": int(overall.get("degraded") or 0),
            "percentage": int(overall.get("percentage") or 0),
            "tiers": {
                "platinum": tiers.get("platinum"),
                "goldOrBetter": tiers.get("goldOrBetter"),
                "silverOrBetter": tiers.get("silverOrBetter"),
            },
            "sortPrev": sort_prev,
            "sortCurr": sort_curr,
            "sortCompare": sort_cmp,
        })

    # Semantic sort: Previous (asc) → Current (asc) → Compare timestamp (asc).
    items.sort(key=lambda x: (x.get("sortPrev", ""), x.get("sortCurr", ""), x.get("sortCompare", "")))

    if limit and limit > 0:
        items = items[:limit]
    return items

@app.route("/api/apps", methods=["GET"])
def api_apps():
    # Accept any case from the client, validate on lowercase.
    raw_domain = (request.args.get("domain") or "").strip()
    domain_l = raw_domain.lower()
    if domain_l not in ("apm", "brum", "mrum"):
        return jsonify({"error": "Invalid domain. Use APM, BRUM, or MRUM."}), 400
    domain_u = domain_l.upper()

    # Optional file override (specific snapshot JSON).
    file_override = (request.args.get("file") or "").strip()
    if file_override:
        json_path = os.path.join(app.config["RESULT_FOLDER"], file_override)
    else:
        # Try helpers with both lower and upper domain for robustness.
        json_path = _last_json_path_for_domain(domain_l) or _last_json_path_for_domain(domain_u)

    logging.info("[/api/apps] domain=%s (l=%s,u=%s) json_path=%s", raw_domain, domain_l, domain_u, json_path)

    # Try JSON first (override or latest).
    if json_path and os.path.isfile(json_path):
        try:
            with open(json_path, "r", encoding="utf-8") as f:
                data = json.load(f)
            names = (data.get("apps", {}) or {}).get("names") or list((data.get("appsIndex") or {}).keys())
            # Return 200 with whatever we found (possibly empty list).
            return jsonify({"domain": domain_u, "apps": names or []})
        except Exception as e:
            logging.warning("[/api/apps] Failed reading JSON %s: %s", json_path, e)

    # Fallback to the last workbook if JSON is missing.
    wb_path = _last_result_path_for_domain(domain_l) or _last_result_path_for_domain(domain_u)
    logging.info("[/api/apps] wb_path=%s", wb_path)

    if wb_path and os.path.isfile(wb_path):
        try:
            df = pd.read_excel(wb_path, sheet_name="Analysis")
            name_col = _find_app_name_column(df)
            if name_col:
                names = (
                    df[name_col]
                    .dropna()
                    .astype(str)
                    .str.strip()
                    .replace({"": None})
                    .dropna()
                    .unique()
                    .tolist()
                )
                return jsonify({"domain": domain_u, "apps": names or []})
        except Exception as e:
            logging.error("[/api/apps] Fallback to workbook failed (%s): %s", wb_path, e, exc_info=True)

    # Final: return 200 with an empty list to avoid frontend errors when no data is available.
    return jsonify({"domain": domain_u, "apps": [], "message": "No data available for domain. Upload a comparison first."})

@app.route("/api/insights", methods=["GET"])
def api_insights():
    domain = (request.args.get("domain") or "").upper()
    app_name = request.args.get("app") or ""
    file = request.args.get("file") or ""  # optional: specific summary filename

    if domain not in ("APM","BRUM","MRUM") or not app_name:
        return jsonify({"error": "Missing domain or app."}), 400

    folder = app.config['RESULT_FOLDER']
    # Choose file: specific or latest for domain.
    if file:
        path = os.path.join(folder, file)
    else:
        # Latest stamped summary for the domain.
        prefix = f"summary_{domain.lower()}_"
        try:
            names = [n for n in os.listdir(folder) if n.startswith(prefix) and n.endswith(".json")]
        except FileNotFoundError:
            names = []
        names.sort(reverse=True)  # stamped names sort newest-first lexicographically
        path = os.path.join(folder, names[0]) if names else ""

    if not path or not os.path.exists(path):
        return jsonify({"error": "Snapshot not found."}), 404

    with open(path, "r", encoding="utf-8") as f:
        data = json.load(f)

    apps_index = (data.get("appsIndex") or {})
    entry = apps_index.get(app_name)
    if not entry:
        # Try a normalized match (some apps differ by whitespace/case)
        key = app_name.strip().lower()
        for k, v in apps_index.items():
            if k.strip().lower() == key:
                entry = v
                break

    if not entry:
        return jsonify({"error": "App not found in snapshot."}), 404

    areas = entry.get("areas", [])
    detail = entry.get("detail", {})

    return jsonify({
        "domain": domain,
        "app": app_name,
        "areas": areas,
        "detail": detail,
        "meta": data.get("meta", {}),
    })

@app.route("/api/history", methods=["GET"])
def api_history():
    domain = (request.args.get("domain") or "").lower()
    if domain not in ("apm", "brum", "mrum"):
        return jsonify({"error": "Invalid domain."}), 400

    folder = app.config['RESULT_FOLDER']
    prefix = f"summary_{domain}_"
    items = []

    for name in sorted(os.listdir(folder)):
        if not (name.startswith(prefix) and name.endswith(".json")):
            continue

        base = name[len(prefix):-5]  # strip prefix and ".json"
        parts = base.split(".")      # stamped: [controller, prev, curr, compare]

        controller = prev = curr = None
        compare = None

        if len(parts) == 4:
            controller, prev, curr, compare = parts
        elif len(parts) == 1:
            # Older files: summary_<domain>_<YYYYMMDD_HHMMSS>.json
            compare = parts[0]
            # Try to read meta from the JSON for richer labels, if available.
            try:
                with open(os.path.join(folder, name), "r", encoding="utf-8") as f:
                    data = json.load(f)
                meta = (data.get("meta") or {})
                controller = meta.get("controller")
                prev = meta.get("previousDate")
                curr = meta.get("currentDate")
            except Exception:
                pass
        else:
            # Unexpected pattern; attempt meta read as a best effort.
            try:
                with open(os.path.join(folder, name), "r", encoding="utf-8") as f:
                    data = json.load(f)
                meta = (data.get("meta") or {})
                controller = meta.get("controller")
                prev = meta.get("previousDate")
                curr = meta.get("currentDate")
                compare = meta.get("compareDate") or compare
            except Exception:
                pass

        items.append({
            "file": name,
            "timestamp": compare or "",
            "controller": controller,
            "prev": prev,
            "curr": curr
        })

    # NEW: optional controller filter
    controller_q = request.args.get("controller")
    if controller_q:
        want = _slug(controller_q)
        filtered = []
        for it in items:
            # Fast path: check the stamped filename's controller slug.
            if it["controller"] and _slug(it["controller"]) == want:
                filtered.append(it)
                continue
            # Fallback: open the JSON and check meta.controller if needed.
            try:
                with open(os.path.join(folder, it["file"]), "r", encoding="utf-8") as f:
                    data = json.load(f)
                meta = (data.get("meta") or {})
                if _slug(meta.get("controller") or "") == want:
                    filtered.append(it)
            except Exception:
                pass
        items = filtered

    items.sort(key=lambda x: x["timestamp"], reverse=True)
    return jsonify({"domain": domain.upper(), "items": items})

@app.route("/api/trends/runs", methods=["GET"])
def api_trends_runs():
    domain = (request.args.get("domain") or "").lower()
    if domain not in ("apm", "brum", "mrum"):
        return jsonify({"error": "Invalid domain."}), 400

    controller = request.args.get("controller")
    try:
        limit = int(request.args.get("limit", "20"))
    except ValueError:
        limit = 20

    baseline = (request.args.get("baseline") or "").lower()

    folder = app.config["RESULT_FOLDER"]
    runs = scan_runs(folder, domain=domain, controller_filter=controller, limit=limit)

    if baseline == "earliestprev":
        prevs = [r["sortPrev"] for r in runs if r.get("sortPrev")]
        if prevs:
            earliest_prev = min(prevs)
            runs = [r for r in runs if r.get("sortPrev") == earliest_prev]

    series = [{
        "compareDate": r["compareDate"],
        "previousDate": r["previousDate"],
        "currentDate": r["currentDate"],
        "improved": r["improved"],
        "degraded": r["degraded"],
        "percentage": r["percentage"],
        "tiers": r["tiers"],
        "file": r["file"],
    } for r in runs]

    label = controller or (runs[0]["controller"] if runs else None)
    return jsonify({"domain": domain.upper(), "controller": label, "count": len(series), "items": series})

def save_workbook(filepath):
    """Open and save the workbook to ensure formulas are recalculated."""
    # Open the workbook
    app = xw.App(visible=False)  # Set visible=False to avoid showing the Excel window
    wb = app.books.open(filepath)
    
    # Save the workbook
    wb.save()
    
    # Close the workbook
    wb.close()
    
    # Quit the application
    app.quit()

# Load configuration from JSON file
def load_config():
    # Determine the current script directory
    script_dir = os.path.dirname(os.path.abspath(__file__))
    
    # Construct the full path to the config.json file
    config_path = os.path.join(script_dir, 'config.json')

    if not os.path.exists(config_path):
        raise FileNotFoundError(f"Config file not found: {config_path}")

    with open(config_path) as f:
        config = json.load(f)
    
    # Convert relative paths to absolute paths
    config['upload_folder'] = os.path.join(script_dir, config.get('upload_folder', 'uploads'))
    config['result_folder'] = os.path.join(script_dir, config.get('result_folder', 'results'))
    config['TEMPLATE_FOLDER'] = os.path.join(script_dir, config.get('TEMPLATE_FOLDER', 'templates'))  # Adding this line
    
    return config

# Load the configuration
config = load_config()

# Configure upload and result folders from config
UPLOAD_FOLDER = config['upload_folder']
RESULT_FOLDER = config['result_folder']

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['RESULT_FOLDER'] = RESULT_FOLDER
logging.basicConfig(level=logging.DEBUG)

# Ensure folders exist
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['RESULT_FOLDER'], exist_ok=True)

###############################################################################################################
###############################               BUILD COMPARISON JSON      ###################################### 
###############################################################################################################
from typing import Optional, Dict, Any, Tuple

def build_comparison_json(
    domain: str,
    comparison_result_path: str,
    current_file_path: str,
    previous_file_path: str,
    result_folder: str,
    meta: Optional[Dict[str, Any]] = None,
) -> Tuple[str, str, Dict[str, Any]]:
    """
    Builds a compact JSON summary AND a per-app index for the selected domain.
    Returns (out_path, out_name, payload).
    """

    result_folder = result_folder or "."

    # ---------- Existing payload construction ----------
    df_analysis = pd.read_excel(comparison_result_path, sheet_name='Analysis')

    AREA_MAP = {
        'APM': [
            'AppAgentsAPM','MachineAgentsAPM','BusinessTransactionsAPM',
            'BackendsAPM','OverheadAPM','ServiceEndpointsAPM',
            'ErrorConfigurationAPM','HealthRulesAndAlertingAPM',
            'DataCollectorsAPM','DashboardsAPM'
        ],
        'BRUM': ['NetworkRequestsBRUM','HealthRulesAndAlertingBRUM'],
        'MRUM': ['NetworkRequestsMRUM','HealthRulesAndAlertingMRUM'],
    }
    DETAIL_SHEETS = {
        'APM': {
            'AppAgentsAPM': 'AppAgentsAPM',
            'MachineAgentsAPM': 'MachineAgentsAPM',
            'BusinessTransactionsAPM': 'BusinessTransactionsAPM',
            'BackendsAPM': 'BackendsAPM',
            'OverheadAPM': 'OverheadAPM',
            'ServiceEndpointsAPM': 'ServiceEndpointsAPM',
            'ErrorConfigurationAPM': 'ErrorConfigurationAPM',
            'HealthRulesAndAlertingAPM': 'HealthRulesAndAlertingAPM',
            'DataCollectorsAPM': 'DataCollectorsAPM',
            'DashboardsAPM': 'DashboardsAPM',
        },
        'BRUM': {
            'NetworkRequestsBRUM': 'NetworkRequestsBRUM',
            'HealthRulesAndAlertingBRUM': 'HealthRulesAndAlertingBRUM'
        },
        'MRUM': {
            'NetworkRequestsMRUM': 'NetworkRequestsMRUM',
            'HealthRulesAndAlertingMRUM': 'HealthRulesAndAlertingMRUM'
        },
    }

    areas = [c for c in AREA_MAP.get(domain, []) if c in df_analysis.columns]

    # Total apps by Analysis.name
    try:
        app_total = df_analysis['name'].dropna().astype(str).str.strip().ne('').sum()
    except Exception:
        app_total = len(df_analysis)

    # Overall improved/degraded from Analysis
    def count_changes(df, col):
        if col not in df.columns:
            return 0, 0, [], []
        s = df[col].astype(str)
        improved = s.str.contains('Upgraded', case=False, na=False)
        degraded = s.str.contains('Downgraded', case=False, na=False)
        imp_names = df.loc[improved, 'name'].astype(str).str.strip().tolist()
        deg_names = df.loc[degraded, 'name'].astype(str).str.strip().tolist()
        return improved.sum(), degraded.sum(), imp_names, deg_names

    overall_imp, overall_deg, _, _ = count_changes(df_analysis, 'OverallAssessment')
    overall_result = 'Increase' if overall_imp > overall_deg else 'Decrease' if overall_deg > overall_imp else 'Even'
    overall_pct = 0 if overall_result == 'Even' else round((overall_imp / max(1, overall_imp + overall_deg)) * 100)

    # Per-area aggregate blocks
    area_blocks = []
    for col in areas:
        imp, deg, imp_names, deg_names = count_changes(df_analysis, col)
        area_blocks.append({
            "name": col,
            "improved": int(imp),
            "degraded": int(deg),
            "improvedApps": imp_names,
            "degradedApps": deg_names
        })

    # Optional tiers from OverallAssessment<domain> in current workbook
    tiers = {}
    sheet_name = f"OverallAssessment{domain}" if domain in ('APM','BRUM','MRUM') else None
    try:
        xls = pd.ExcelFile(current_file_path) if current_file_path else None
        if xls and sheet_name in xls.sheet_names:
            df_overall = pd.read_excel(current_file_path, sheet_name=sheet_name)
            def last_pct(col):
                if col in df_overall.columns:
                    s = df_overall[col].astype(str).str.strip().str.replace('%','')
                    vals = pd.to_numeric(s, errors='coerce').dropna()
                    return f"{vals.iloc[-1]:.1f}%" if len(vals) else None
                return None
            tiers = {
                "platinum": last_pct("percentageTotalPlatinum"),
                "goldOrBetter": last_pct("percentageTotalGoldOrBetter"),
                "silverOrBetter": last_pct("percentageTotalSilverOrBetter"),
            }
            tiers = {k:v for k,v in tiers.items() if v is not None}
    except Exception:
        tiers = {}

    # Build per-app index (areas + detail metrics from detail sheets)
    appsIndex = {}
    app_names = df_analysis['name'].dropna().astype(str).str.strip().unique().tolist()

    # Read all detail sheets once
    detail_frames = {}
    for area_col, sheet in DETAIL_SHEETS.get(domain, {}).items():
        try:
            detail_frames[area_col] = pd.read_excel(comparison_result_path, sheet_name=sheet)
        except Exception:
            detail_frames[area_col] = None

    def normalize_status(val):
        s = str(val or "").lower()
        if "upgraded" in s: return "Upgraded"
        if "downgraded" in s: return "Downgraded"
        return "No Change"

    for app in app_names:
        row = df_analysis[df_analysis["name"].astype(str).str.strip() == app]
        per_app_areas = []
        per_app_detail = {}
        for area_col in areas:
            status = normalize_status(row[area_col].iloc[0] if not row.empty else None)
            per_app_areas.append({"name": area_col, "status": status})

            df_detail = detail_frames.get(area_col)
            if df_detail is not None and len(df_detail.columns) > 0:
                # Try 'application' column; fallback to first column
                app_col = 'application' if 'application' in df_detail.columns else df_detail.columns[0]
                r = df_detail[df_detail[app_col].astype(str).str.strip() == app]
                if not r.empty:
                    vals = {str(c): ("" if pd.isna(r.iloc[0][c]) else str(r.iloc[0][c])) for c in df_detail.columns}
                    per_app_detail[area_col] = vals

        appsIndex[app] = {"areas": per_app_areas, "detail": per_app_detail}

    payload = {
        "domain": domain,
        "generatedAt": dt.datetime.utcnow().replace(microsecond=0).isoformat() + "Z",
        "apps": { "total": int(app_total), "names": app_names },
        "overall": {
            "improved": int(overall_imp),
            "degraded": int(overall_deg),
            "result": overall_result,
            "percentage": int(overall_pct)
        },
        "tiers": tiers,
        "areas": area_blocks,
        "appsIndex": appsIndex
    }

    # ---------- Meta completion + stamped filename ----------
    # Helpers (local to keep imports optional)
    def _yyyymmdd(d: Optional[dt.datetime]) -> Optional[str]:
        return d.strftime("%Y%m%d") if d else None

    def _guess_workbook_date(path: Optional[str]) -> Optional[str]:
        if not path:
            return None
        try:
            from openpyxl import load_workbook
            wb = load_workbook(path, read_only=True, data_only=True)
            created = getattr(wb.properties, "created", None)
            modified = getattr(wb.properties, "modified", None)
            wb.close()
            return _yyyymmdd(created or modified)
        except Exception:
            pass
        try:
            ts = os.path.getmtime(path)
            return dt.datetime.utcfromtimestamp(ts).strftime("%Y%m%d")
        except Exception:
            return None

    def _safe_slug(s: Optional[str]) -> str:
        import re
        s = (s or "unknown").lower()
        s = re.sub(r"[^a-z0-9]+", "-", s)
        s = re.sub(r"-+", "-", s).strip("-")
        return s or "unknown"

    # Start with provided meta and fill gaps.
    m: Dict[str, Any] = dict(meta or {})

    # Controller: pull from any detail sheet that has a 'controller' column.
    if not m.get("controller"):
        controller_val = None
        for df in detail_frames.values():
            if df is not None and "controller" in df.columns:
                s = df["controller"].dropna().astype(str).str.strip()
                if len(s):
                    controller_val = s.iloc[0]
                    break
        m["controller"] = controller_val or "Unknown"

    # Dates: prefer workbook properties; fall back to filesystem times.
    if not m.get("previousDate"):
        m["previousDate"] = _guess_workbook_date(previous_file_path) or "Unknown"
    if not m.get("currentDate"):
        m["currentDate"] = _guess_workbook_date(current_file_path) or "Unknown"

    # Compare date: stamp now if missing.
    if not m.get("compareDate"):
        m["compareDate"] = dt.datetime.utcnow().strftime("%Y%m%d_%H%M%S")

    # Attach to payload.
    payload["meta"] = m

    # Filename stamped with controller + dates.
    controller_slug = _safe_slug(m.get("controller"))
    base = f"{controller_slug}.{m['previousDate']}.{m['currentDate']}.{m['compareDate']}"
    out_name = f"summary_{domain.lower()}_{base}.json"

    # Ensure output folder exists and write JSON.
    os.makedirs(result_folder, exist_ok=True)
    out_path = os.path.join(result_folder, out_name)
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(payload, f, ensure_ascii=False, indent=2)

    return out_path, out_name, payload


def check_controllers_match(previous_file_path, current_file_path):
    # Load the previous and current workbooks
    previous_workbook = pd.read_excel(previous_file_path, sheet_name='Analysis')
    current_workbook = pd.read_excel(current_file_path, sheet_name='Analysis')

    # Extract the 'controller' column from both workbooks and strip whitespaces
    previous_controllers = previous_workbook['controller'].dropna().str.strip().unique()
    current_controllers = current_workbook['controller'].dropna().str.strip().unique()

    # Log the exact controller values for debugging
    logging.debug(f"Previous controller(s): {previous_controllers}")
    logging.debug(f"Current controller(s): {current_controllers}")

    # Check if the controllers match
    if not (len(previous_controllers) == 1 and len(current_controllers) == 1 and previous_controllers[0] == current_controllers[0]):
        logging.error(f"Controllers do not match. Previous controller: {previous_controllers}, Current controller: {current_controllers}")
        return False

    return True

# Solution from Circuit
def parse_percent_to_float(val):
    """Convert values like '42.9%' or numeric into float; return None if not parsable."""
    if pd.isna(val):
        return None
    if isinstance(val, (int, float)):
        return float(val)
    s = str(val).strip()
    # If value is a change string like '11.4% → 42.9%', take the right-most token.
    if '→' in s:
        s = s.split('→')[-1].strip()
    s = s.replace('%', '')
    try:
        return float(s)
    except Exception:
        return None

def find_percent_column(df, keywords):
    """Find a column whose name contains the keywords (e.g., ['gold']) and likely holds a percentage."""
    if df is None:
        return None
    kw = [k.lower() for k in keywords]
    # Prefer names that imply percent/rate.
    candidates = []
    for col in df.columns:
        name = col.lower()
        if all(k in name for k in kw) and ('percent' in name or 'percentage' in name or '%' in name or 'rate' in name or 'ratio' in name):
            candidates.append(col)
    if candidates:
        return candidates[0]
    # Fallback: any column matching keywords.
    for col in df.columns:
        name = col.lower()
        if all(k in name for k in kw):
            return col
    return None

def change_arrow(curr, prev):
    """Return an arrow based on change direction; neutral if missing."""
    if curr is None or prev is None:
        return '→'
    if curr > prev:
        return '↑'
    if curr < prev:
        return '↓'
    return '→'

def add_key_callouts_slide(prs, current_summary_df, previous_summary_df, find_table_placeholder_by_name, insert_table_at_placeholder):
    """Create/populate the 'Assessment Result - Key Callouts' table (Table 5)."""
    # Try to locate an existing slide by title; else add one.
    target_slide = None
    for s in prs.slides:
        for shape in s.shapes:
            if getattr(shape, 'has_text_frame', False) and 'Assessment Result - Key Callouts' in shape.text:
                target_slide = s
                break
        if target_slide:
            break
    if target_slide is None:
        layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        target_slide = prs.slides.add_slide(layout)

    # Title.
    if getattr(target_slide.shapes, 'title', None):
        target_slide.shapes.title.text = "Assessment Result - Key Callouts"

    # Subtitle.
    subtitle_box = target_slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(9.0), Inches(0.6))
    subtitle_box.text_frame.text = "The table below shows the Overall Assessment APM result."

    # Extract Gold/Platinum percentages.
    gold_col_curr = find_percent_column(current_summary_df, ["gold"])
    gold_col_prev = find_percent_column(previous_summary_df, ["gold"])
    plat_col_curr = find_percent_column(current_summary_df, ["platinum"])
    plat_col_prev = find_percent_column(previous_summary_df, ["platinum"])

    curr_gold = parse_percent_to_float(current_summary_df[gold_col_curr].iloc[0]) if gold_col_curr else None
    prev_gold = parse_percent_to_float(previous_summary_df[gold_col_prev].iloc[0]) if gold_col_prev else None
    curr_platinum = parse_percent_to_float(current_summary_df[plat_col_curr].iloc[0]) if plat_col_curr else None
    prev_platinum = parse_percent_to_float(previous_summary_df[plat_col_prev].iloc[0]) if plat_col_prev else None

    # Rows to match your slide.
    rows = [
        [
            "B/S/G/P Model Adoption & Maturity Status",
            "Insert commentry",
            "Insert Outcomes",
            "↑",
        ],
        [
            "Gold Status Apps",
            "Insert commentry",
            (
                f"Increase from {prev_gold:.1f}%→{curr_gold:.1f}%"
                if (curr_gold is not None and prev_gold is not None)
                else "Increase observed."
            ),
            change_arrow(curr_gold, prev_gold),
        ],
        [
            "Platinum Status Apps",
            "Insert commentry",
            (
                f"{curr_platinum:.1f}% platinum"
                if curr_platinum is not None
                else "Platinum presence observed."
            ),
            change_arrow(curr_platinum, prev_platinum),
        ],
        [
            "Maturity Partnership",
            "Insert Commentry",
            "",
            "↑",
        ],
    ]

    # Insert table by placeholder name 'Table 5' or add new one.
    placeholder = find_table_placeholder_by_name(target_slide, "Table 5")
    if placeholder:
        table = insert_table_at_placeholder(target_slide, "Table 5", len(rows) + 1, 4)
    else:
        table = target_slide.shapes.add_table(len(rows) + 1, 4, Inches(0.6), Inches(2.1), Inches(9.0), Inches(4.0)).table

    headers = [
        "AppD Maturity Progression & Engagement",
        "Commentary",
        "Outcomes",
        "Change/Status Since Last",
    ]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(12)

    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            if c_idx == 3 and value in ("↑", "↓", "→"):
                run = p.runs[0] if p.runs else p.add_run()
                if value == "↑":
                    run.font.color.rgb = RGBColor(0, 176, 80)
                elif value == "↓":
                    run.font.color.rgb = RGBColor(192, 0, 0)
                else:
                    run.font.color.rgb = RGBColor(255, 255, 255)


PINK = RGBColor(255, 20, 147)  # Deep pink.

def set_arrow_cell(cell, direction, color=PINK, size_pt=36, font_name=None):
    """
    Render a large, filled symbol centered in the table cell, colored in `color`.
    direction: one of '↑', '↓', '→'.
    - Uses ▲ ▼ for up/down.
    - Uses ▶ with Variation Selector-15 (U+FE0E) for neutral to force text (avoid emoji 'play button').
    """
    # Force text presentation for neutral: '▶' + U+FE0E.
    glyphs = {'↑': '▲', '↓': '▼', '→': '▶\ufe0e'}
    sym = glyphs.get(direction, '▶\ufe0e')

    # Default to a very common Office font.
    font_name = font_name or 'Calibri'

    tf = cell.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    run = p.add_run()
    run.text = sym
    run.font.bold = True
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    run.font.name = font_name

    # Fallbacks if the environment substitutes a font without these glyphs
    for fname in (font_name, 'Arial', 'Helvetica', 'Arial Unicode MS'):
        if fname:
            run.font.name = fname
            break

def autosize_col_to_header(table, col_idx, header_pt=12, padding_in=0.5, avg_char_em=0.55):
    """
    Approximate a column width so it fits the header text.
    - header_pt: header font size in points.
    - padding_in: extra inches to add for padding.
    - avg_char_em: average character width as a fraction of font size.
    """
    header_text = table.cell(0, col_idx).text_frame.text or ""
    width_in = (len(header_text) * header_pt * avg_char_em) / 72.0 + padding_in  # points -> inches
    # Minimum sensible width for visibility.
    width_in = max(width_in, 1.2)
    try:
        table.columns[col_idx].width = Inches(width_in)
    except Exception:
        pass

# -------------------- NEW: maturity badge helpers (module-level, reusable) --------------------
MATURITY_COLORS = {
    'Bronze':   RGBColor(205, 127, 50),
    'Silver':   RGBColor(166, 166, 166),
    'Gold':     RGBColor(255, 192, 0),
    'Platinum': RGBColor(190, 190, 200),
}
_RANK = {'bronze': 0, 'silver': 1, 'gold': 2, 'platinum': 3}

def overall_maturity_from_df(df, grade_func, col='OverallAssessment'):
    """Compute majority tier from df[col]; ties prefer higher maturity."""
    if df is None or col not in df.columns:
        return None
    counts = {'bronze': 0, 'silver': 0, 'gold': 0, 'platinum': 0}
    for v in df[col]:
        t = (grade_func(v) or "").lower()
        if t in counts:
            counts[t] += 1
    total = sum(counts.values())
    if total == 0:
        return None
    best = max(counts.items(), key=lambda kv: (kv[1], _RANK[kv[0]]))
    return best[0].title()

def _ideal_text_rgb(bg_rgb):
    """Choose white or dark gray for readability over the fill color."""
    r, g, b = bg_rgb[0], bg_rgb[1], bg_rgb[2]
    brightness = (r * 299 + g * 587 + b * 114) / 1000.0
    return RGBColor(255, 255, 255) if brightness < 140 else RGBColor(31, 31, 31)

def color_oval_for_maturity(slide, shape_name, tier, update_text=False):
    """Fill the named oval with the tier color and set text color."""
    if tier not in MATURITY_COLORS:
        return
    target = next((sh for sh in slide.shapes if sh.name == shape_name), None)
    if target is None:
        return
    target.fill.solid()
    target.fill.fore_color.rgb = MATURITY_COLORS[tier]
    if update_text and hasattr(target, "text_frame") and target.text_frame is not None:
        tf = target.text_frame
        tf.clear()
        run = tf.paragraphs[0].add_run()
        run.text = tier
        run.font.bold = True
    if hasattr(target, "text_frame") and target.text_frame is not None:
        fg = _ideal_text_rgb(MATURITY_COLORS[tier])
        for para in target.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = fg


def build_key_callouts_for_domain(
    slide,
    domain: str,
    comparison_result_path: str,
    current_file_path: str,
    previous_file_path: str,
    find_table_placeholder_by_name,
    insert_table_at_placeholder,
    table_placeholder_name: str = "Table Placeholder 1",
    oval_name: str = "Oval 10",
):
    """
    Build the 'Assessment Result - Key Callouts' table for APM/BRUM/MRUM.
    Mirrors your APM logic (coverage, tier deltas, next focus, notes, and oval color).
    """
    domain = domain.upper()

    # Helpers local to this builder (do not override your APM ones)
    def _parse_percent_to_float(val):
        if pd.isna(val):
            return None
        if isinstance(val, (int, float)):
            return float(val)
        s = str(val).strip()
        if '→' in s:
            s = s.split('→')[-1].strip()
        s = s.replace('%', '')
        try:
            return float(s)
        except Exception:
            return None

    def _get_tier_percent(df, tier):
        name_map = {c.lower(): c for c in df.columns}
        candidates = [f"{tier.lower()} %", f"{tier.lower()}%", f"percentage{tier.lower()}", f"{tier.lower()}percentage"]
        for cand in candidates:
            if cand in name_map:
                return _parse_percent_to_float(df[name_map[cand]].iloc[0])
        needed = ['bronze', 'silver', 'gold', 'platinum']
        if all(k in name_map for k in needed):
            try:
                total = 0.0
                counts = {}
                for k in needed:
                    val = pd.to_numeric(df[name_map[k]].iloc[0], errors='coerce')
                    counts[k] = 0.0 if pd.isna(val) else float(val)
                    total += counts[k]
                if total > 0:
                    return (counts[tier.lower()] / total) * 100.0
            except Exception:
                return None
        return None

    def _trend_word(curr, prev):
        if curr is None or prev is None:
            return "held steady"
        if curr > prev:
            return "increased"
        if curr < prev:
            return "decreased"
        return "held steady"

    def _delta_pp(curr, prev):
        if curr is None or prev is None:
            return None
        return round(curr - prev, 1)

    def _arrow(curr, prev):
        if curr is None or prev is None:
            return '→'
        if curr > prev:
            return '↑'
        if curr < prev:
            return '↓'
        return '→'

    def _arrow_threshold(curr, prev, threshold_pp=0.5):
        if curr is None or prev is None:
            return '→'
        delta = curr - prev
        if delta >= threshold_pp:
            return '↑'
        if delta <= -threshold_pp:
            return '↓'
        return '→'

    def _fmt_pp_delta(prev, curr):
        if prev is None or curr is None:
            return None
        d = curr - prev
        sign = '+' if d > 0 else '−' if d < 0 else '±'
        return f"{prev:.1f}%→{curr:.1f}% ({sign}{abs(d):.1f} pp)."

    def _grade_token_local(s: str):
        if not s:
            return None
        m = re.search(r'(platinum|gold|silver|bronze)', str(s), re.I)
        return m.group(1).lower() if m else None

    def _apps_coverage(path):
        try:
            df = pd.read_excel(path, sheet_name='Analysis')
            total = int(df['name'].dropna().astype(str).str.strip().ne('').sum())
            if total == 0:
                return (0, 0, 0.0)
            rated = int(df['OverallAssessment'].apply(_grade_token_local).notna().sum())
            pct = (rated / total) * 100.0
            return (total, rated, pct)
        except Exception:
            return (0, 0, None)

    # Load frames for this domain
    try:
        curr_summary = pd.read_excel(current_file_path, sheet_name='Summary')
        prev_summary = pd.read_excel(previous_file_path, sheet_name='Summary')
        df_cmp = pd.read_excel(comparison_result_path, sheet_name='Analysis')
        df_curr_analysis = pd.read_excel(current_file_path, sheet_name='Analysis')
    except Exception as e:
        logging.error("[Key Callouts %s] Failed to load sheets: %s", domain, e, exc_info=True)
        return

    # Tier percentages (Gold/Platinum).
    curr_gold = _get_tier_percent(curr_summary, 'Gold')
    prev_gold = _get_tier_percent(prev_summary, 'Gold')
    curr_plat = _get_tier_percent(curr_summary, 'Platinum')
    prev_plat = _get_tier_percent(prev_summary, 'Platinum')
    delta_gold = _delta_pp(curr_gold, prev_gold)
    delta_plat = _delta_pp(curr_plat, prev_plat)

    # Coverage (B/S/G/P applied to this domain).
    _, _, cov_prev = _apps_coverage(previous_file_path)
    total_curr, rated_curr, cov_curr = _apps_coverage(current_file_path)
    cov_arrow = _arrow_threshold(cov_curr, cov_prev)
    cov_outcome = (
        f"Coverage: {cov_curr:.1f}% of apps rated ({rated_curr}/{total_curr})."
        if cov_curr is not None and total_curr > 0 else "Coverage data not available."
    )
    cov_prev_curr = _fmt_pp_delta(cov_prev, cov_curr)
    if cov_prev_curr:
        cov_outcome = f"{cov_outcome} {cov_prev_curr}"

    # Overall result and next focus from compare Analysis.
    def _count_changes(df, col):
        if df is None or col not in df.columns:
            return 0, 0
        s = df[col].astype(str)
        up = s.str.contains('Upgraded', case=False, na=False).sum()
        down = s.str.contains('Downgraded', case=False, na=False).sum()
        return int(up), int(down)

    up_overall, down_overall = _count_changes(df_cmp, 'OverallAssessment')
    overall_result_text = "Increase" if up_overall > down_overall else "Decrease" if down_overall > up_overall else "Even"

    if domain == "APM":
        area_cols = [
            'AppAgentsAPM','MachineAgentsAPM','BusinessTransactionsAPM',
            'BackendsAPM','OverheadAPM','ServiceEndpointsAPM',
            'ErrorConfigurationAPM','HealthRulesAndAlertingAPM',
            'DataCollectorsAPM','DashboardsAPM'
        ]
        pretty = {
            'AppAgentsAPM': 'App Agents',
            'MachineAgentsAPM': 'Machine Agents',
            'BusinessTransactionsAPM': 'Business Transactions',
            'BackendsAPM': 'Backends',
            'OverheadAPM': 'Overhead',
            'ServiceEndpointsAPM': 'Service Endpoints',
            'ErrorConfigurationAPM': 'Error Configuration',
            'HealthRulesAndAlertingAPM': 'Health Rules & Alerting',
            'DataCollectorsAPM': 'Data Collectors',
            'DashboardsAPM': 'Dashboards',
        }
    elif domain == "BRUM":
        area_cols = ['NetworkRequestsBRUM','HealthRulesAndAlertingBRUM']
        pretty = {
            'NetworkRequestsBRUM': 'Network Requests',
            'HealthRulesAndAlertingBRUM': 'Health Rules & Alerting',
        }
    else:  # MRUM
        area_cols = ['NetworkRequestsMRUM','HealthRulesAndAlertingMRUM']
        pretty = {
            'NetworkRequestsMRUM': 'Network Requests',
            'HealthRulesAndAlertingMRUM': 'Health Rules & Alerting',
        }

    downgraded_counts = []
    for col in area_cols:
        if col in df_cmp.columns:
            s = df_cmp[col].astype(str)
            cnt = s.str.contains('Downgraded', case=False, na=False).sum()
            downgraded_counts.append((col, int(cnt)))
    downgraded_counts.sort(key=lambda x: x[1], reverse=True)
    focus_list = [pretty[c] for c, n in downgraded_counts if n > 0][:2]
    next_focus_text = ", ".join(focus_list) if focus_list else "Maintain current progress"

    def _fmt_outcome(prev, curr, delta):
        if prev is None or curr is None:
            return "Data not available."
        sign = "+" if delta is not None and delta > 0 else "−" if delta is not None and delta < 0 else "±"
        if delta is None:
            return f"{prev:.1f}%→{curr:.1f}%"
        return f"{prev:.1f}%→{curr:.1f}% ({sign}{abs(delta):.1f} pp)."

    headers = [
        "AppD Maturity Progression & Engagement",
        "Commentary",
        "Outcomes",
        "Change/Status Since Last",
    ]

    rows = [
        [
            "B/S/G/P Model Adoption & Maturity Status",
            f"B/S/G/P model applied to {domain}; assessment covered {int(total_curr)} apps.",
            cov_outcome,
            cov_arrow,
        ],
        [
            "Gold Status Apps",
            f"Gold-or-better coverage {_trend_word(curr_gold, prev_gold)} across the portfolio.",
            _fmt_outcome(prev_gold, curr_gold, delta_gold),
            _arrow(curr_gold, prev_gold),
        ],
        [
            "Platinum Status Apps",
            f"Platinum presence {_trend_word(curr_plat, prev_plat)}; teams progressing on prerequisites.",
            _fmt_outcome(prev_plat, curr_plat, delta_plat),
            _arrow(curr_plat, prev_plat),
        ],
        [
            "Maturity Partnership",
            "Working cadence in place; recommendations implemented during this period.",
            f"Overall result: {overall_result_text}. Next focus: {next_focus_text}.",
            "↑" if overall_result_text == "Increase" else "↓" if overall_result_text == "Decrease" else "→",
        ],
    ]

    # Insert/update table
    key_callouts_ph = find_table_placeholder_by_name(slide, table_placeholder_name)
    if key_callouts_ph:
        table = insert_table_at_placeholder(slide, table_placeholder_name, len(rows) + 1, len(headers))
    else:
        table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.6), Inches(2.1), Inches(9.0), Inches(4.0)).table

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(12)
        cell.text_frame.word_wrap = False
    autosize_col_to_header(table, 3, header_pt=12, padding_in=0.6, avg_char_em=0.55)

    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            if c_idx == 3 and value in ("↑", "↓", "→"):
                set_arrow_cell(cell, value, color=PINK, size_pt=36)
            else:
                cell.text = str(value)
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(12)

    # Badge color
    tier = overall_maturity_from_df(df_curr_analysis, grade_func=_grade_token_local)
    if tier:
        color_oval_for_maturity(slide, shape_name=oval_name, tier=tier, update_text=False)

    # Notes
    def _tier_counts(df):
        counts = {'bronze': 0, 'silver': 0, 'gold': 0, 'platinum': 0}
        col = 'OverallAssessment'
        if df is None or col not in df.columns:
            return counts, 0
        for v in df[col]:
            t = _grade_token_local(v)
            if t in counts:
                counts[t] += 1
        total = sum(counts.values())
        return counts, total

    def _pct(n, d): return (n / d) * 100.0 if d else 0.0

    tier_counts, tier_total = _tier_counts(df_curr_analysis)
    b, s, g, p = (tier_counts['bronze'], tier_counts['silver'], tier_counts['gold'], tier_counts['platinum'])
    pb, ps, pg, pp_ = (_pct(b, tier_total), _pct(s, tier_total), _pct(g, tier_total), _pct(p, tier_total))

    rationale = (
        f"Status is {tier} because it has the largest share of rated apps this run. "
        f"Distribution — Platinum {pp_:.1f}% ({p}), Gold {pg:.1f}% ({g}), "
        f"Silver {ps:.1f}% ({s}), Bronze {pb:.1f}% ({b})."
    )
    coverage_note = f"Rated coverage: {cov_curr:.1f}% ({rated_curr}/{total_curr})." if cov_curr is not None and total_curr > 0 else "Rated coverage: n/a."

    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.clear()
    tf.paragraphs[0].text = "Overall tier selection: majority of app ratings in Analysis; ties prefer the higher tier."
    p2 = tf.add_paragraph(); p2.text = rationale
    p3 = tf.add_paragraph(); p3.text = f"{coverage_note} Next focus: {next_focus_text}."


# ---------------------------------------------------------------------------------------------

def generate_powerpoint_from_analysis(comparison_result_path, powerpoint_output_path, current_file_path, previous_file_path):
    logging.debug("Generating PowerPoint presentation...")

    try:
        # prs = Presentation(template_path)  # Open the template

        # Define the relative path for the template using the TEMPLATE_FOLDER
        template_folder = config.get('TEMPLATE_FOLDER', 'templates')  # 'templates' is the default folder name
        template_path = os.path.join(template_folder, 'template.pptx')

        # Load the 'Analysis' sheet from the current workbook (uploaded by the user)
        df_current_analysis = pd.read_excel(current_file_path, sheet_name='Analysis')

        # Count the number of valid applications by counting rows where 'name' column is not NaN or empty
        number_of_apps = df_current_analysis['name'].dropna().str.strip().ne('').sum()

        # Log the number of valid applications
        logging.info(f"Number of applications in the current 'Analysis' sheet: {number_of_apps}")

        # Check if the template exists, otherwise, ask the user for input or use environment variables
        if not os.path.exists(template_path):
            template_path = os.getenv('TEMPLATE_PATH', template_path)  # Allow user to set this via an environment variable
            if not os.path.exists(template_path):
                template_path = input("Template not found! Please provide the full path to the template: ")

        # Load the template
        prs = Presentation(template_path)
        logging.debug(f"Template loaded from: {template_path}")

        # Load Summary sheets (current and previous) to drive Key Callouts on Slide 2.
        current_summary_df = pd.read_excel(current_file_path, sheet_name='Summary')
        previous_summary_df = pd.read_excel(previous_file_path, sheet_name='Summary')

        # Load the Summary sheet
        summary_df = pd.read_excel(comparison_result_path, sheet_name='Summary')
        logging.debug("Loaded Summary sheet successfully.")
        logging.debug(f"Summary DataFrame head:\n{summary_df.head()}")

        # Load Summary sheets to drive the Key Callouts slide.
        current_summary_df = pd.read_excel(current_file_path, sheet_name='Summary')
        previous_summary_df = pd.read_excel(previous_file_path, sheet_name='Summary')

        # Load the Analysis sheet
        df_analysis = pd.read_excel(comparison_result_path, sheet_name='Analysis')
        # Load the 'AppAgentsAPM' sheet from the Excel file
        df_app_agents = pd.read_excel(comparison_result_path, sheet_name='AppAgentsAPM')
        # Load the 'MachineAgentsAPM' sheet from the Excel file
        df_machine_agents = pd.read_excel(comparison_result_path, sheet_name='MachineAgentsAPM')
        # Load the 'BusinessTransactionsAPM' sheet from the Excel file
        df_BTs = pd.read_excel(comparison_result_path, sheet_name='BusinessTransactionsAPM')
        # Load the 'BackendsAPM' sheet from the Excel file
        df_Backends = pd.read_excel(comparison_result_path, sheet_name='BackendsAPM')
        # Load the 'OverheadAPM' sheet from the Excel file
        df_Overhead = pd.read_excel(comparison_result_path, sheet_name='OverheadAPM')
        # Load the 'ServiceEndpointsAPM' sheet from the Excel file
        df_ServiceEndpoints = pd.read_excel(comparison_result_path, sheet_name='ServiceEndpointsAPM')
        # Load the 'ErrorConfigurationAPM' sheet from the Excel file
        df_ErrorConfiguration = pd.read_excel(comparison_result_path, sheet_name='ErrorConfigurationAPM')
        # Load the 'HealthRulesAndAlertingAPM' sheet from the Excel file
        df_HealthRulesAndAlerting = pd.read_excel(comparison_result_path, sheet_name='HealthRulesAndAlertingAPM')
        # Load the 'DataCollectorsAPM' sheet from the Excel file
        df_DataCollectors = pd.read_excel(comparison_result_path, sheet_name='DataCollectorsAPM')
        # Load the 'DashboardsAPM' sheet from the Excel file
        df_Dashboards = pd.read_excel(comparison_result_path, sheet_name='DashboardsAPM')

        # Function to find table placeholders by name
        def find_table_placeholder_by_name(slide, name):
            for shape in slide.shapes:
                if shape.is_placeholder and shape.name == name:
                    return shape
            return None  # Return None if not found

        def insert_table_at_placeholder(slide, placeholder_name, rows, cols):
            """Insert a table at the position of a placeholder."""
            placeholder = find_table_placeholder_by_name(slide, placeholder_name)
            
            if not placeholder:
                logging.error(f"Placeholder '{placeholder_name}' not found on the slide.")
                return None

            # Get placeholder dimensions
            left = placeholder.left
            top = placeholder.top
            width = placeholder.width
            height = placeholder.height

            logging.debug(f"Inserting table at placeholder position: left={left}, top={top}, width={width}, height={height}")

            # Insert table at the placeholder's position
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            return table_shape.table  # Return the inserted table

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Slide 2 (index 1) — Assessment Result - Key Callouts **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 

        slide = prs.slides[1]  # Slide 2

        # Existing local helpers to parse percentages and arrows.
        def _parse_percent_to_float(val):
            if pd.isna(val):
                return None
            if isinstance(val, (int, float)):
                return float(val)
            s = str(val).strip()
            if '→' in s:
                s = s.split('→')[-1].strip()
            s = s.replace('%', '')
            try:
                return float(s)
            except Exception:
                return None

        def _get_tier_percent(df, tier):
            name_map = {c.lower(): c for c in df.columns}
            candidates = [f"{tier.lower()} %", f"{tier.lower()}%", f"percentage{tier.lower()}", f"{tier.lower()}percentage"]
            for cand in candidates:
                if cand in name_map:
                    return _parse_percent_to_float(df[name_map[cand]].iloc[0])
            needed = ['bronze', 'silver', 'gold', 'platinum']
            if all(k in name_map for k in needed):
                try:
                    total = 0.0
                    counts = {}
                    for k in needed:
                        val = pd.to_numeric(df[name_map[k]].iloc[0], errors='coerce')
                        counts[k] = 0.0 if pd.isna(val) else float(val)
                        total += counts[k]
                    if total > 0:
                        return (counts[tier.lower()] / total) * 100.0
                except Exception:
                    return None
            return None

        def _arrow(curr, prev):
            if curr is None or prev is None:
                return '→'
            if curr > prev:
                return '↑'
            if curr < prev:
                return '↓'
            return '→'

        def _trend_word(curr, prev):
            if curr is None or prev is None:
                return "held steady"
            if curr > prev:
                return "increased"
            if curr < prev:
                return "decreased"
            return "held steady"

        def _delta_pp(curr, prev):
            if curr is None or prev is None:
                return None
            return round(curr - prev, 1)

        # NEW: Coverage helpers for APM adoption (local grade_token here; module-level functions consume it).
        def _grade_token(s: str):
            if not s:
                return None
            m = re.search(r'(platinum|gold|silver|bronze)', str(s), re.I)
            return m.group(1).lower() if m else None

        def _apps_coverage(path):
            try:
                df = pd.read_excel(path, sheet_name='Analysis')
                total = int(df['name'].dropna().astype(str).str.strip().ne('').sum())
                if total == 0:
                    return (0, 0, 0.0)
                rated = int(df['OverallAssessment'].apply(_grade_token).notna().sum())
                pct = (rated / total) * 100.0
                return (total, rated, pct)
            except Exception:
                return (0, 0, None)

        def _arrow_threshold(curr, prev, threshold_pp=0.5):
            if curr is None or prev is None:
                return '→'
            delta = curr - prev
            if delta >= threshold_pp:
                return '↑'
            if delta <= -threshold_pp:
                return '↓'
            return '→'

        def _fmt_pp_delta(prev, curr):
            if prev is None or curr is None:
                return None
            d = curr - prev
            sign = '+' if d > 0 else '−' if d < 0 else '±'
            return f"{prev:.1f}%→{curr:.1f}% ({sign}{abs(d):.1f} pp)."

        # Derive Gold/Platinum percentages from Summary sheets.
        curr_gold = _get_tier_percent(current_summary_df, 'Gold')
        prev_gold = _get_tier_percent(previous_summary_df, 'Gold')
        curr_plat = _get_tier_percent(current_summary_df, 'Platinum')
        prev_plat = _get_tier_percent(previous_summary_df, 'Platinum')

        # NEW: Coverage from the original APM workbooks (for the first row).
        total_prev, rated_prev, cov_prev = _apps_coverage(previous_file_path)
        total_curr, rated_curr, cov_curr = _apps_coverage(current_file_path)
        cov_arrow = _arrow_threshold(cov_curr, cov_prev)
        cov_outcome = (
            f"Coverage: {cov_curr:.1f}% of apps rated ({rated_curr}/{total_curr})."
            if cov_curr is not None and total_curr > 0
            else "Coverage data not available."
        )
        cov_prev_curr = _fmt_pp_delta(cov_prev, cov_curr)
        if cov_prev_curr:
            cov_outcome = f"{cov_outcome} {cov_prev_curr}"

        # Overall result and next focus (unchanged from your prior logic).
        try:
            df_cmp = pd.read_excel(comparison_result_path, sheet_name='Analysis')
        except Exception:
            df_cmp = None

        def _count_changes(df, col):
            if df is None or col not in df.columns:
                return 0, 0
            s = df[col].astype(str)
            up = s.str.contains('Upgraded', case=False, na=False).sum()
            down = s.str.contains('Downgraded', case=False, na=False).sum()
            return int(up), int(down)

        up_overall, down_overall = _count_changes(df_cmp, 'OverallAssessment')
        overall_result_text = "Increase" if up_overall > down_overall else "Decrease" if down_overall > up_overall else "Even"

        area_cols = [
            'AppAgentsAPM','MachineAgentsAPM','BusinessTransactionsAPM',
            'BackendsAPM','OverheadAPM','ServiceEndpointsAPM',
            'ErrorConfigurationAPM','HealthRulesAndAlertingAPM',
            'DataCollectorsAPM','DashboardsAPM'
        ]
        pretty = {
            'AppAgentsAPM': 'App Agents',
            'MachineAgentsAPM': 'Machine Agents',
            'BusinessTransactionsAPM': 'Business Transactions',
            'BackendsAPM': 'Backends',
            'OverheadAPM': 'Overhead',
            'ServiceEndpointsAPM': 'Service Endpoints',
            'ErrorConfigurationAPM': 'Error Configuration',
            'HealthRulesAndAlertingAPM': 'Health Rules & Alerting',
            'DataCollectorsAPM': 'Data Collectors',
            'DashboardsAPM': 'Dashboards',
        }
        downgraded_counts = []
        if df_cmp is not None:
            for col in area_cols:
                if col in df_cmp.columns:
                    s = df_cmp[col].astype(str)
                    cnt = s.str.contains('Downgraded', case=False, na=False).sum()
                    downgraded_counts.append((col, int(cnt)))
        downgraded_counts.sort(key=lambda x: x[1], reverse=True)
        focus_list = [pretty[c] for c, n in downgraded_counts if n > 0][:2]
        next_focus_text = ", ".join(focus_list) if focus_list else "Maintain current progress"

        delta_gold = _delta_pp(curr_gold, prev_gold)
        delta_plat = _delta_pp(curr_plat, prev_plat)

        def _fmt_outcome(prev, curr, delta):
            if prev is None or curr is None:
                return "Data not available."
            sign = "+" if delta is not None and delta > 0 else "−" if delta is not None and delta < 0 else "±"
            if delta is None:
                return f"{prev:.1f}%→{curr:.1f}%"
            return f"{prev:.1f}%→{curr:.1f}% ({sign}{abs(delta):.1f} pp)."

        # Table headers and rows.
        headers = [
            "AppD Maturity Progression & Engagement",
            "Commentary",
            "Outcomes",
            "Change/Status Since Last",
        ]

        rows = [
            # UPDATED first row: APM-only coverage.
            [
                "B/S/G/P Model Adoption & Maturity Status",
                f"B/S/G/P model applied to APM; assessment covered {int(total_curr)} apps.",
                cov_outcome,
                cov_arrow,
            ],
            [
                "Gold Status Apps",
                f"Gold-or-better coverage {_trend_word(curr_gold, prev_gold)} across the portfolio.",
                _fmt_outcome(prev_gold, curr_gold, delta_gold),
                _arrow(curr_gold, prev_gold),
            ],
            [
                "Platinum Status Apps",
                f"Platinum presence {_trend_word(curr_plat, prev_plat)}; teams progressing on prerequisites.",
                _fmt_outcome(prev_plat, curr_plat, delta_plat),
                _arrow(curr_plat, prev_plat),
            ],
            [
                "Maturity Partnership",
                "Working cadence in place; recommendations implemented during this period.",
                f"Overall result: {overall_result_text}. Next focus: {next_focus_text}.",
                "↑" if overall_result_text == "Increase" else "↓" if overall_result_text == "Decrease" else "→",
            ],
        ]

        # Insert table, set header styles, autosize arrow column, and render rows.
        key_callouts_ph = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if key_callouts_ph:
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(rows) + 1, len(headers))
        else:
            table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.6), Inches(2.1), Inches(9.0), Inches(4.0)).table

        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(12)
            cell.text_frame.word_wrap = False

        autosize_col_to_header(table, 3, header_pt=12, padding_in=0.6, avg_char_em=0.55)

        for r_idx, row in enumerate(rows, start=1):
            for c_idx, value in enumerate(row):
                cell = table.cell(r_idx, c_idx)
                if c_idx == 3 and value in ("↑", "↓", "→"):
                    set_arrow_cell(cell, value, color=PINK, size_pt=36)
                else:
                    cell.text = str(value)
                    p = cell.text_frame.paragraphs[0]
                    p.font.size = Pt(12)

        # NEW: color the oval named "Oval 10" to reflect overall maturity (APM current Analysis).
        overall_tier = overall_maturity_from_df(df_current_analysis, grade_func=_grade_token)
        if overall_tier:
            color_oval_for_maturity(slide, shape_name="Oval 10", tier=overall_tier, update_text=False)

    # --- Add notes explaining why the status is Silver (or other tier) ---
        def _tier_counts(df):
            counts = {'bronze': 0, 'silver': 0, 'gold': 0, 'platinum': 0}
            col = 'OverallAssessment'
            if df is None or col not in df.columns:
                return counts, 0
            for v in df[col]:
                t = _grade_token(v)
                if t in counts:
                    counts[t] += 1
            total = sum(counts.values())
            return counts, total

        def _pct(n, d):
            return (n / d) * 100.0 if d else 0.0

        tier_counts, tier_total = _tier_counts(df_current_analysis)
        b, s, g, p = (tier_counts['bronze'], tier_counts['silver'], tier_counts['gold'], tier_counts['platinum'])
        pb, ps, pg, pp_ = (_pct(b, tier_total), _pct(s, tier_total), _pct(g, tier_total), _pct(p, tier_total))

        # Build a clear, single-paragraph rationale line.
        rationale = (
            f"Status is {overall_tier} because it has the largest share of rated apps this run. "
            f"Distribution — Platinum {pp_:.1f}% ({p}), Gold {pg:.1f}% ({g}), "
            f"Silver {ps:.1f}% ({s}), Bronze {pb:.1f}% ({b})."
        )

        coverage_note = (
            f"Rated coverage: {cov_curr:.1f}% ({rated_curr}/{total_curr})."
            if cov_curr is not None and total_curr > 0 else "Rated coverage: n/a."
        )

        next_focus_note = f"Next focus: {next_focus_text}."

        # Write to slide notes.
        notes = slide.notes_slide  # creates one if it doesn't exist
        tf = notes.notes_text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = "Overall tier selection: majority of app ratings in Analysis; ties prefer the higher tier."
        p2 = tf.add_paragraph()
        p2.text = rationale
        p3 = tf.add_paragraph()
        p3.text = f"{coverage_note} {next_focus_note}"
        
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Now handle Slide 4 table with "Upgraded" applications **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        slide = prs.slides[3]  # Slide 4 (index 3)
        upgraded_apps = df_analysis[df_analysis['OverallAssessment'].str.contains('upgraded', case=False, na=False)]['name'].tolist()

        # Count the number of applications in the current "Analysis" sheet
        current_analysis_df = pd.read_excel(current_file_path, sheet_name='Analysis')  # Load the current "Analysis" sheet
        number_of_apps = len(current_analysis_df)

        # Insert the count into TextBox 7
        textbox_7 = None
        for shape in slide.shapes:
            if shape.name == "TextBox 7":
                textbox_7 = shape
                break

        if textbox_7:
            textbox_7.text = f"{number_of_apps}"  # Set the text with the count
        else:
            logging.warning("TextBox 8 not found on Slide 3.")

        # Insert Upgraded Applications Table onto Slide 3 (Slide index 2) - using Table Placeholder 1
        upgraded_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")  # We are now using the same placeholder
        if upgraded_placeholder:
            logging.debug("Found Upgraded Applications table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(upgraded_apps) + 1, 1)
        else:
            logging.warning("Upgraded Applications table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(upgraded_apps) + 1, 1, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  

        # Add header for the new table
        table.cell(0, 0).text = "Applications with Upgraded Metrics"
        table.cell(0, 0).text_frame.paragraphs[0].font.size = Pt(12)

        # Populate the table with upgraded applications
        for idx, app in enumerate(upgraded_apps):
            table.cell(idx + 1, 0).text = app
            table.cell(idx + 1, 0).text_frame.paragraphs[0].font.size = Pt(12)

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert Summary Table onto Slide 5 (Slide index 4) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*
        slide = prs.slides[4]  # Slide 5 (index 4)
        summary_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")  # Placeholder for Summary Table

        if summary_placeholder:
            logging.debug("Found Summary table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(summary_df) + 1, len(summary_df.columns))
        else:
            logging.warning("Summary table placeholder not found. Adding manually.")
            # Explicitly add a new table with defined dimensions for Slide 5
            table = slide.shapes.add_table(len(summary_df) + 1, len(summary_df.columns), Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  

        # Set column headers for the Summary table
        for col_idx, column in enumerate(summary_df.columns):
            table.cell(0, col_idx).text = str(column)
            table.cell(0, col_idx).text_frame.paragraphs[0].font.size = Pt(12)

        # Populate table with Summary data
        for row_idx, row in summary_df.iterrows():
            for col_idx, value in enumerate(row):
                table.cell(row_idx + 1, col_idx).text = str(value)
                table.cell(row_idx + 1, col_idx).text_frame.paragraphs[0].font.size = Pt(12)

        # Add the title for Slide 4 (Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Comparison Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white


        # Load the Analysis sheet
        df = pd.read_excel(comparison_result_path, sheet_name='Analysis')

        columns = [
            'AppAgentsAPM', 'MachineAgentsAPM', 'BusinessTransactionsAPM',
            'BackendsAPM', 'OverheadAPM', 'ServiceEndpointsAPM',
            'ErrorConfigurationAPM', 'HealthRulesAndAlertingAPM',
            'DataCollectorsAPM', 'DashboardsAPM', 'OverallAssessment'
        ]

        results = {}
        total_applications = len(df)
        
        for col in columns:
            df[col] = df[col].astype(str)
            upgraded_count = df[col].str.contains('upgraded', case=False, na=False).sum()
            downgraded_count = df[col].str.contains('downgraded', case=False, na=False).sum()

            # Total applications is the length of the column
            total_applications = len(df[col])

            overall_result = "Increase" if upgraded_count > downgraded_count else "Decrease" if downgraded_count > upgraded_count else "Even"
            percentage_value = 0 if overall_result == "Even" else round((upgraded_count / total_applications) * 100)

            # Log the results for each column
            # logging.debug(f"Column: {col}")
            # logging.debug(f"Upgraded Count: {upgraded_count}")
            # logging.debug(f"Total Applications: {total_applications}")
            # logging.debug(f"Overall Result: {overall_result}")
            # logging.debug(f"Percentage: {percentage_value}%")

            results[col] = {
                'upgraded': upgraded_count,
                'downgraded': downgraded_count,
                'overall_result': overall_result,
                'percentage': percentage_value
            }

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert Summary Table from Previous Workbook onto Slide 5 (Table Placeholder 4) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        slide = prs.slides[4]  # Slide 5 (index 4)

        # Load the previous summary data
        previous_summary_df = pd.read_excel(previous_file_path, sheet_name='Summary')

        # Add to Table Placeholder 4 (for previous summary)
        summary_placeholder_previous = find_table_placeholder_by_name(slide, "Table Placeholder 4")  # Placeholder for Previous Summary Table
        if summary_placeholder_previous:
            logging.debug("Found Table Placeholder 4. Inserting table for previous summary.")
            table_previous = insert_table_at_placeholder(slide, "Table Placeholder 4", len(previous_summary_df) + 1, len(previous_summary_df.columns))
        else:
            logging.warning("Table Placeholder 4 not found. Adding manually.")
            table_previous = slide.shapes.add_table(len(previous_summary_df) + 1, len(previous_summary_df.columns), Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  

        # Set column headers for the previous summary table
        for col_idx, column in enumerate(previous_summary_df.columns):
            table_previous.cell(0, col_idx).text = str(column)
            table_previous.cell(0, col_idx).text_frame.paragraphs[0].font.size = Pt(12)

        # Populate the table with previous summary data
        for row_idx, row in previous_summary_df.iterrows():
            for col_idx, value in enumerate(row):
                table_previous.cell(row_idx + 1, col_idx).text = str(value)
                table_previous.cell(row_idx + 1, col_idx).text_frame.paragraphs[0].font.size = Pt(12)

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert Summary Table from Current Workbook onto Slide 5 (Table Placeholder 4)
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # Load the current summary data
        current_summary_df = pd.read_excel(current_file_path, sheet_name='Summary')

        # Add to Table Placeholder 3 (for current summary)
        summary_placeholder_current = find_table_placeholder_by_name(slide, "Table Placeholder 3")  # Placeholder for Current Summary Table
        if summary_placeholder_current:
            logging.debug("Found Table Placeholder 3. Inserting table for current summary.")
            table_current = insert_table_at_placeholder(slide, "Table Placeholder 3", len(current_summary_df) + 1, len(current_summary_df.columns))
        else:
            logging.warning("Table Placeholder 3 not found. Adding manually.")
            table_current = slide.shapes.add_table(len(current_summary_df) + 1, len(current_summary_df.columns), Inches(0.5), Inches(6), Inches(9), Inches(4)).table  

        # Set column headers for the current summary table
        for col_idx, column in enumerate(current_summary_df.columns):
            table_current.cell(0, col_idx).text = str(column)
            table_current.cell(0, col_idx).text_frame.paragraphs[0].font.size = Pt(12)

        # Populate the table with current summary data
        for row_idx, row in current_summary_df.iterrows():
            for col_idx, value in enumerate(row):
                table_current.cell(row_idx + 1, col_idx).text = str(value)
                table_current.cell(row_idx + 1, col_idx).text_frame.paragraphs[0].font.size = Pt(12)

        # Add the title for Slide 4 (Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Comparison Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white


        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*
        # Insert Overall Assessment Table onto Slide 7 (Slide index 6)
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*
        slide = prs.slides[6]  # Slide 7
        overall_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")  # Placeholder name

        if overall_placeholder:
            # logging.debug("Found Overall Assessment table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", 2, 5)
        else:
            # logging.warning("Overall Assessment table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(2, 5, Inches(0.5), Inches(1.5), Inches(9), Inches(1.5)).table  

        headers = ['Metric', '# of Apps Improved', '# Apps Degraded', 'Overall Result', 'Percentage Value']
        for col_idx, header in enumerate(headers):
            table.cell(0, col_idx).text = header
            table.cell(0, col_idx).text_frame.paragraphs[0].font.size = Pt(14)

        overall_assessment = results['OverallAssessment']
        table.cell(1, 0).text = 'OverallAssessment'
        table.cell(1, 1).text = str(overall_assessment['upgraded'])
        table.cell(1, 2).text = str(overall_assessment['downgraded'])
        table.cell(1, 3).text = overall_assessment['overall_result']
        table.cell(1, 4).text = f"{overall_assessment['percentage']}%"

        if overall_assessment['overall_result'] == "Increase":
            table.cell(1, 4).fill.solid()
            table.cell(1, 4).fill.fore_color.rgb = RGBColor(0, 255, 0)  # Green

        # Add the title for Slide 7
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Overall Assessment Result"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*
        # Insert Status Table onto Slide 8 (Slide index 7)
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*
        slide = prs.slides[7]  # Slide 8 (index 7)
        status_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")  # Placeholder name

        num_rows = len(columns)  # Should match the expected row count
        num_cols = 5  

        if status_placeholder:
            # logging.debug("Found Status table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", num_rows, num_cols)
        else:
            # logging.warning("Status table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  

        headers = ['Metric', '# of Apps Improved', '# Apps Degraded', 'Overall Result', 'Percentage Value']
        for col_idx, header in enumerate(headers):
            table.cell(0, col_idx).text = header
            table.cell(0, col_idx).text_frame.paragraphs[0].font.size = Pt(14)

        for i, col in enumerate(columns[:-1]):  
            table.cell(i + 1, 0).text = col
            table.cell(i + 1, 1).text = str(results[col]['upgraded'])
            table.cell(i + 1, 2).text = str(results[col]['downgraded'])
            table.cell(i + 1, 3).text = results[col]['overall_result']
            table.cell(i + 1, 4).text = f"{results[col]['percentage']}%"

            if results[col]['overall_result'] == "Increase":
                table.cell(i + 1, 4).fill.solid()
                table.cell(i + 1, 4).fill.fore_color.rgb = RGBColor(0, 255, 0)  # Green

        # Add the title for Slide 8
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "APM Maturity Assessment Result"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert APM AGENT Downgrade Table onto Slide 12 (Slide index 11) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[11]  # Slide 12 (index 11)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the AppAgentsAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['AppAgentsAPM']
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 12 (Slide index 11)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            logging.debug(f"Grade: {grade}, Applications: {', '.join(downgrade_data[grade]['applications'])}")

            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = ', '.join(downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 12 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "APM Agent - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'percentAgentsLessThan1YearOld': 'Rectangle 11',
            'metricLimitNotHit': 'Rectangle 10',
            'percentAgentsLessThan2YearsOld': 'Rectangle 12',
            'percentAgentsReportingData': 'Rectangle 13',
            'percentAgentsRunningSameVersion': 'Rectangle 14'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_app_agents.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'declined' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        # for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 12
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 15.") 


        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert MACHINE AGENT Downgrade Table onto Slide 13 (Slide index 12) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[12]  # Slide 13 (index 12)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the MachineAgentsAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['MachineAgentsAPM']  # Use MachineAgentsAPM column for comparison
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 13 (Slide index 12)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            logging.debug(f"Grade: {grade}, Applications: {', '.join(downgrade_data[grade]['applications'])}")
            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = ', '.join(downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 11 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Machine Agent - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'percentAgentsLessThan1YearOld': 'Rectangle 8',
            'percentAgentsLessThan2YearsOld': 'Rectangle 9',
            'percentAgentsReportingData': 'Rectangle 10',
            'percentAgentsRunningSameVersion': 'Rectangle 11',
            'percentAgentsInstalledAlongsideAppAgents': 'Rectangle 12'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_machine_agents.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'declined' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        # for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 13
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 13.") 

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert BT Downgrade Table onto Slide 14 (Slide index 13) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[13]  # Slide 14 (index 13)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the BusinessTransactionsAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['BusinessTransactionsAPM']  # Use BusinessTransactionsAPM column for comparison
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 14 (Slide index 13)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            logging.debug(f"Grade: {grade}, Applications: {', '.join(downgrade_data[grade]['applications'])}")
            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = ', '.join(downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 14 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Business Transactions - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'numberOfBTs': 'Rectangle 17',
            'percentBTsWithLoad': 'Rectangle 18',
            'btLockdownEnabled': 'Rectangle 19',
            'numberCustomMatchRules': 'Rectangle 20'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_BTs.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'decreased' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        # for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 14
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 14.") 

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert Backend Downgrade Table onto Slide 15 (Slide index 14) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[14]  # Slide 15 (index 14)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the BackendsAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['BackendsAPM']  # Use BackendsAPM column for comparison
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 15 (Slide index 14)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            # Convert all items in the applications list to strings
            applications_str = ', '.join(str(app) for app in downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            
            # Log the grade and applications
            logging.debug(f"Grade: {grade}, Applications: {applications_str}")
            
            # Populate the table with the data
            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = applications_str  # Display the application names
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 15 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Backends - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'percentBackendsWithLoad': 'Rectangle 10',
            'backendLimitNotHit': 'Rectangle 11',
            'numberOfCustomBackendRules': 'Rectangle 12'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_Backends.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'decreased' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        # for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 15
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 15.") 

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert SEP Downgrade Table onto Slide 16 (Slide index 15) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[15]  # Slide 16 (index 15)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the ServiceEndpointsAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['ServiceEndpointsAPM']  # Use ServiceEndpointsAPM column for comparison
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 16 (Slide index 15)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            logging.debug(f"Grade: {grade}, Applications: {', '.join(downgrade_data[grade]['applications'])}")
            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = ', '.join(downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 14 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Service Endpoints - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'numberOfCustomServiceEndpointRules': 'Rectangle 10',
            'serviceEndpointLimitNotHit': 'Rectangle 11',
            'percentServiceEndpointsWithLoadOrDisabled': 'Rectangle 12'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_ServiceEndpoints.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'decreased' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        #for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 16
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 14.") 

        
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert ERROR CONFIG Downgrade Table onto Slide 17 (Slide index 16) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[16]  # Slide 17 (index 16)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the ErrorConfigurationAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['ErrorConfigurationAPM']  # Use ErrorConfigurationAPM column for comparison
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 15 (Slide index 14)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            logging.debug(f"Grade: {grade}, Applications: {', '.join(downgrade_data[grade]['applications'])}")
            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = ', '.join(downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 17 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Error Configuration - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'successPercentageOfWorstTransaction': 'Rectangle 10',
            'numberOfCustomRules': 'Rectangle 11'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_ErrorConfiguration.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'decreased' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        #for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 17
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 15.") 

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert HR & ALERTS Downgrade Table onto Slide 18 (Slide index 17) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[17]  # Slide 18 (index 17)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the HealthRulesAndAlertingAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['HealthRulesAndAlertingAPM']  # Use HealthRulesAndAlertingAPM column for comparison
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 18 (Slide index 17)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            logging.debug(f"Grade: {grade}, Applications: {', '.join(downgrade_data[grade]['applications'])}")
            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = ', '.join(downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 18 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Health Rules & Alerting - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'numberOfHealthRuleViolations': 'Rectangle 10',
            'numberOfDefaultHealthRulesModified': 'Rectangle 11',
            'numberOfActionsBoundToEnabledPolicies': 'Rectangle 12',
            'numberOfCustomHealthRules': 'Rectangle 13'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_HealthRulesAndAlerting.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'decreased' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        #for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 18
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 16.") 

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert DATA COLLECTORS Downgrade Table onto Slide 19 (Slide index 18) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[18]  # Slide 19 (index 18)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the DataCollectorsAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['DataCollectorsAPM']  # Use DataCollectorsAPM column for comparison
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 19 (Slide index 18)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            logging.debug(f"Grade: {grade}, Applications: {', '.join(downgrade_data[grade]['applications'])}")
            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = ', '.join(downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 19 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Data Collectors - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'numberOfDataCollectorFieldsConfigured': 'Rectangle 10',
            'numberOfDataCollectorFieldsCollectedInSnapshots': 'Rectangle 11',
            'numberOfDataCollectorFieldsCollectedInAnalytics': 'Rectangle 12',
            'biqEnabled': 'Rectangle 13'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_DataCollectors.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'decreased' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        #for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 19
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 17.") 

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert DASHBOARDS Downgrade Table onto Slide 20 (Slide index 19) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[19]  # Slide 20 (index 19)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the DashboardsAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['DashboardsAPM']  # Use DashboardsAPM column for comparison
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 20 (Slide index 19)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            logging.debug(f"Grade: {grade}, Applications: {', '.join(downgrade_data[grade]['applications'])}")
            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = ', '.join(downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 20 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Dashboards - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

                # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'numberOfDashboards': 'Rectangle 10',
            'percentageOfDashboardsModifiedLast6Months': 'Rectangle 11',
            'numberOfDashboardsUsingBiQ': 'Rectangle 12'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_Dashboards.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'decreased' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        #for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 20
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 20.") 

        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-* 
        # ** Insert OVERHEAD Downgrade Table onto Slide 21 (Slide index 20) **
        # *-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*

        slide = prs.slides[20]  # Slide 21 (index 20)

        # Define grades for comparison, including 'platinum' for the downgrade logic
        all_grades = ['platinum', 'gold', 'silver', 'bronze']

        # Define grades for the table (no platinum included as requested)
        grades_for_table = ['gold', 'silver', 'bronze']
        downgrade_data = {grade: {'applications': [], 'number_of_apps': 0, 'percentage': 0} for grade in grades_for_table}

        # Iterate through the OverheadAPM column to detect downgrades
        for index, row in df_analysis.iterrows():
            current_value = row['OverheadAPM']  # Use OverheadAPM column for comparison
            app_name = row['name']  # Capture the application name

            # Log the application name and values for debugging
            logging.debug(f"Checking App: {app_name} - Current Value: {current_value}")

            # Check if the value contains a downgrade based on '→'
            if '→' in str(current_value):
                logging.debug(f"Found potential Downgrade in {app_name}: {current_value}")

                try:
                    # Split the value into previous and current grades based on '→'
                    previous_value, current_grade = current_value.split('→')
                    previous_value = previous_value.strip().lower()  # Ensure case-insensitive comparison
                    current_grade = current_grade.strip().lower().split(' ')[0]  # Get only the grade name

                    # Log the extracted values for debugging
                    logging.debug(f"Extracted: Previous Value: {previous_value}, Current Grade: {current_grade}")
                    logging.debug(f"App Name: {app_name}")

                    # Ensure that both grades are valid and are in the `all_grades` list
                    if previous_value in all_grades and current_grade in all_grades:
                        # Check if downgrade is valid (previous grade > current grade)
                        if all_grades.index(previous_value) < all_grades.index(current_grade):
                            logging.debug(f"Adding {app_name} to {current_grade} downgrade list")
                            downgrade_data[current_grade]['applications'].append(app_name)
                            downgrade_data[current_grade]['number_of_apps'] += 1  # Increment the number of applications
                            logging.debug(f"Current Downgrade Data: {downgrade_data}")
                        else:
                            logging.debug(f"Not a downgrade for {app_name}: {previous_value} → {current_grade}")
                    else:
                        logging.debug(f"Invalid grades detected for downgrade comparison: {previous_value}, {current_grade}")
                except Exception as e:
                    logging.error(f"Error processing downgrade for {app_name}: {e}")
            else:
                logging.debug(f"No Downgrade for App: {app_name} - Current Value: {current_value}")

        # Log applications for each grade to check population
        for grade in grades_for_table:
            logging.debug(f"Applications for {grade}: {downgrade_data[grade]['applications']}")

        # Calculate the percentage of downgrades for each grade
        total_apps = len(df_analysis)
        for grade in grades_for_table:
            downgrade_data[grade]['percentage'] = len(downgrade_data[grade]['applications']) / total_apps * 100

        # Log the percentage of downgrades
        logging.debug(f"Downgrade Percentages: {downgrade_data}")

        # Insert Downgrade Summary Table onto Slide 21 (Slide index 20)
        downgrade_placeholder = find_table_placeholder_by_name(slide, "Table Placeholder 1")
        if downgrade_placeholder:
            logging.debug("Found Downgrade table placeholder. Inserting table.")
            table = insert_table_at_placeholder(slide, "Table Placeholder 1", len(grades_for_table) + 1, 4)  # Increase columns to 4
        else:
            logging.warning("Downgrade table placeholder not found. Adding manually.")
            table = slide.shapes.add_table(len(grades_for_table) + 1, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table  # Increase columns to 4

        # Set table headers
        table.cell(0, 0).text = "Grade"
        table.cell(0, 1).text = "Application Names"
        table.cell(0, 2).text = "Number of Applications"
        table.cell(0, 3).text = "Percentage Declined"

        # Populate the table with downgrade data
        logging.debug(f"Populating table with downgrade data:")
        for i, grade in enumerate(grades_for_table):
            logging.debug(f"Grade: {grade}, Applications: {', '.join(downgrade_data[grade]['applications'])}")
            table.cell(i + 1, 0).text = grade.capitalize()  # Capitalize the grade names for display
            table.cell(i + 1, 1).text = ', '.join(downgrade_data[grade]['applications']) if downgrade_data[grade]['applications'] else "None"
            table.cell(i + 1, 2).text = str(downgrade_data[grade]['number_of_apps'])  # Number of Applications
            table.cell(i + 1, 3).text = f"{downgrade_data[grade]['percentage']:.2f}%"

        # Add the title for Slide 20 (Downgrade Summary Slide)
        title_placeholder = find_table_placeholder_by_name(slide, "Title 2")
        if title_placeholder:
            title_placeholder.text = "Overhead - Downgrade Summary"
            # Set text color to white
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

        # Create a dictionary of column names and their corresponding rectangle names
        columns_and_rectangles = {
            'developerModeNotEnabledForAnyBT': 'Rectangle 10',
            'findEntryPointsNotEnabled': 'Rectangle 11',
            'aggressiveSnapshottingNotEnabled': 'Rectangle 12',
            'developerModeNotEnabledForApplication': 'Rectangle 13'
        }

        # Initialize a dictionary to store 'Declined' counts for each column
        declined_counts = {key: 0 for key in columns_and_rectangles}

        # Iterate through the rows and count 'Declined' in each relevant column
        for index, row in df_Overhead.iterrows():
            for column, rectangle in columns_and_rectangles.items():
                # If 'Declined' is found in the current column (case insensitive)
                if 'changed' in str(row[column]).lower():
                    declined_counts[column] += 1

        # Log the counts for debugging
        #for column, count in declined_counts.items():
        #    logging.debug(f"Number of 'Declined' cells in {column}: {count}")

        # Insert the 'Declined' count into the corresponding rectangles on Slide 20
        for column, rectangle_name in columns_and_rectangles.items():
            # Find the rectangle shape by name
            rectangle = None
            for shape in slide.shapes:
                if shape.name == rectangle_name:
                    rectangle = shape
                    break
            
            # Update the text of the rectangle if found
            if rectangle:
                rectangle.text = f"{declined_counts[column]}"  # Set the text with the count
            else:
                logging.warning(f"{rectangle_name} not found on Slide 20.")

        # Save the PowerPoint
        prs.save(powerpoint_output_path)
        logging.debug(f"PowerPoint saved to {powerpoint_output_path}.")

    except Exception as e:
        logging.error(f"Error generating PowerPoint: {e}", exc_info=True)
        raise

    ################################################################################### 
    ######################           BRUM PPT                       ###################
    ###################################################################################

def generate_powerpoint_from_brum(comparison_result_path, powerpoint_output_path, current_file_path, previous_file_path):
    logging.debug("[BRUM] Generating PowerPoint presentation...")

    try:
        # Match APM path logic: TEMPLATE_FOLDER + filename, with the same fallback flow.
        template_folder = config.get('TEMPLATE_FOLDER', 'templates')
        template_path = os.path.join(template_folder, 'template_brum.pptx')

        # Load the 'Analysis' sheet from the current BRUM workbook.
        df_current_analysis = pd.read_excel(current_file_path, sheet_name='Analysis')

        # Count valid apps by non-empty 'name' values (same as APM).
        number_of_apps = df_current_analysis['name'].dropna().str.strip().ne('').sum()
        logging.info(f"[BRUM] Number of applications in the current 'Analysis' sheet: {number_of_apps}")

        # Same fallback behavior used by APM.
        if not os.path.exists(template_path):
            template_path = os.getenv('TEMPLATE_PATH', template_path)
            if not os.path.exists(template_path):
                template_path = input("Template not found! Please provide the full path to the template: ")

        # Load the template.
        prs = Presentation(template_path)
        logging.debug(f"[BRUM] Template loaded from: {template_path}")

        # Load Summary sheets (current and previous) for Key Callouts slide(s).
        current_summary_df = pd.read_excel(current_file_path, sheet_name='Summary')
        previous_summary_df = pd.read_excel(previous_file_path, sheet_name='Summary')

        # Load the Summary sheet from the comparison result (for comparison tables).
        summary_df = pd.read_excel(comparison_result_path, sheet_name='Summary')
        logging.debug("[BRUM] Loaded Summary sheet successfully.")
        logging.debug(f"[BRUM] Summary DataFrame head:\n{summary_df.head()}")

        # Load the Analysis sheet from the comparison result (annotations like Upgraded/Downgraded).
        df_analysis = pd.read_excel(comparison_result_path, sheet_name='Analysis')

        # Load BRUM-specific sheets from the comparison result.
        df_network_requests = pd.read_excel(comparison_result_path, sheet_name='NetworkRequestsBRUM')
        df_health_rules    = pd.read_excel(comparison_result_path, sheet_name='HealthRulesAndAlertingBRUM')
        df_overall_brum    = pd.read_excel(comparison_result_path, sheet_name='OverallAssessmentBRUM')

        # Placeholder helpers — identical pattern to APM.
        def find_table_placeholder_by_name(slide, name):
            for shape in slide.shapes:
                if shape.is_placeholder and shape.name == name:
                    return shape
            return None

        def insert_table_at_placeholder(slide, placeholder_name, rows, cols):
            placeholder = find_table_placeholder_by_name(slide, placeholder_name)
            if not placeholder:
                logging.error(f"[BRUM] Placeholder '{placeholder_name}' not found on the slide.")
                return None
            left, top, width, height = placeholder.left, placeholder.top, placeholder.width, placeholder.height
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            return table_shape.table

        # Utility helpers.
        def find_slides_with_placeholder(prs_obj, placeholder_name):
            return [s for s in prs_obj.slides if find_table_placeholder_by_name(s, placeholder_name)]

        def slide_title_text(slide):
            for shp in slide.shapes:
                if getattr(shp, "name", "") == "Title 2" and hasattr(shp, "text_frame") and shp.text_frame:
                    return shp.text_frame.text.strip()
            for shp in slide.shapes:
                if hasattr(shp, "text_frame") and shp.text_frame and shp.text_frame.text:
                    return shp.text_frame.text.strip()
            return ""

        # Track slides we have already populated using hashable IDs.
        used_slide_ids = set()

        def choose_slide_for_section(prefer_titles, required_placeholders=("Table Placeholder 1",), exclude_ids=None):
            exclude_ids = exclude_ids or set()
            # Try by title first.
            for s in prs.slides:
                if id(s) in exclude_ids:
                    continue
                title = slide_title_text(s)
                if any(pt.lower() in title.lower() for pt in prefer_titles):
                    if all(find_table_placeholder_by_name(s, p) for p in required_placeholders):
                        return s
            # Fallback: first slide with required placeholders not excluded.
            for s in prs.slides:
                if id(s) in exclude_ids:
                    continue
                if all(find_table_placeholder_by_name(s, p) for p in required_placeholders):
                    return s
            return None

        # ============================
        # BRUM Key Callouts (Table Placeholder 1)
        # ============================

        # Read previous/current OverallAssessmentBRUM to compute percentages.
        try:
            curr_overall_df = pd.read_excel(current_file_path, sheet_name='OverallAssessmentBRUM')
            prev_overall_df = pd.read_excel(previous_file_path, sheet_name='OverallAssessmentBRUM')
        except Exception:
            curr_overall_df = pd.DataFrame()
            prev_overall_df = pd.DataFrame()

        def last_percent(df, col):
            if df.empty or col not in df.columns:
                return None
            s = pd.to_numeric(df[col].astype(str).str.replace('%', ''), errors='coerce').dropna()
            return float(s.iloc[-1]) if not s.empty else None

        curr_gold  = last_percent(curr_overall_df, 'percentageTotalGoldOrBetter')
        prev_gold  = last_percent(prev_overall_df, 'percentageTotalGoldOrBetter')
        curr_plat  = last_percent(curr_overall_df, 'percentageTotalPlatinum')
        prev_plat  = last_percent(prev_overall_df, 'percentageTotalPlatinum')

        # Count upgrades/downgrades from Analysis for Overall and key BRUM areas.
        def count_changes(df, col):
            if col not in df.columns:
                return 0, 0
            s = df[col].astype(str)
            return (
                s.str.contains('Upgraded', case=False, na=False).sum(),
                s.str.contains('Downgraded', case=False, na=False).sum()
            )

        oa_up, oa_down   = count_changes(df_analysis, 'OverallAssessment')
        net_up, net_down = count_changes(df_analysis, 'NetworkRequestsBRUM')
        hra_up, hra_down = count_changes(df_analysis, 'HealthRulesAndAlertingBRUM')

        def arrow(curr, prev):
            if curr is None or prev is None:
                return '→'
            return '↑' if curr > prev else '↓' if curr < prev else '→'

        def fmt_change(prev, curr, suffix='%'):
            if prev is None or curr is None:
                return "Change observed."
            return (
                f"Increase from {prev:.1f}{suffix}→{curr:.1f}{suffix}"
                if curr > prev else
                f"Decrease from {prev:.1f}{suffix}→{curr:.1f}{suffix}"
                if curr < prev else
                f"No change ({curr:.1f}{suffix})."
            )

        slides_with_ph = find_slides_with_placeholder(prs, "Table Placeholder 1")
        key_callouts_slide = slides_with_ph[0] if slides_with_ph else (prs.slides[0] if len(prs.slides) else None)
        if key_callouts_slide:
            used_slide_ids.add(id(key_callouts_slide))

        if key_callouts_slide:
            headers = [
                "AppD Maturity Progression & Engagement",
                "Commentary",
                "Outcomes",
                "Change/Status Since Last"
            ]

            rows = [
                [
                    "B/S/G/P Model Adoption & Maturity Status (BRUM).",
                    f"BRUM analysis coverage across {int(number_of_apps)} applications.",
                    f"Overall BRUM upgrades: {oa_up}, downgrades: {oa_down}.",
                    '↑' if oa_up > oa_down else '↓' if oa_down > oa_up else '→'
                ],
                [
                    "Gold Status Apps (BRUM).",
                    "Change in gold-or-better coverage across the portfolio.",
                    fmt_change(prev_gold, curr_gold),
                    arrow(prev_gold, curr_gold)  # note: arrow(prev, curr) or arrow(curr, prev) per your utility
                ],
                [
                    "Platinum Status Apps (BRUM).",
                    "Top-tier BRUM maturity presence across applications.",
                    f"{curr_plat:.1f}% platinum." if curr_plat is not None else "Platinum presence observed.",
                    arrow(curr_plat, prev_plat)
                ],
                [
                    "Maturity Partnership (BRUM).",
                    "Ongoing improvements in BRUM instrumentation and alerting.",
                    f"Network Requests (↑{net_up}/↓{net_down}), Health Rules (↑{hra_up}/↓{hra_down}).",
                    '↑' if (net_up + hra_up) > (net_down + hra_down) else '↓' if (net_down + hra_down) > (net_up + hra_up) else '→'
                ]
            ]

            table = insert_table_at_placeholder(
                key_callouts_slide, "Table Placeholder 1", len(rows) + 1, len(headers)
            )
            if table:
                from pptx.util import Pt
                from pptx.dml.color import RGBColor

                # Header row
                for i, h in enumerate(headers):
                    cell = table.cell(0, i)
                    cell.text = h
                    p = cell.text_frame.paragraphs[0]
                    p.font.bold = True
                    p.font.size = Pt(12)

                # Make sure the status column is wide enough (optional)
                try:
                    autosize_col_to_header(table, 3, header_pt=12, padding_in=0.6, avg_char_em=0.55)
                except Exception:
                    pass

                # Data rows — use set_arrow_cell to render filled triangles like APM
                for r_idx, row in enumerate(rows, start=1):
                    for c_idx, value in enumerate(row):
                        cell = table.cell(r_idx, c_idx)
                        if c_idx == 3 and value in ('↑', '↓', '→'):
                            # Triangles with AppD pink (or pick green/red by direction if you prefer)
                            set_arrow_cell(cell, value, color=PINK, size_pt=36, font_name='Calibri')
                        else:
                            cell.text = str(value)
                            p = cell.text_frame.paragraphs[0]
                            p.font.size = Pt(12)
            else:
                logging.error("[BRUM] Could not insert Key Callouts table; placeholder missing.")
        else:
            logging.error("[BRUM] No slides available to place Key Callouts.")

        # --- BRUM: maturity badge + notes (self-contained, safe) ---
        try:
            # Local helpers (do not rely on module-level state).
            def _brum_grade_token(v):
                if not isinstance(v, str):
                    v = "" if pd.isna(v) else str(v)
                m = re.search(r'(platinum|gold|silver|bronze)', v, re.I)
                return m.group(1).lower() if m else None

            def _overall_maturity_from_df_brum(df):
                col = 'OverallAssessment'
                if df is None or col not in df.columns:
                    return None, {'bronze':0,'silver':0,'gold':0,'platinum':0}, 0
                counts = {'bronze':0,'silver':0,'gold':0,'platinum':0}
                rated = 0
                for v in df[col]:
                    t = _brum_grade_token(v)
                    if t in counts:
                        counts[t] += 1
                        rated += 1
                if rated == 0:
                    return None, counts, 0
                # Majority; tie breaks to higher maturity.
                rank = {'bronze':0,'silver':1,'gold':2,'platinum':3}
                best = max(counts.items(), key=lambda kv: (kv[1], rank[kv[0]]))
                return best[0].title(), counts, rated

            def _ideal_text_rgb_local(rgb):
                r, g, b = rgb[0], rgb[1], rgb[2]
                bright = (r*299 + g*587 + b*114)/1000.0
                return RGBColor(255,255,255) if bright < 140 else RGBColor(31,31,31)

            def _color_oval_for_maturity_local(slide_obj, shape_name, tier):
                palette = {
                    'Bronze':   RGBColor(205,127,50),
                    'Silver':   RGBColor(166,166,166),
                    'Gold':     RGBColor(255,192,0),
                    'Platinum': RGBColor(190,190,200),
                }
                if slide_obj is None or tier not in palette:
                    return
                target = next((sh for sh in slide_obj.shapes if getattr(sh, "name", "") == shape_name), None)
                if not target:
                    return
                target.fill.solid()
                target.fill.fore_color.rgb = palette[tier]
                if hasattr(target, "text_frame") and target.text_frame:
                    fg = _ideal_text_rgb_local(palette[tier])
                    for p in target.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.color.rgb = fg

            # Compute majority tier and coverage based on current Analysis.
            overall_tier, counts, rated = _overall_maturity_from_df_brum(df_current_analysis)
            if overall_tier and key_callouts_slide is not None:
                _color_oval_for_maturity_local(key_callouts_slide, "Oval 10", overall_tier)

                total_apps = int(number_of_apps) if number_of_apps else 0
                coverage = (rated / total_apps * 100.0) if total_apps else 0.0
                # Distribution over rated apps.
                def pct(n): return (n / rated * 100.0) if rated else 0.0
                pb, ps, pg, pp_ = (
                    pct(counts['bronze']), pct(counts['silver']),
                    pct(counts['gold']),   pct(counts['platinum'])
                )

                # Notes: explain rule + data that led to the tier.
                notes = key_callouts_slide.notes_slide
                tf = notes.notes_text_frame
                tf.clear()
                tf.paragraphs[0].text = "Overall tier selection: majority of app ratings in Analysis; ties prefer higher tier (Platinum > Gold > Silver > Bronze)."
                p2 = tf.add_paragraph()
                p2.text = (
                    f"Status is {overall_tier} based on rated distribution — "
                    f"Platinum {pp_:.1f}% ({counts['platinum']}), "
                    f"Gold {pg:.1f}% ({counts['gold']}), "
                    f"Silver {ps:.1f}% ({counts['silver']}), "
                    f"Bronze {pb:.1f}% ({counts['bronze']})."
                )
                p3 = tf.add_paragraph()
                p3.text = f"Rated coverage this run: {coverage:.1f}% ({rated}/{total_apps})."
        except Exception as e:
            logging.warning("[BRUM] Maturity badge/notes skipped: %s", e)

        # ============================
        # Applications Improved (Table Placeholder 1 on the next slide with same name)
        # ============================

        improved = []
        cols_map = [
            ('NetworkRequestsBRUM', 'Network Requests'),
            ('HealthRulesAndAlertingBRUM', 'Health Rules & Alerting'),
            ('OverallAssessment', 'Overall')
        ]

        for _, r in df_analysis.iterrows():
            app = str(r.get('name', '') or '').strip()
            if not app:
                continue
            areas = []
            for col, label in cols_map:
                if col in df_analysis.columns:
                    val = r.get(col, '')
                    if isinstance(val, str) and 'upgraded' in val.lower():
                        areas.append(label)
            if areas:
                improved.append((app, ', '.join(areas)))

        improved.sort(key=lambda x: x[0].lower())

        improved_slide = None
        for s in slides_with_ph:
            if id(s) not in used_slide_ids:
                improved_slide = s
                break
        if improved_slide:
            used_slide_ids.add(id(improved_slide))
            headers = ["Application", "Improvement Areas"]
            row_count = max(1, len(improved)) + 1
            table = insert_table_at_placeholder(improved_slide, "Table Placeholder 1", row_count, len(headers))
            if table:
                from pptx.util import Pt
                # Header
                for c, h in enumerate(headers):
                    table.cell(0, c).text = h
                    table.cell(0, c).text_frame.paragraphs[0].font.size = Pt(12)
                # Rows
                if improved:
                    for idx, (app, areas) in enumerate(improved, start=1):
                        table.cell(idx, 0).text = app
                        table.cell(idx, 1).text = areas
                        table.cell(idx, 0).text_frame.paragraphs[0].font.size = Pt(12)
                        table.cell(idx, 1).text_frame.paragraphs[0].font.size = Pt(12)
                else:
                    table.cell(1, 0).text = "No applications improved in this period."
                    table.cell(1, 1).text = ""
            else:
                logging.error("[BRUM] Could not insert Improved Applications table; placeholder missing.")
        else:
            logging.warning("[BRUM] No second slide with 'Table Placeholder 1' found for Improved Applications.")

        # ============================
        # Summary slide: Previous (Table Placeholder 4), Current (Table Placeholder 3), Comparison (Table Placeholder 1)
        # ============================

        def find_slide_with_all_placeholders(prs_obj, names):
            for s in prs_obj.slides:
                if all(find_table_placeholder_by_name(s, n) for n in names):
                    return s
            return None

        summary_slide = find_slide_with_all_placeholders(prs, ["Table Placeholder 1", "Table Placeholder 3", "Table Placeholder 4"])
        if not summary_slide:
            for s in prs.slides:
                if any(find_table_placeholder_by_name(s, n) for n in ["Table Placeholder 1", "Table Placeholder 3", "Table Placeholder 4"]):
                    summary_slide = s
                    break
        if summary_slide:
            used_slide_ids.add(id(summary_slide))

        def fill_table_from_df(table, df):
            from pptx.util import Pt
            # Header
            for c, col in enumerate(df.columns):
                table.cell(0, c).text = str(col)
                table.cell(0, c).text_frame.paragraphs[0].font.size = Pt(12)
            # Rows
            for r_idx, row in df.iterrows():
                for c_idx, val in enumerate(row):
                    table.cell(r_idx + 1, c_idx).text = "" if pd.isna(val) else str(val)
                    table.cell(r_idx + 1, c_idx).text_frame.paragraphs[0].font.size = Pt(12)

        if summary_slide:
            ph4 = find_table_placeholder_by_name(summary_slide, "Table Placeholder 4")
            if ph4:
                table_prev = insert_table_at_placeholder(summary_slide, "Table Placeholder 4", len(previous_summary_df) + 1, len(previous_summary_df.columns))
                if table_prev:
                    fill_table_from_df(table_prev, previous_summary_df)
            else:
                logging.warning("[BRUM] 'Table Placeholder 4' not found for Previous Summary.")

            ph3 = find_table_placeholder_by_name(summary_slide, "Table Placeholder 3")
            if ph3:
                table_curr = insert_table_at_placeholder(summary_slide, "Table Placeholder 3", len(current_summary_df) + 1, len(current_summary_df.columns))
                if table_curr:
                    fill_table_from_df(table_curr, current_summary_df)
            else:
                logging.warning("[BRUM] 'Table Placeholder 3' not found for Current Summary.")

            ph1 = find_table_placeholder_by_name(summary_slide, "Table Placeholder 1")
            if ph1:
                table_comp = insert_table_at_placeholder(summary_slide, "Table Placeholder 1", len(summary_df) + 1, len(summary_df.columns))
                if table_comp:
                    fill_table_from_df(table_comp, summary_df)
            else:
                logging.warning("[BRUM] 'Table Placeholder 1' not found for Comparison Summary.")
        else:
            logging.error("[BRUM] No suitable slide found for Summary tables.")

        # ============================
        # Overall Assessment slide (Table Placeholder 1)
        # ============================

        overall_slide = choose_slide_for_section(
            prefer_titles=["Overall Assessment", "Overall BRUM Assessment"],
            required_placeholders=("Table Placeholder 1",),
            exclude_ids=used_slide_ids
        )
        if overall_slide:
            used_slide_ids.add(id(overall_slide))
            headers = ['Metric', '# of Apps Improved', '# Apps Degraded', 'Overall Result', 'Percentage Value']
            table = insert_table_at_placeholder(overall_slide, "Table Placeholder 1", 2, len(headers))
            if table:
                from pptx.util import Pt
                from pptx.dml.color import RGBColor
                for i, h in enumerate(headers):
                    table.cell(0, i).text = h
                    table.cell(0, i).text_frame.paragraphs[0].font.size = Pt(14)

                overall_result = 'Increase' if oa_up > oa_down else 'Decrease' if oa_down > oa_up else 'Even'
                percentage_value = 0 if overall_result == 'Even' else round((oa_up / max(1, oa_up + oa_down)) * 100)

                table.cell(1, 0).text = 'Overall BRUM Assessment'
                table.cell(1, 1).text = str(oa_up)
                table.cell(1, 2).text = str(oa_down)
                table.cell(1, 3).text = overall_result
                table.cell(1, 4).text = f"{percentage_value}%"

                run = table.cell(1, 4).text_frame.paragraphs[0].runs[0] if table.cell(1, 4).text_frame.paragraphs[0].runs else table.cell(1, 4).text_frame.paragraphs[0].add_run()
                if overall_result == "Increase":
                    run.font.color.rgb = RGBColor(0, 176, 80)
                elif overall_result == "Decrease":
                    run.font.color.rgb = RGBColor(192, 0, 0)
        else:
            logging.warning("[BRUM] Could not find slide for Overall Assessment.")

        # ============================
        # BRUM Entity Comparison (Slide 8) — Table Placeholder 1
        # Headers: Metric | # of Apps Improved | # Apps Degraded | Overall Result | Percentage Value
        # Rows: NetworkRequestsBRUM, HealthRulesAndAlertingBRUM
        # ============================

        def result_and_percentage(up, down):
            # Match overall slide logic: Percentage = improved ratio.
            if up > down:
                return "Increase", round((up / max(1, up + down)) * 100)
            if down > up:
                return "Decrease", round((up / max(1, up + down)) * 100)
            return "Even", 0

        # Prefer a slide by title; otherwise fall back to slide index 7 (8th slide),
        # or the next slide that has "Table Placeholder 1" not yet used.
        entity_slide = choose_slide_for_section(
            prefer_titles=["BRUM Comparison", "Entity Comparison", "Comparison Result", "BRUM Entity Comparison"],
            required_placeholders=("Table Placeholder 1",),
            exclude_ids=used_slide_ids
        )
        if not entity_slide and len(prs.slides) > 7:
            entity_slide = prs.slides[7]
            # Ensure it has the table placeholder; if not, try any unused slide with the placeholder.
            if not find_table_placeholder_by_name(entity_slide, "Table Placeholder 1"):
                for s in prs.slides:
                    if id(s) in used_slide_ids:
                        continue
                    if find_table_placeholder_by_name(s, "Table Placeholder 1"):
                        entity_slide = s
                        break

        if entity_slide and find_table_placeholder_by_name(entity_slide, "Table Placeholder 1"):
            used_slide_ids.add(id(entity_slide))

            headers = ['Metric', '# of Apps Improved', '# Apps Degraded', 'Overall Result', 'Percentage Value']
            # Compute results for each BRUM area.
            net_result, net_pct = result_and_percentage(net_up, net_down)
            hra_result, hra_pct = result_and_percentage(hra_up, hra_down)

            rows = [
                ['NetworkRequestsBRUM', str(net_up), str(net_down), net_result, f"{net_pct}%"],
                ['HealthRulesAndAlertingBRUM', str(hra_up), str(hra_down), hra_result, f"{hra_pct}%"]
            ]

            table = insert_table_at_placeholder(entity_slide, "Table Placeholder 1", len(rows) + 1, len(headers))
            if table:
                from pptx.util import Pt
                from pptx.dml.color import RGBColor

                # Header row
                for i, h in enumerate(headers):
                    cell = table.cell(0, i)
                    cell.text = h
                    p = cell.text_frame.paragraphs[0]
                    p.font.bold = True
                    p.font.size = Pt(12)

                # Data rows
                for r_idx, row in enumerate(rows, start=1):
                    for c_idx, value in enumerate(row):
                        cell = table.cell(r_idx, c_idx)
                        cell.text = value
                        p = cell.text_frame.paragraphs[0]
                        p.font.size = Pt(12)
                        # Color Percentage Value based on Overall Result.
                        if c_idx == 4:
                            result_text = rows[r_idx - 1][3]
                            run = p.runs[0] if p.runs else p.add_run()
                            if result_text == "Increase":
                                run.font.color.rgb = RGBColor(0, 176, 80)   # Green
                            elif result_text == "Decrease":
                                run.font.color.rgb = RGBColor(192, 0, 0)   # Red
        else:
            logging.warning("[BRUM] Entity comparison slide not found or missing 'Table Placeholder 1'.")

        # ============================
        # Slide 11: Network Requests Deep Dive (BRUM)
        # - Rectangles: 8..12 with key metrics
        # - Table Placeholder 1: Grade summary with declined/downgraded apps
        # ============================

        # Locate slide 11 (0-based index 10), with fallback by title/placeholder.
        deep_dive_slide = prs.slides[10] if len(prs.slides) > 10 else None
        if deep_dive_slide is None or not find_table_placeholder_by_name(deep_dive_slide, "Table Placeholder 1"):
            deep_dive_slide = choose_slide_for_section(
                prefer_titles=["Network Requests Deep Dive", "Network Requests", "BRUM Network Requests"],
                required_placeholders=("Table Placeholder 1",),
                exclude_ids=used_slide_ids
            )
        if deep_dive_slide:
            used_slide_ids.add(id(deep_dive_slide))
        else:
            logging.warning("[BRUM] Network Requests Deep Dive slide not found.")
            # Skip this section if we cannot find the slide.
            deep_dive_slide = None

        # Helper to safely set text for a named rectangle/textbox on the slide.
        def set_shape_text(slide, shape_name, text):
            if slide is None:
                return False
            for shp in slide.shapes:
                if getattr(shp, "name", "") == shape_name and hasattr(shp, "text_frame") and shp.text_frame:
                    shp.text_frame.clear()
                    shp.text_frame.text = str(text)
                    return True
            logging.debug("[BRUM] Shape '%s' not found on deep dive slide.", shape_name)
            return False

        # Column resolution helpers.
        def first_present_col(df, candidates):
            for c in candidates:
                if c in df.columns:
                    return c
            return None

        # ============================
        # Rectangles 8..12 — count declines per metric from df_network_requests (independent of Analysis).
        # Handles transitions like "True → False (Declined)" or "5 → 0 (Declined)" and generic "decreased/reduced".
        # ============================
        if deep_dive_slide:

            logging.debug("[BRUM][Slide11] df_network_requests columns: %s", list(df_network_requests.columns))

            # Resolve application column.
            def resolve_app_col(df):
                col = first_present_col(df, ["name", "Name", "applicationName", "Application Name", "Application"])
                if col:
                    return col
                for c in df.columns:
                    if str(c).lower() in ("app", "application", "application name"):
                        return c
                # fallback: first object/text column
                for c in df.columns:
                    try:
                        if df[c].dtype == object:
                            return c
                    except Exception:
                        continue
                return None

            app_col_nr_eff = resolve_app_col(df_network_requests)
            logging.debug("[BRUM][Slide11] Rectangles: resolved app_col=%s", app_col_nr_eff)

            # Transition parsing helpers: "prev → curr (Declined)" etc.
            def parse_transition_tokens(val):
                s = str(val or "").strip()
                if "→" in s:
                    prev, curr = s.split("→", 1)
                    return prev.strip(), curr.strip().split("(")[0].strip()
                return None, None

            def token_to_bool(tok):
                t = str(tok or "").strip().lower()
                if t in {"true", "yes", "y", "1"}: return True
                if t in {"false", "no", "n", "0"}: return False
                return None

            def token_to_num(tok):
                try:
                    return float(str(tok).strip())
                except Exception:
                    return None

            def is_bool_decline_cell(val):
                prev, curr = parse_transition_tokens(val)
                if prev is not None and curr is not None:
                    pb = token_to_bool(prev)
                    cb = token_to_bool(curr)
                    if pb is True and cb is False:
                        return True
                s = str(val or "").lower()
                return ("declined" in s or "downgraded" in s or "decreased" in s or "reduced" in s) and "false" in s

            def is_num_decline_cell(val):
                prev, curr = parse_transition_tokens(val)
                if prev is not None and curr is not None:
                    pn = token_to_num(prev); cn = token_to_num(curr)
                    if pn is not None and cn is not None and cn < pn:
                        return True
                s = str(val or "").lower()
                return any(k in s for k in ("declined", "decreased", "reduced", "down", "↓"))

            # Case-insensitive resolver for metric column names.
            def resolve_metric_col(df, candidates):
                # Try exact candidate names.
                col = first_present_col(df, candidates)
                if col:
                    return col
                # Case/spacing-insensitive.
                norm = lambda x: "".join(str(x).lower().split())
                cand_norms = {norm(c): c for c in candidates}
                for c in df.columns:
                    if norm(c) in cand_norms:
                        return c
                # Heuristic fallback by keyword overlap.
                key_tokens = set()
                for c in candidates:
                    key_tokens.update([t for t in str(c).lower().replace("#", "").replace("_", " ").split() if t])
                best = None; best_score = 0
                for c in df.columns:
                    lc = str(c).lower()
                    score = sum(1 for t in key_tokens if t in lc)
                    if score > best_score:
                        best, best_score = c, score
                return best if best_score >= max(2, int(len(key_tokens) * 0.4)) else None

            # Resolve metric columns.
            col_collecting   = resolve_metric_col(df_network_requests, [
                "CollectingDataPastOneDay", "CollectingDataPast1Day", "CollectingDataPastDay", "CollectingData"
            ])
            col_limit_nothit = resolve_metric_col(df_network_requests, [
                "NetworkRequestLimitNotHit", "NetworkRequestsLimitNotHit", "LimitNotHit", "RequestLimitNotHit", "networkRequestLimitNotHit"
            ])
            col_custom_rules = resolve_metric_col(df_network_requests, [
                "CustomMatchRulesCount", "# Custom Match Rules", "NumCustomMatchRules", "CustomMatchRules", "customMatchRulesCount"
            ])
            col_bt_corr      = resolve_metric_col(df_network_requests, [
                "HasBTCorrelation", "BTCorrelation", "BusinessTransactionCorrelation"
            ])
            col_ces_include  = resolve_metric_col(df_network_requests, [
                "HasCustomEventServiceIncludeRule", "CustomEventServiceIncludeRule", "HasCESIncludeRule"
            ])

            logging.debug("[BRUM][Slide11] Rectangles: metric columns -> collecting=%s, limitNotHit=%s, customRules=%s, btCorr=%s, cesInclude=%s",
                          col_collecting, col_limit_nothit, col_custom_rules, col_bt_corr, col_ces_include)

            # Count declines per metric by scanning rows directly (no Analysis gating).
            def count_metric_declines(df, app_col, metric_col, is_bool, label):
                if not app_col or not metric_col:
                    logging.warning("[BRUM][Slide11] Metric '%s' missing column (%s) or app_col (%s).", label, metric_col, app_col)
                    return 0, []
                apps = []
                for _, r in df.iterrows():
                    app = str(r.get(app_col, "") or "").strip()
                    val = r.get(metric_col, "")
                    if (is_bool and is_bool_decline_cell(val)) or ((not is_bool) and is_num_decline_cell(val)):
                        apps.append(app)
                logging.info("[BRUM][Slide11] Rectangles: %s declines=%d (apps sample: %s)", label, len(apps), apps[:10])
                return len(apps), apps

            collecting_cnt, collecting_apps       = count_metric_declines(df_network_requests, app_col_nr_eff, col_collecting,   True,  "Collecting Data Past One Day")
            limit_not_hit_cnt, limit_not_hit_apps = count_metric_declines(df_network_requests, app_col_nr_eff, col_limit_nothit, True,  "Network Request Limit Not Hit")
            custom_rules_cnt, custom_rules_apps   = count_metric_declines(df_network_requests, app_col_nr_eff, col_custom_rules, False, "# Custom Match Rules")
            bt_corr_cnt, bt_corr_apps             = count_metric_declines(df_network_requests, app_col_nr_eff, col_bt_corr,      True,  "Has BT Correlation")
            ces_include_cnt, ces_include_apps     = count_metric_declines(df_network_requests, app_col_nr_eff, col_ces_include,  True,  "Has Custom Event Service Include Rule")

            # Explicit debug for Concerto if present.
            if app_col_nr_eff and col_limit_nothit:
                concerto_rows = df_network_requests[df_network_requests[app_col_nr_eff].astype(str).str.strip().str.lower() == "concerto"]
                if not concerto_rows.empty:
                    logging.debug("[BRUM][Slide11] Concerto networkRequestLimitNotHit cell(s): %s",
                                  concerto_rows[col_limit_nothit].astype(str).tolist())

            # Write counts into the rectangles.
            set_shape_text(deep_dive_slide, "Rectangle 8",  str(collecting_cnt))
            set_shape_text(deep_dive_slide, "Rectangle 9",  str(limit_not_hit_cnt))
            set_shape_text(deep_dive_slide, "Rectangle 10", str(custom_rules_cnt))
            set_shape_text(deep_dive_slide, "Rectangle 11", str(bt_corr_cnt))
            set_shape_text(deep_dive_slide, "Rectangle 12", str(ces_include_cnt))

        # ============================
        # Table — Declined-only (Network Requests) in Table Placeholder 1.
        # Detect downgrades from Analysis ("→", "declined/downgraded to <grade>").
        # Fallback: include apps that show metric-level declines even if Analysis is missing.
        # Group apps under destination grade (Gold/Silver/Bronze) and compute per-grade percentages.
        # ============================
        if deep_dive_slide and find_table_placeholder_by_name(deep_dive_slide, "Table Placeholder 1"):

            logging.debug("[BRUM][Slide11] df_network_requests columns: %s", list(df_network_requests.columns))
            logging.debug("[BRUM][Slide11] df_analysis columns: %s", list(df_analysis.columns))

            # Canonical grade order for rank comparisons (higher rank has lower index).
            all_grades = ['platinum', 'gold', 'silver', 'bronze']
            table_grades = ['Gold', 'Silver', 'Bronze']  # display order

            def norm_grade(s):
                s = str(s).strip().lower()
                for g in all_grades:
                    if g in s:
                        return g
                return None

            def parse_transition(val):
                """
                Returns (prev_grade_norm, curr_grade_norm) or (None, None) if not parsable.
                Supports:
                  - 'Gold → Silver'
                  - 'Declined to Bronze', 'Downgraded to Silver', 'Now Gold'
                """
                s = str(val or "").strip()
                if not s:
                    return (None, None)

                # Arrow format
                if '→' in s:
                    parts = s.split('→', 1)
                    prev = norm_grade(parts[0])
                    curr = norm_grade(parts[1])
                    return (prev, curr)

                # Phrasal formats
                low = s.lower()
                import re
                m = re.search(r'(?:declined|downgraded)\s+(?:to\s+)?(platinum|gold|silver|bronze)', low)
                if m:
                    return (None, m.group(1))  # only current known
                m = re.search(r'(?:now|is\s+now|became)\s+(platinum|gold|silver|bronze)', low)
                if m:
                    return (None, m.group(1))
                # If any grade appears, treat as current grade.
                g = norm_grade(low)
                if g:
                    return (None, g)
                return (None, None)

            # Resolve app and grade columns and also prepare row lookup.
            app_col_nr = first_present_col(df_network_requests, ["name", "Name", "applicationName", "Application Name", "Application"]) \
                         or next((c for c in df_network_requests.columns if str(c).lower() in ("app", "application", "application name")), None)
            grade_col_nr = first_present_col(df_network_requests, [
                "NetworkRequestsGrade", "networkRequestsGrade", "BRUMNetworkRequestsGrade",
                "Network Requests Grade", "Grade", "grade"
            ])

            # Row lookup by app.
            def row_for_app(app):
                if not app_col_nr:
                    return None
                match = df_network_requests[df_network_requests[app_col_nr].astype(str).str.strip() == str(app)]
                return match.iloc[0] if not match.empty else None

            # Fallback grade resolver that scans all values in the NR row.
            def resolve_grade_for_app(app):
                # Prefer explicit grade column.
                if grade_col_nr:
                    r = row_for_app(app)
                    if r is not None:
                        g = norm_grade(r.get(grade_col_nr, ""))
                        if g:
                            return g.capitalize()
                # Scan entire row values for grade keywords.
                r = row_for_app(app)
                if r is not None:
                    for v in r.values:
                        g = norm_grade(v)
                        if g:
                            return g.capitalize()
                # Scan Analysis text for grade keywords.
                if "NetworkRequestsBRUM" in df_analysis.columns:
                    txt = df_analysis.loc[df_analysis["name"].astype(str).str.strip() == str(app), "NetworkRequestsBRUM"]
                    if not txt.empty:
                        _, cg = parse_transition(txt.iloc[0])
                        if cg:
                            return cg.capitalize()
                return None

            logging.debug("[BRUM][Slide11] Resolved app_col_nr=%s, grade_col_nr=%s", app_col_nr, grade_col_nr)

            # Build per-grade totals from NR sheet if possible; else infer from row scan or Analysis.
            totals_by_grade = {g: 0 for g in table_grades}
            grade_by_app = {}

            if app_col_nr:
                for _, r in df_network_requests.iterrows():
                    app = str(r.get(app_col_nr, "") or "").strip()
                    g = None
                    if grade_col_nr:
                        g = norm_grade(r.get(grade_col_nr, ""))
                    if not g:
                        # scan row values
                        for v in r.values:
                            g = norm_grade(v)
                            if g:
                                break
                    disp = g.capitalize() if g else None
                    if app and disp in totals_by_grade:
                        totals_by_grade[disp] += 1
                        grade_by_app[app] = disp

            if not any(v > 0 for v in totals_by_grade.values()) and "NetworkRequestsBRUM" in df_analysis.columns:
                # Infer totals from Analysis mentions.
                inferred_totals = {g: 0 for g in table_grades}
                for _, r in df_analysis.iterrows():
                    _, cg = parse_transition(r.get("NetworkRequestsBRUM", ""))
                    disp = cg.capitalize() if cg else None
                    if disp in inferred_totals:
                        inferred_totals[disp] += 1
                if any(inferred_totals.values()):
                    totals_by_grade = inferred_totals
                    logging.warning("[BRUM][Slide11] Using inferred per-grade totals from Analysis: %s", totals_by_grade)

            # A) Declines detected from Analysis.
            declined_apps_analysis = set()
            if "NetworkRequestsBRUM" in df_analysis.columns:
                for _, r in df_analysis.iterrows():
                    app = str(r.get("name", "") or "").strip()
                    val = r.get("NetworkRequestsBRUM", "")
                    prev_g, curr_g = parse_transition(val)
                    is_dg = False
                    low = str(val).lower()
                    if 'declined' in low or 'downgraded' in low:
                        is_dg = True
                    elif prev_g and curr_g and prev_g in all_grades and curr_g in all_grades:
                        is_dg = all_grades.index(prev_g) < all_grades.index(curr_g)
                    if is_dg and app:
                        declined_apps_analysis.add(app)

            logging.info("[BRUM][Slide11] Declined apps from Analysis: %d", len(declined_apps_analysis))
            logging.debug("[BRUM][Slide11] Declined apps (Analysis) sample: %s", list(sorted(declined_apps_analysis))[:20])

            # B) Declines detected from metric columns (same detectors as rectangles).
            metric_cols_and_types = [
                (col_collecting,   True),
                (col_limit_nothit, True),
                (col_custom_rules, False),
                (col_bt_corr,      True),
                (col_ces_include,  True),
            ]
            metric_declined_apps = set()
            if app_col_nr:
                for col, is_bool in metric_cols_and_types:
                    if not col:
                        continue
                    for _, r in df_network_requests.iterrows():
                        app = str(r.get(app_col_nr, "") or "").strip()
                        val = r.get(col, "")
                        if (is_bool and is_bool_decline_cell(val)) or ((not is_bool) and is_num_decline_cell(val)):
                            metric_declined_apps.add(app)

            logging.info("[BRUM][Slide11] Declined apps from NR metrics: %d", len(metric_declined_apps))
            logging.debug("[BRUM][Slide11] Declined apps (Metrics) sample: %s", list(sorted(metric_declined_apps))[:20])

            # Union of Analysis- and Metric-detected declines.
            declined_union = declined_apps_analysis.union(metric_declined_apps)
            logging.info("[BRUM][Slide11] Total declined apps (union): %d", len(declined_union))

            # Group apps under destination grade.
            declined_by_grade = {g: [] for g in table_grades}
            missing_grade = []
            for app in sorted(declined_union):
                dest = grade_by_app.get(app)
                if not dest:
                    dest = resolve_grade_for_app(app)
                if dest in declined_by_grade:
                    declined_by_grade[dest].append(app)
                else:
                    missing_grade.append(app)

            if missing_grade:
                logging.warning("[BRUM][Slide11] %d declined apps have no resolvable grade (not shown in table): %s",
                                len(missing_grade), missing_grade[:20])

            logging.debug("[BRUM][Slide11] Declined-by-grade counts: %s",
                          {k: len(v) for k, v in declined_by_grade.items()})
            logging.debug("[BRUM][Slide11] Final totals_by_grade: %s", totals_by_grade)

            # Build table rows in Gold, Silver, Bronze order.
            headers = ["Grade", "Application Names", "Number of Applications", "Percentage Declined"]
            rows = []
            for g in table_grades:
                names = sorted(declined_by_grade[g], key=str.lower)
                num_apps = len(names)
                denom = totals_by_grade.get(g, 0)
                pct = round((num_apps / denom) * 100) if denom > 0 else 0
                logging.info("[BRUM][Slide11] Grade=%s declined=%d total=%d pct=%d%%", g, num_apps, denom, pct)
                rows.append([g, "\n".join(names) if names else "—", str(num_apps), f"{pct}%"])

            # Insert and render the single table.
            table = insert_table_at_placeholder(deep_dive_slide, "Table Placeholder 1", len(rows) + 1, len(headers))
            if table:
                from pptx.util import Pt
                from pptx.dml.color import RGBColor
                for i, h in enumerate(headers):
                    cell = table.cell(0, i); cell.text = h
                    p = cell.text_frame.paragraphs[0]; p.font.bold = True; p.font.size = Pt(12)
                for r_idx, row in enumerate(rows, start=1):
                    for c_idx, val in enumerate(row):
                        cell = table.cell(r_idx, c_idx); cell.text = val
                        p = cell.text_frame.paragraphs[0]; p.font.size = Pt(12)
                        if c_idx == 3:
                            run = p.runs[0] if p.runs else p.add_run()
                            pct_num = int(str(val).replace("%", "") or 0)
                            run.font.color.rgb = RGBColor(192, 0, 0) if pct_num > 0 else RGBColor(0, 176, 80)
            else:
                logging.error("[BRUM][Slide11] Could not insert Network Requests Declined table; placeholder missing.") 

        # ============================
        # Slide 12: Health Rules & Alerting Deep Dive (BRUM)
        # - Rectangles: 10..12 with key metrics
        # - Table Placeholder 1: Grade summary with declined/downgraded apps
        # ============================

        # Locate slide 12 (0-based index 11), fallback by title/placeholder.
        hra_deep_dive_slide = prs.slides[11] if len(prs.slides) > 11 else None
        if hra_deep_dive_slide is None or not find_table_placeholder_by_name(hra_deep_dive_slide, "Table Placeholder 1"):
            hra_deep_dive_slide = choose_slide_for_section(
                prefer_titles=["Health Rules & Alerting Deep Dive", "Health Rules Deep Dive", "Health Rules & Alerting", "Health Rules"],
                required_placeholders=("Table Placeholder 1",),
                exclude_ids=used_slide_ids
            )
        if hra_deep_dive_slide:
            used_slide_ids.add(id(hra_deep_dive_slide))
        else:
            logging.warning("[BRUM] Health Rules & Alerting Deep Dive slide not found.")
            hra_deep_dive_slide = None

        # Helpers (reuse patterns from Slide 11).
        def set_shape_text(slide, shape_name, text):
            if slide is None:
                return False
            for shp in slide.shapes:
                if getattr(shp, "name", "") == shape_name and hasattr(shp, "text_frame") and shp.text_frame:
                    shp.text_frame.clear()
                    shp.text_frame.text = str(text)
                    return True
            logging.debug("[BRUM] Shape '%s' not found on HRA deep dive slide.", shape_name)
            return False

        def first_present_col(df, candidates):
            for c in candidates:
                if c in df.columns:
                    return c
            return None

        # ============================
        # Rectangles 10..12 — count declines per metric from df_health_rules (independent of Analysis).
        # Handles transitions like "5 → 0 (Declined)" and generic "decreased/reduced".
        # ============================
        if hra_deep_dive_slide:

            logging.debug("[BRUM][Slide12] df_health_rules columns: %s", list(df_health_rules.columns))

            # Resolve application column.
            def resolve_app_col(df):
                col = first_present_col(df, ["name", "Name", "applicationName", "Application Name", "Application"])
                if col:
                    return col
                for c in df.columns:
                    if str(c).lower() in ("app", "application", "application name"):
                        return c
                # fallback: first text/object column
                for c in df.columns:
                    try:
                        if df[c].dtype == object:
                            return c
                    except Exception:
                        continue
                return None

            app_col_hr_eff = resolve_app_col(df_health_rules)
            logging.debug("[BRUM][Slide12] Rectangles: resolved app_col=%s", app_col_hr_eff)

            # Transition parsing helpers: "prev → curr (Declined)" etc.
            def parse_transition_tokens(val):
                s = str(val or "").strip()
                if "→" in s:
                    prev, curr = s.split("→", 1)
                    return prev.strip(), curr.strip().split("(")[0].strip()
                return None, None

            def token_to_num(tok):
                try:
                    return float(str(tok).strip())
                except Exception:
                    return None

            def is_num_decline_cell(val):
                prev, curr = parse_transition_tokens(val)
                if prev is not None and curr is not None:
                    pn = token_to_num(prev); cn = token_to_num(curr)
                    if pn is not None and cn is not None and cn < pn:
                        return True
                s = str(val or "").lower()
                return any(k in s for k in ("declined", "decreased", "reduced", "down", "↓"))

            # Case-insensitive resolver for metric column names.
            def resolve_metric_col(df, candidates):
                col = first_present_col(df, candidates)
                if col:
                    return col
                # Case/spacing-insensitive match.
                norm = lambda x: "".join(str(x).lower().split())
                cand_norms = {norm(c): c for c in candidates}
                for c in df.columns:
                    if norm(c) in cand_norms:
                        return c
                # Heuristic fallback by keyword overlap.
                key_tokens = set()
                for c in candidates:
                    key_tokens.update([t for t in str(c).lower().replace("#", "").replace("_", " ").split() if t])
                best = None; best_score = 0
                for c in df.columns:
                    lc = str(c).lower()
                    score = sum(1 for t in key_tokens if t in lc)
                    if score > best_score:
                        best, best_score = c, score
                return best if best_score >= max(2, int(len(key_tokens) * 0.4)) else None

            # Resolve metric columns (common headers + variants).
            col_violations   = resolve_metric_col(df_health_rules, [
                "NumberOfHealthRuleViolations", "# Of Health Rule Violations",
                "HealthRuleViolations", "HealthRulesViolations", "numHealthRuleViolations"
            ])
            col_actions_bound = resolve_metric_col(df_health_rules, [
                "NumberOfActionsBoundToEnabledPolicies", "# Of Actions Bound To Enabled Policies",
                "ActionsBoundToEnabledPolicies", "ActionsBoundEnabledPolicies", "ActionsBoundEnabledPoliciesCount"
            ])
            col_custom_rules = resolve_metric_col(df_health_rules, [
                "NumberOfCustomHealthRules", "# Of Custom Health Rules",
                "CustomHealthRulesCount", "CustomHealthRules"
            ])

            logging.debug("[BRUM][Slide12] Rectangles: metric columns -> violations=%s, actionsBound=%s, customRules=%s",
                          col_violations, col_actions_bound, col_custom_rules)

            # Count numeric declines per metric.
            def count_metric_declines(df, app_col, metric_col, label):
                if not app_col or not metric_col:
                    logging.warning("[BRUM][Slide12] Metric '%s' missing column (%s) or app_col (%s).", label, metric_col, app_col)
                    return 0, []
                apps = []
                for _, r in df.iterrows():
                    app = str(r.get(app_col, "") or "").strip()
                    val = r.get(metric_col, "")
                    if is_num_decline_cell(val):
                        apps.append(app)
                logging.info("[BRUM][Slide12] Rectangles: %s declines=%d (apps sample: %s)", label, len(apps), apps[:10])
                return len(apps), apps

            violations_cnt, violations_apps       = count_metric_declines(df_health_rules, app_col_hr_eff, col_violations,   "Number Of Health Rule Violations")
            actions_bound_cnt, actions_bound_apps = count_metric_declines(df_health_rules, app_col_hr_eff, col_actions_bound, "Number Of Actions Bound To Enabled Policies")
            custom_rules_cnt, custom_rules_apps   = count_metric_declines(df_health_rules, app_col_hr_eff, col_custom_rules,  "Number Of Custom Health Rules")

            # Write counts into the rectangles.
            set_shape_text(hra_deep_dive_slide, "Rectangle 10", str(violations_cnt))
            set_shape_text(hra_deep_dive_slide, "Rectangle 11", str(actions_bound_cnt))
            set_shape_text(hra_deep_dive_slide, "Rectangle 12", str(custom_rules_cnt))

        # ============================
        # Table — Declined-only (Health Rules & Alerting) in Table Placeholder 1.
        # Detect downgrades from Analysis ("→", "declined/downgraded to <grade>").
        # Fallback: include apps that show metric-level declines even if Analysis is missing.
        # Group apps under destination grade (Gold/Silver/Bronze) and compute per-grade percentages.
        # ============================
        if hra_deep_dive_slide and find_table_placeholder_by_name(hra_deep_dive_slide, "Table Placeholder 1"):

            logging.debug("[BRUM][Slide12] df_health_rules columns: %s", list(df_health_rules.columns))
            logging.debug("[BRUM][Slide12] df_analysis columns: %s", list(df_analysis.columns))

            # Canonical grade order for rank comparisons.
            all_grades   = ['platinum', 'gold', 'silver', 'bronze']
            table_grades = ['Gold', 'Silver', 'Bronze']

            def norm_grade(s):
                s = str(s).strip().lower()
                for g in all_grades:
                    if g in s:
                        return g
                return None

            def parse_transition(val):
                s = str(val or "").strip()
                if not s:
                    return (None, None)
                if '→' in s:
                    parts = s.split('→', 1)
                    prev = norm_grade(parts[0])
                    curr = norm_grade(parts[1])
                    return (prev, curr)
                low = s.lower()
                import re
                m = re.search(r'(?:declined|downgraded)\s+(?:to\s+)?(platinum|gold|silver|bronze)', low)
                if m:
                    return (None, m.group(1))
                m = re.search(r'(?:now|is\s+now|became)\s+(platinum|gold|silver|bronze)', low)
                if m:
                    return (None, m.group(1))
                g = norm_grade(low)
                if g:
                    return (None, g)
                return (None, None)

            # Resolve app and grade columns and prepare row lookup.
            app_col_hr = first_present_col(df_health_rules, ["name", "Name", "applicationName", "Application Name", "Application"]) \
                         or next((c for c in df_health_rules.columns if str(c).lower() in ("app", "application", "application name")), None)
            grade_col_hr = first_present_col(df_health_rules, [
                "HealthRulesAndAlertingGrade", "HealthRulesGrade", "BRUMHealthRulesGrade",
                "Health Rules Grade", "Grade", "grade"
            ])

            def row_for_app_hr(app):
                if not app_col_hr:
                    return None
                match = df_health_rules[df_health_rules[app_col_hr].astype(str).str.strip() == str(app)]
                return match.iloc[0] if not match.empty else None

            def resolve_grade_for_app_hr(app):
                # Prefer explicit grade column.
                if grade_col_hr:
                    r = row_for_app_hr(app)
                    if r is not None:
                        g = norm_grade(r.get(grade_col_hr, ""))
                        if g:
                            return g.capitalize()
                # Scan entire row values for grade keywords.
                r = row_for_app_hr(app)
                if r is not None:
                    for v in r.values:
                        g = norm_grade(v)
                        if g:
                            return g.capitalize()
                # Scan Analysis text for grade keywords.
                if "HealthRulesAndAlertingBRUM" in df_analysis.columns:
                    txt = df_analysis.loc[df_analysis["name"].astype(str).str.strip() == str(app), "HealthRulesAndAlertingBRUM"]
                    if not txt.empty:
                        _, cg = parse_transition(txt.iloc[0])
                        if cg:
                            return cg.capitalize()
                return None

            logging.debug("[BRUM][Slide12] Resolved app_col_hr=%s, grade_col_hr=%s", app_col_hr, grade_col_hr)

            # Build per-grade totals from HRA sheet if possible; else infer.
            totals_by_grade_hr = {g: 0 for g in table_grades}
            grade_by_app_hr = {}

            if app_col_hr:
                for _, r in df_health_rules.iterrows():
                    app = str(r.get(app_col_hr, "") or "").strip()
                    g = None
                    if grade_col_hr:
                        g = norm_grade(r.get(grade_col_hr, ""))
                    if not g:
                        for v in r.values:
                            g = norm_grade(v)
                            if g:
                                break
                    disp = g.capitalize() if g else None
                    if app and disp in totals_by_grade_hr:
                        totals_by_grade_hr[disp] += 1
                        grade_by_app_hr[app] = disp

            if not any(v > 0 for v in totals_by_grade_hr.values()) and "HealthRulesAndAlertingBRUM" in df_analysis.columns:
                inferred_totals = {g: 0 for g in table_grades}
                for _, r in df_analysis.iterrows():
                    _, cg = parse_transition(r.get("HealthRulesAndAlertingBRUM", ""))
                    disp = cg.capitalize() if cg else None
                    if disp in inferred_totals:
                        inferred_totals[disp] += 1
                if any(inferred_totals.values()):
                    totals_by_grade_hr = inferred_totals
                    logging.warning("[BRUM][Slide12] Using inferred per-grade totals from Analysis: %s", totals_by_grade_hr)

            # A) Declines detected from Analysis.
            declined_apps_analysis_hr = set()
            if "HealthRulesAndAlertingBRUM" in df_analysis.columns:
                for _, r in df_analysis.iterrows():
                    app = str(r.get("name", "") or "").strip()
                    val = r.get("HealthRulesAndAlertingBRUM", "")
                    prev_g, curr_g = parse_transition(val)
                    is_dg = False
                    low = str(val).lower()
                    if 'declined' in low or 'downgraded' in low:
                        is_dg = True
                    elif prev_g and curr_g and prev_g in all_grades and curr_g in all_grades:
                        is_dg = all_grades.index(prev_g) < all_grades.index(curr_g)
                    if is_dg and app:
                        declined_apps_analysis_hr.add(app)

            logging.info("[BRUM][Slide12] Declined apps from Analysis: %d", len(declined_apps_analysis_hr))
            logging.debug("[BRUM][Slide12] Declined apps (Analysis) sample: %s", list(sorted(declined_apps_analysis_hr))[:20])

            # B) Declines detected from HRA metric columns (use same parser as rectangles).
            metric_cols_hr = [col_violations, col_actions_bound, col_custom_rules]
            metric_declined_apps_hr = set()
            if app_col_hr:
                for col in metric_cols_hr:
                    if not col:
                        continue
                    for _, r in df_health_rules.iterrows():
                        app = str(r.get(app_col_hr, "") or "").strip()
                        val = r.get(col, "")
                        if is_num_decline_cell(val):
                            metric_declined_apps_hr.add(app)

            logging.info("[BRUM][Slide12] Declined apps from HRA metrics: %d", len(metric_declined_apps_hr))
            logging.debug("[BRUM][Slide12] Declined apps (Metrics) sample: %s", list(sorted(metric_declined_apps_hr))[:20])

            # Union of Analysis- and Metric-detected declines.
            declined_union_hr = declined_apps_analysis_hr.union(metric_declined_apps_hr)
            logging.info("[BRUM][Slide12] Total declined apps (union): %d", len(declined_union_hr))

            # Group apps under destination grade.
            declined_by_grade_hr = {g: [] for g in table_grades}
            missing_grade_hr = []
            for app in sorted(declined_union_hr):
                dest = grade_by_app_hr.get(app)
                if not dest:
                    dest = resolve_grade_for_app_hr(app)
                if dest in declined_by_grade_hr:
                    declined_by_grade_hr[dest].append(app)
                else:
                    missing_grade_hr.append(app)

            if missing_grade_hr:
                logging.warning("[BRUM][Slide12] %d declined apps have no resolvable grade (not shown in table): %s",
                                len(missing_grade_hr), missing_grade_hr[:20])

            logging.debug("[BRUM][Slide12] Declined-by-grade counts: %s",
                          {k: len(v) for k, v in declined_by_grade_hr.items()})
            logging.debug("[BRUM][Slide12] Final totals_by_grade_hr: %s", totals_by_grade_hr)

            # Build table in Gold, Silver, Bronze order.
            headers = ["Grade", "Application Names", "Number of Applications", "Percentage Declined"]
            rows = []
            for g in table_grades:
                names = sorted(declined_by_grade_hr[g], key=str.lower)
                num_apps = len(names)
                denom = totals_by_grade_hr.get(g, 0)
                pct = round((num_apps / denom) * 100) if denom > 0 else 0
                logging.info("[BRUM][Slide12] Grade=%s declined=%d total=%d pct=%d%%", g, num_apps, denom, pct)
                rows.append([g, "\n".join(names) if names else "—", str(num_apps), f"{pct}%"])

            table = insert_table_at_placeholder(hra_deep_dive_slide, "Table Placeholder 1", len(rows) + 1, len(headers))
            if table:
                from pptx.util import Pt
                from pptx.dml.color import RGBColor
                for i, h in enumerate(headers):
                    cell = table.cell(0, i); cell.text = h
                    p = cell.text_frame.paragraphs[0]; p.font.bold = True; p.font.size = Pt(12)
                for r_idx, row in enumerate(rows, start=1):
                    for c_idx, val in enumerate(row):
                        cell = table.cell(r_idx, c_idx); cell.text = val
                        p = cell.text_frame.paragraphs[0]; p.font.size = Pt(12)
                        if c_idx == 3:
                            run = p.runs[0] if p.runs else p.add_run()
                            pct_num = int(str(val).replace("%", "") or 0)
                            run.font.color.rgb = RGBColor(192, 0, 0) if pct_num > 0 else RGBColor(0, 176, 80)
            else:
                logging.error("[BRUM][Slide12] Could not insert Health Rules & Alerting Declined table; placeholder missing.")

        # ============================
        # Populate "TextBox 7" with number of BRUM applications (slide index 5)
        # ============================

        def set_textbox_value(prs_obj, shape_name, text, fallback_slide_index=5):
            for s in prs_obj.slides:
                for shp in s.shapes:
                    if getattr(shp, "name", "") == shape_name and hasattr(shp, "text_frame"):
                        shp.text_frame.clear()
                        shp.text_frame.text = str(text)
                        return True
            if len(prs_obj.slides) > fallback_slide_index:
                s = prs_obj.slides[fallback_slide_index]
                for shp in s.shapes:
                    if getattr(shp, "name", "") == shape_name and hasattr(shp, "text_frame"):
                        shp.text_frame.clear()
                        shp.text_frame.text = str(text)
                        return True
            return False

        if not set_textbox_value(prs, "TextBox 7", number_of_apps, fallback_slide_index=5):
            logging.warning("[BRUM] 'TextBox 7' not found; BRUM application count not written.")

        # Save the presentation when all slide content is populated.
        prs.save(powerpoint_output_path)
        logging.debug(f"[BRUM] PowerPoint saved to: {powerpoint_output_path}")

    except Exception as e:
        logging.error(f"[BRUM] Error generating PowerPoint: {e}", exc_info=True)
        raise

################################################################################### 
############### MRUM PPT                                   ########################
###################################################################################

def generate_powerpoint_from_mrum(comparison_result_path, powerpoint_output_path, current_file_path, previous_file_path):
    logging.debug("[MRUM] Generating PowerPoint presentation...")

    try:
        # Common imports used throughout.
        from pptx.util import Pt
        from pptx.dml.color import RGBColor

        # Template
        template_folder = config.get('TEMPLATE_FOLDER', 'templates')
        template_path = os.path.join(template_folder, 'template_mrum.pptx')
        if not os.path.exists(template_path):
            template_path = os.getenv('MRUM_TEMPLATE_PATH', template_path)
            if not os.path.exists(template_path):
                template_path = input("MRUM template not found! Provide full path to template_mrum.pptx: ")

        prs = Presentation(template_path)
        logging.debug(f"[MRUM] Template loaded from: {template_path}")

        # Load current/previous Summary (for Key Callouts), and comparison workbook sheets.
        current_summary_df = pd.read_excel(current_file_path, sheet_name='Summary')
        previous_summary_df = pd.read_excel(previous_file_path, sheet_name='Summary')

        summary_df = pd.read_excel(comparison_result_path, sheet_name='Summary')
        df_analysis = pd.read_excel(comparison_result_path, sheet_name='Analysis')
        df_network_requests = pd.read_excel(comparison_result_path, sheet_name='NetworkRequestsMRUM')
        df_health_rules = pd.read_excel(comparison_result_path, sheet_name='HealthRulesAndAlertingMRUM')
        try:
            df_overall_mrum = pd.read_excel(comparison_result_path, sheet_name='OverallAssessmentMRUM')
        except Exception:
            df_overall_mrum = pd.DataFrame()

        # Count apps in current analysis for widgets.
        try:
            df_current_analysis = pd.read_excel(current_file_path, sheet_name='Analysis')
            number_of_apps = df_current_analysis['name'].dropna().str.strip().ne('').sum()
        except Exception:
            number_of_apps = 0

        # Helpers
        def find_table_placeholder_by_name(slide, name):
            for shape in slide.shapes:
                if shape.is_placeholder and shape.name == name:
                    return shape
            return None

        def insert_table_at_placeholder(slide, placeholder_name, rows, cols):
            placeholder = find_table_placeholder_by_name(slide, placeholder_name)
            if not placeholder:
                logging.error(f"[MRUM] Placeholder '{placeholder_name}' not found on the slide.")
                return None
            left, top, width, height = placeholder.left, placeholder.top, placeholder.width, placeholder.height
            return slide.shapes.add_table(rows, cols, left, top, width, height).table

        def slide_title_text(slide):
            for shp in slide.shapes:
                if getattr(shp, "name", "") == "Title 2" and hasattr(shp, "text_frame") and shp.text_frame:
                    return shp.text_frame.text.strip()
            for shp in slide.shapes:
                if hasattr(shp, "text_frame") and shp.text_frame and shp.text_frame.text:
                    return shp.text_frame.text.strip()
            return ""

        used_slide_ids = set()

        def choose_slide_for_section(prefer_titles, required_placeholders=("Table Placeholder 1",), exclude_ids=None):
            exclude_ids = exclude_ids or set()
            for s in prs.slides:
                if id(s) in exclude_ids:
                    continue
                title = slide_title_text(s)
                if any(pt.lower() in title.lower() for pt in prefer_titles):
                    if all(find_table_placeholder_by_name(s, p) for p in required_placeholders):
                        return s
            for s in prs.slides:
                if id(s) in exclude_ids:
                    continue
                if all(find_table_placeholder_by_name(s, p) for p in required_placeholders):
                    return s
            return None

        # ============================
        # Key Callouts (Table Placeholder 1)
        # ============================
        def last_percent(df, col):
            if df.empty or col not in df.columns:
                return None
            s = pd.to_numeric(df[col].astype(str).str.replace('%', ''), errors='coerce').dropna()
            return float(s.iloc[-1]) if not s.empty else None

        curr_gold = last_percent(df_overall_mrum, 'percentageTotalGoldOrBetter')
        prev_gold = last_percent(pd.read_excel(previous_file_path, sheet_name='OverallAssessmentMRUM') if 'OverallAssessmentMRUM' in pd.ExcelFile(previous_file_path).sheet_names else pd.DataFrame(), 'percentageTotalGoldOrBetter')
        curr_plat = last_percent(df_overall_mrum, 'percentageTotalPlatinum')
        prev_plat = last_percent(pd.read_excel(previous_file_path, sheet_name='OverallAssessmentMRUM') if 'OverallAssessmentMRUM' in pd.ExcelFile(previous_file_path).sheet_names else pd.DataFrame(), 'percentageTotalPlatinum')

        def count_changes(df, col):
            if col not in df.columns:
                return 0, 0
            s = df[col].astype(str)
            return (
                s.str.contains('Upgraded', case=False, na=False).sum(),
                s.str.contains('Downgraded', case=False, na=False).sum()
            )

        oa_up, oa_down   = count_changes(df_analysis, 'OverallAssessment')
        net_up, net_down = count_changes(df_analysis, 'NetworkRequestsMRUM')
        hra_up, hra_down = count_changes(df_analysis, 'HealthRulesAndAlertingMRUM')

        def arrow(curr, prev):
            if curr is None or prev is None:
                return '→'
            return '↑' if curr > prev else '↓' if curr < prev else '→'

        def fmt_change(prev, curr, suffix='%'):
            if prev is None or curr is None:
                return "Change observed."
            return (
                f"Increase from {prev:.1f}{suffix}→{curr:.1f}{suffix}" if curr > prev else
                f"Decrease from {prev:.1f}{suffix}→{curr:.1f}{suffix}" if curr < prev else
                f"No change ({curr:.1f}{suffix})."
            )

        slides_with_ph = [s for s in prs.slides if find_table_placeholder_by_name(s, "Table Placeholder 1")]
        key_callouts_slide = slides_with_ph[0] if slides_with_ph else (prs.slides[0] if len(prs.slides) else None)
        if key_callouts_slide:
            used_slide_ids.add(id(key_callouts_slide))
            headers = ["AppD Maturity Progression & Engagement", "Commentary", "Outcomes", "Change/Status Since Last"]
            rows = [
                [
                    "B/S/G/P Model Adoption & Maturity Status (MRUM).",
                    f"MRUM analysis coverage across {int(number_of_apps)} applications.",
                    f"Overall MRUM upgrades: {oa_up}, downgrades: {oa_down}.",
                    '↑' if oa_up > oa_down else '↓' if oa_down > oa_up else '→'
                ],
                [
                    "Gold Status Apps (MRUM).",
                    "Change in gold-or-better coverage across the portfolio.",
                    fmt_change(prev_gold, curr_gold),
                    arrow(curr_gold, prev_gold)
                ],
                [
                    "Platinum Status Apps (MRUM).",
                    "Top-tier MRUM maturity presence across applications.",
                    f"{curr_plat:.1f}% platinum." if curr_plat is not None else "Platinum presence observed.",
                    arrow(curr_plat, prev_plat)
                ],
                [
                    "Maturity Partnership (MRUM).",
                    "Ongoing improvements in MRUM instrumentation and alerting.",
                    f"Network Requests (↑{net_up}/↓{net_down}), Health Rules (↑{hra_up}/↓{hra_down}).",
                    '↑' if (net_up + hra_up) > (net_down + hra_down) else '↓' if (net_down + hra_down) > (net_up + hra_up) else '→'
                ]
            ]
            table = insert_table_at_placeholder(key_callouts_slide, "Table Placeholder 1", len(rows) + 1, len(headers))
            if table:
                # Header row
                for i, h in enumerate(headers):
                    cell = table.cell(0, i); cell.text = h
                    p = cell.text_frame.paragraphs[0]; p.font.bold = True; p.font.size = Pt(12)

                # Optional: widen status column if helper exists
                try:
                    autosize_col_to_header(table, 3, header_pt=12, padding_in=0.6, avg_char_em=0.55)
                except Exception:
                    pass

                # Data rows — render filled triangles with text-variant neutral to avoid emoji
                for r_idx, row in enumerate(rows, start=1):
                    for c_idx, value in enumerate(row):
                        cell = table.cell(r_idx, c_idx)
                        if c_idx == 3 and value in ('↑', '↓', '→'):
                            set_arrow_cell(cell, value, color=PINK, size_pt=36, font_name='Calibri')
                        else:
                            cell.text = str(value)
                            p = cell.text_frame.paragraphs[0]
                            p.font.size = Pt(12)

        # --- MRUM: maturity badge + notes (self-contained, safe) ---
        try:
            def _mrum_grade_token(v):
                if not isinstance(v, str):
                    v = "" if pd.isna(v) else str(v)
                m = re.search(r'(platinum|gold|silver|bronze)', v, re.I)
                return m.group(1).lower() if m else None

            def _overall_maturity_from_df_mrum(df):
                col = 'OverallAssessment'
                if df is None or col not in df.columns:
                    return None, {'bronze':0,'silver':0,'gold':0,'platinum':0}, 0
                counts = {'bronze':0,'silver':0,'gold':0,'platinum':0}
                rated = 0
                for v in df[col]:
                    t = _mrum_grade_token(v)
                    if t in counts:
                        counts[t] += 1
                        rated += 1
                if rated == 0:
                    return None, counts, 0
                rank = {'bronze':0,'silver':1,'gold':2,'platinum':3}
                best = max(counts.items(), key=lambda kv: (kv[1], rank[kv[0]]))
                return best[0].title(), counts, rated

            def _ideal_text_rgb_local(rgb):
                r, g, b = rgb[0], rgb[1], rgb[2]
                bright = (r*299 + g*587 + b*114)/1000.0
                return RGBColor(255,255,255) if bright < 140 else RGBColor(31,31,31)

            def _color_oval_for_maturity_local(slide_obj, shape_name, tier):
                palette = {
                    'Bronze':   RGBColor(205,127,50),
                    'Silver':   RGBColor(166,166,166),
                    'Gold':     RGBColor(255,192,0),
                    'Platinum': RGBColor(190,190,200),
                }
                if slide_obj is None or tier not in palette:
                    return
                target = next((sh for sh in slide_obj.shapes if getattr(sh, "name", "") == shape_name), None)
                if not target:
                    return
                target.fill.solid()
                target.fill.fore_color.rgb = palette[tier]
                if hasattr(target, "text_frame") and target.text_frame:
                    fg = _ideal_text_rgb_local(palette[tier])
                    for p in target.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.color.rgb = fg

            # Compute majority tier and coverage.
            overall_tier, counts, rated = _overall_maturity_from_df_mrum(df_current_analysis)
            if overall_tier and key_callouts_slide is not None:
                _color_oval_for_maturity_local(key_callouts_slide, "Oval 10", overall_tier)

                total_apps = int(number_of_apps) if number_of_apps else 0
                coverage = (rated / total_apps * 100.0) if total_apps else 0.0
                def pct(n): return (n / rated * 100.0) if rated else 0.0
                pb, ps, pg, pp_ = (
                    pct(counts['bronze']), pct(counts['silver']),
                    pct(counts['gold']),   pct(counts['platinum'])
                )

                notes = key_callouts_slide.notes_slide
                tf = notes.notes_text_frame
                tf.clear()
                tf.paragraphs[0].text = "Overall tier selection: majority of app ratings in Analysis; ties prefer higher tier (Platinum > Gold > Silver > Bronze)."
                p2 = tf.add_paragraph()
                p2.text = (
                    f"Status is {overall_tier} based on rated distribution — "
                    f"Platinum {pp_:.1f}% ({counts['platinum']}), "
                    f"Gold {pg:.1f}% ({counts['gold']}), "
                    f"Silver {ps:.1f}% ({counts['silver']}), "
                    f"Bronze {pb:.1f}% ({counts['bronze']})."
                )
                p3 = tf.add_paragraph()
                p3.text = f"Rated coverage this run: {coverage:.1f}% ({rated}/{total_apps})."
        except Exception as e:
            logging.warning("[MRUM] Maturity badge/notes skipped: %s", e)

        # ============================
        # Applications Improved (Table Placeholder 1 on next slide)
        # ============================
        improved = []
        cols_map = [
            ('NetworkRequestsMRUM', 'Network Requests'),
            ('HealthRulesAndAlertingMRUM', 'Health Rules & Alerting'),
            ('OverallAssessment', 'Overall')
        ]
        for _, r in df_analysis.iterrows():
            app = str(r.get('name', '') or '').strip()
            if not app:
                continue
            areas = []
            for col, label in cols_map:
                if col in df_analysis.columns:
                    val = r.get(col, '')
                    if isinstance(val, str) and 'upgraded' in val.lower():
                        areas.append(label)
            if areas:
                improved.append((app, ', '.join(areas)))
        improved.sort(key=lambda x: x[0].lower())

        improved_slide = None
        for s in slides_with_ph:
            if id(s) not in used_slide_ids:
                improved_slide = s
                break
        if improved_slide:
            used_slide_ids.add(id(improved_slide))
            headers = ["Application", "Improvement Areas"]
            row_count = max(1, len(improved)) + 1
            table = insert_table_at_placeholder(improved_slide, "Table Placeholder 1", row_count, len(headers))
            if table:
                for c, h in enumerate(headers):
                    table.cell(0, c).text = h
                    table.cell(0, c).text_frame.paragraphs[0].font.size = Pt(12)
                if improved:
                    for idx, (app, areas) in enumerate(improved, start=1):
                        table.cell(idx, 0).text = app; table.cell(idx, 1).text = areas
                        table.cell(idx, 0).text_frame.paragraphs[0].font.size = Pt(12)
                        table.cell(idx, 1).text_frame.paragraphs[0].font.size = Pt(12)
                else:
                    table.cell(1, 0).text = "No applications improved in this period."
                    table.cell(1, 1).text = ""

        # ============================
        # Summary slide: Previous (TP 4), Current (TP 3), Comparison (TP 1)
        # ============================
        def find_slide_with_all_placeholders(prs_obj, names):
            for s in prs_obj.slides:
                if all(find_table_placeholder_by_name(s, n) for n in names):
                    return s
            return None

        summary_slide = find_slide_with_all_placeholders(prs, ["Table Placeholder 1", "Table Placeholder 3", "Table Placeholder 4"])
        if not summary_slide:
            for s in prs.slides:
                if any(find_table_placeholder_by_name(s, n) for n in ["Table Placeholder 1", "Table Placeholder 3", "Table Placeholder 4"]):
                    summary_slide = s
                    break
        if summary_slide:
            used_slide_ids.add(id(summary_slide))

            def fill_table_from_df(table, df):
                for c, col in enumerate(df.columns):
                    table.cell(0, c).text = str(col)
                    table.cell(0, c).text_frame.paragraphs[0].font.size = Pt(12)
                for r_idx, row in df.iterrows():
                    for c_idx, val in enumerate(row):
                        table.cell(r_idx + 1, c_idx).text = "" if pd.isna(val) else str(val)
                        table.cell(r_idx + 1, c_idx).text_frame.paragraphs[0].font.size = Pt(12)

            ph4 = find_table_placeholder_by_name(summary_slide, "Table Placeholder 4")
            if ph4:
                table_prev = insert_table_at_placeholder(summary_slide, "Table Placeholder 4", len(previous_summary_df) + 1, len(previous_summary_df.columns))
                if table_prev: fill_table_from_df(table_prev, previous_summary_df)

            ph3 = find_table_placeholder_by_name(summary_slide, "Table Placeholder 3")
            if ph3:
                table_curr = insert_table_at_placeholder(summary_slide, "Table Placeholder 3", len(current_summary_df) + 1, len(current_summary_df.columns))
                if table_curr: fill_table_from_df(table_curr, current_summary_df)

            ph1 = find_table_placeholder_by_name(summary_slide, "Table Placeholder 1")
            if ph1:
                table_comp = insert_table_at_placeholder(summary_slide, "Table Placeholder 1", len(summary_df) + 1, len(summary_df.columns))
                if table_comp: fill_table_from_df(table_comp, summary_df)

        # ============================
        # Overall Assessment slide (Table Placeholder 1)
        # ============================
        overall_slide = choose_slide_for_section(
            prefer_titles=["Overall Assessment", "Overall MRUM Assessment"],
            required_placeholders=("Table Placeholder 1",),
            exclude_ids=used_slide_ids
        )
        if overall_slide:
            used_slide_ids.add(id(overall_slide))
            headers = ['Metric', '# of Apps Improved', '# Apps Degraded', 'Overall Result', 'Percentage Value']
            table = insert_table_at_placeholder(overall_slide, "Table Placeholder 1", 2, len(headers))
            if table:
                for i, h in enumerate(headers):
                    table.cell(0, i).text = h
                    table.cell(0, i).text_frame.paragraphs[0].font.size = Pt(14)

                overall_result = 'Increase' if oa_up > oa_down else 'Decrease' if oa_down > oa_up else 'Even'
                percentage_value = 0 if overall_result == 'Even' else round((oa_up / max(1, oa_up + oa_down)) * 100)

                table.cell(1, 0).text = 'Overall MRUM Assessment'
                table.cell(1, 1).text = str(oa_up)
                table.cell(1, 2).text = str(oa_down)
                table.cell(1, 3).text = overall_result
                table.cell(1, 4).text = f"{percentage_value}%"

                run = table.cell(1, 4).text_frame.paragraphs[0].runs[0] if table.cell(1, 4).text_frame.paragraphs[0].runs else table.cell(1, 4).text_frame.paragraphs[0].add_run()
                if overall_result == "Increase":
                    run.font.color.rgb = RGBColor(0, 176, 80)
                elif overall_result == "Decrease":
                    run.font.color.rgb = RGBColor(192, 0, 0)

        # ============================
        # MRUM Entity Comparison (Table Placeholder 1)
        # ============================
        def result_and_percentage(up, down):
            if up > down:
                return "Increase", round((up / max(1, up + down)) * 100)
            if down > up:
                return "Decrease", round((up / max(1, up + down)) * 100)
            return "Even", 0

        entity_slide = choose_slide_for_section(
            prefer_titles=["MRUM Comparison", "Entity Comparison", "Comparison Result", "MRUM Entity Comparison"],
            required_placeholders=("Table Placeholder 1",),
            exclude_ids=used_slide_ids
        )
        if not entity_slide and len(prs.slides) > 7:
            entity_slide = prs.slides[7]
            if not find_table_placeholder_by_name(entity_slide, "Table Placeholder 1"):
                for s in prs.slides:
                    if id(s) in used_slide_ids:
                        continue
                    if find_table_placeholder_by_name(s, "Table Placeholder 1"):
                        entity_slide = s
                        break

        if entity_slide and find_table_placeholder_by_name(entity_slide, "Table Placeholder 1"):
            used_slide_ids.add(id(entity_slide))
            headers = ['Metric', '# of Apps Improved', '# Apps Degraded', 'Overall Result', 'Percentage Value']
            net_result, net_pct = result_and_percentage(net_up, net_down)
            hra_result, hra_pct = result_and_percentage(hra_up, hra_down)
            rows = [
                ['NetworkRequestsMRUM', str(net_up), str(net_down), net_result, f"{net_pct}%"],
                ['HealthRulesAndAlertingMRUM', str(hra_up), str(hra_down), hra_result, f"{hra_pct}%"]
            ]
            table = insert_table_at_placeholder(entity_slide, "Table Placeholder 1", len(rows) + 1, len(headers))
            if table:
                for i, h in enumerate(headers):
                    cell = table.cell(0, i); cell.text = h
                    p = cell.text_frame.paragraphs[0]; p.font.bold = True; p.font.size = Pt(12)
                for r_idx, row in enumerate(rows, start=1):
                    for c_idx, value in enumerate(row):
                        cell = table.cell(r_idx, c_idx); cell.text = value
                        p = cell.text_frame.paragraphs[0]; p.font.size = Pt(12)
                        if c_idx == 4:
                            result_text = rows[r_idx - 1][3]
                            run = p.runs[0] if p.runs else p.add_run()
                            if result_text == "Increase":
                                run.font.color.rgb = RGBColor(0, 176, 80)
                            elif result_text == "Decrease":
                                run.font.color.rgb = RGBColor(192, 0, 0)

        # ============================
        # Slide 11 equivalent: MRUM Network Requests Deep Dive
        # ============================
        deep_dive_slide = prs.slides[10] if len(prs.slides) > 10 else None
        if deep_dive_slide is None or not find_table_placeholder_by_name(deep_dive_slide, "Table Placeholder 1"):
            deep_dive_slide = choose_slide_for_section(
                prefer_titles=["Mobile Network Requests Deep Dive", "MRUM Network Requests", "Mobile Network Requests"],
                required_placeholders=("Table Placeholder 1",),
                exclude_ids=used_slide_ids
            )
        if deep_dive_slide:
            used_slide_ids.add(id(deep_dive_slide))
        else:
            logging.warning("[MRUM] Network Requests Deep Dive slide not found.")
            deep_dive_slide = None

        def set_shape_text(slide, shape_name, text):
            if slide is None:
                return False
            for shp in slide.shapes:
                if getattr(shp, "name", "") == shape_name and hasattr(shp, "text_frame") and shp.text_frame:
                    shp.text_frame.clear()
                    shp.text_frame.text = str(text)
                    return True
            logging.debug("[MRUM] Shape '%s' not found on deep dive slide.", shape_name)
            return False

        def first_present_col(df, candidates):
            for c in candidates:
                if c in df.columns:
                    return c
            return None

        # Decline detectors
        def parse_transition_tokens(val):
            s = str(val or "").strip()
            if "→" in s:
                prev, curr = s.split("→", 1)
                return prev.strip(), curr.strip().split("(")[0].strip()
            return None, None

        def token_to_bool(tok):
            t = str(tok or "").strip().lower()
            if t in {"true", "yes", "y", "1"}: return True
            if t in {"false", "no", "n", "0"}: return False
            return None

        def token_to_num(tok):
            try: return float(str(tok).strip())
            except Exception: return None

        def is_bool_decline_cell(val):
            prev, curr = parse_transition_tokens(val)
            if prev is not None and curr is not None:
                pb = token_to_bool(prev); cb = token_to_bool(curr)
                if pb is True and cb is False: return True
            s = str(val or "").lower()
            return ("declined" in s or "downgraded" in s or "decreased" in s or "reduced" in s) and "false" in s

        def is_num_decline_cell(val):
            prev, curr = parse_transition_tokens(val)
            if prev is not None and curr is not None:
                pn = token_to_num(prev); cn = token_to_num(curr)
                if pn is not None and cn is not None and cn < pn: return True
            s = str(val or "").lower()
            return any(k in s for k in ("declined", "decreased", "reduced", "down", "↓"))

        # Rectangles 8..12
        if deep_dive_slide:
            col_collecting   = first_present_col(df_network_requests, ["collectingDataPastOneDay", "CollectingDataPastOneDay"])
            col_limit_nothit = first_present_col(df_network_requests, ["networkRequestLimitNotHit", "NetworkRequestLimitNotHit"])
            col_custom_rules = first_present_col(df_network_requests, ["numberCustomMatchRules", "NumberCustomMatchRules"])
            col_bt_corr      = first_present_col(df_network_requests, ["hasBtCorrelation", "HasBTCorrelation"])
            col_ces_include  = first_present_col(df_network_requests, ["hasCustomEventServiceIncludeRule", "HasCustomEventServiceIncludeRule"])

            def count_metric_declines(df, metric_col, is_bool):
                if not metric_col: return 0
                cnt = 0
                for _, r in df.iterrows():
                    val = r.get(metric_col, "")
                    if (is_bool and is_bool_decline_cell(val)) or ((not is_bool) and is_num_decline_cell(val)):
                        cnt += 1
                return cnt

            set_shape_text(deep_dive_slide, "Rectangle 8",  str(count_metric_declines(df_network_requests, col_collecting,   True)))
            set_shape_text(deep_dive_slide, "Rectangle 9",  str(count_metric_declines(df_network_requests, col_limit_nothit, True)))
            set_shape_text(deep_dive_slide, "Rectangle 10", str(count_metric_declines(df_network_requests, col_custom_rules, False)))
            set_shape_text(deep_dive_slide, "Rectangle 11", str(count_metric_declines(df_network_requests, col_bt_corr,      True)))
            set_shape_text(deep_dive_slide, "Rectangle 12", str(count_metric_declines(df_network_requests, col_ces_include,  True)))

        # Declined-only table for Network Requests (like BRUM)
        if deep_dive_slide and find_table_placeholder_by_name(deep_dive_slide, "Table Placeholder 1"):
            all_grades = ['platinum', 'gold', 'silver', 'bronze']
            table_grades = ['Gold', 'Silver', 'Bronze']

            def norm_grade(s):
                s = str(s).strip().lower()
                for g in all_grades:
                    if g in s: return g
                return None

            def parse_transition(val):
                s = str(val or "").strip()
                if not s: return (None, None)
                if '→' in s:
                    parts = s.split('→', 1)
                    return norm_grade(parts[0]), norm_grade(parts[1])
                low = s.lower()
                import re
                m = re.search(r'(?:declined|downgraded)\s+(?:to\s+)?(platinum|gold|silver|bronze)', low)
                if m: return (None, m.group(1))
                m = re.search(r'(?:now|is\s+now|became)\s+(platinum|gold|silver|bronze)', low)
                if m: return (None, m.group(1))
                g = norm_grade(low)
                return (None, g) if g else (None, None)

            app_col_nr = first_present_col(df_network_requests, ["name", "Name", "applicationName", "Application Name", "Application"]) \
                         or next((c for c in df_network_requests.columns if str(c).lower() in ("app", "application", "application name")), None)

            def row_for_app(app):
                if not app_col_nr: return None
                m = df_network_requests[df_network_requests[app_col_nr].astype(str).str.strip() == str(app)]
                return m.iloc[0] if not m.empty else None

            totals_by_grade = {g: 0 for g in table_grades}
            grade_by_app = {}
            if "NetworkRequestsMRUM" in df_analysis.columns:
                for _, r in df_analysis.iterrows():
                    app = str(r.get("name", "") or "").strip()
                    _, cg = parse_transition(r.get("NetworkRequestsMRUM", ""))
                    disp = cg.capitalize() if cg else None
                    if app and disp in totals_by_grade:
                        totals_by_grade[disp] += 1
                        grade_by_app[app] = disp

            declined_from_analysis = set()
            if "NetworkRequestsMRUM" in df_analysis.columns:
                for _, r in df_analysis.iterrows():
                    app = str(r.get("name", "") or "").strip()
                    val = r.get("NetworkRequestsMRUM", "")
                    prev_g, curr_g = parse_transition(val)
                    low = str(val).lower()
                    is_dg = ('declined' in low or 'downgraded' in low) or (
                        prev_g and curr_g and prev_g in all_grades and curr_g in all_grades and all_grades.index(prev_g) < all_grades.index(curr_g)
                    )
                    if is_dg and app:
                        declined_from_analysis.add(app)

            metric_cols_and_types = [
                (first_present_col(df_network_requests, ["collectingDataPastOneDay", "CollectingDataPastOneDay"]), True),
                (first_present_col(df_network_requests, ["networkRequestLimitNotHit", "NetworkRequestLimitNotHit"]), True),
                (first_present_col(df_network_requests, ["numberCustomMatchRules", "NumberCustomMatchRules"]), False),
                (first_present_col(df_network_requests, ["hasBtCorrelation", "HasBTCorrelation"]), True),
                (first_present_col(df_network_requests, ["hasCustomEventServiceIncludeRule", "HasCustomEventServiceIncludeRule"]), True),
            ]
            metric_declined_set = set()
            if app_col_nr:
                for col, is_bool in metric_cols_and_types:
                    if not col: continue
                    for _, r in df_network_requests.iterrows():
                        app = str(r.get(app_col_nr, "") or "").strip()
                        val = r.get(col, "")
                        if (is_bool and is_bool_decline_cell(val)) or ((not is_bool) and is_num_decline_cell(val)):
                            metric_declined_set.add(app)

            declined_union = declined_from_analysis.union(metric_declined_set)

            def resolve_grade_for_app(app):
                dest = grade_by_app.get(app)
                if dest: return dest
                r = row_for_app(app)
                if r is not None:
                    for v in r.values:
                        g = norm_grade(v)
                        if g: return g.capitalize()
                if "NetworkRequestsMRUM" in df_analysis.columns:
                    txt = df_analysis.loc[df_analysis["name"].astype(str).str.strip() == str(app), "NetworkRequestsMRUM"]
                    if not txt.empty:
                        _, cg = parse_transition(txt.iloc[0])
                        if cg: return cg.capitalize()
                return None

            declined_by_grade = {g: [] for g in table_grades}
            missing_grade = []
            for app in sorted(declined_union):
                dest = resolve_grade_for_app(app)
                if dest in declined_by_grade:
                    declined_by_grade[dest].append(app)
                else:
                    missing_grade.append(app)
            if missing_grade:
                logging.warning("[MRUM][Slide11] %d declined apps have no resolvable grade: %s", len(missing_grade), missing_grade[:20])

            headers = ["Grade", "Application Names", "Number of Applications", "Percentage Declined"]
            rows = []
            for g in table_grades:
                names = sorted(declined_by_grade[g], key=str.lower)
                num_apps = len(names)
                denom = totals_by_grade.get(g, 0)
                pct = round((num_apps / denom) * 100) if denom > 0 else 0
                rows.append([g, "\n".join(names) if names else "—", str(num_apps), f"{pct}%"])

            table = insert_table_at_placeholder(deep_dive_slide, "Table Placeholder 1", len(rows) + 1, len(headers))
            if table:
                for i, h in enumerate(headers):
                    cell = table.cell(0, i); cell.text = h
                    p = cell.text_frame.paragraphs[0]; p.font.bold = True; p.font.size = Pt(12)
                for r_idx, row in enumerate(rows, start=1):
                    for c_idx, val in enumerate(row):
                        cell = table.cell(r_idx, c_idx); cell.text = val
                        p = cell.text_frame.paragraphs[0]; p.font.size = Pt(12)
                        if c_idx == 3:
                            run = p.runs[0] if p.runs else p.add_run()
                            pct_num = int(str(val).replace("%", "") or 0)
                            run.font.color.rgb = RGBColor(192, 0, 0) if pct_num > 0 else RGBColor(0, 176, 80)

        # ============================
        # Slide 12 equivalent: MRUM Health Rules & Alerting Deep Dive
        # ============================
        hra_slide = prs.slides[11] if len(prs.slides) > 11 else None
        if hra_slide is None or not find_table_placeholder_by_name(hra_slide, "Table Placeholder 1"):
            hra_slide = choose_slide_for_section(
                prefer_titles=["Mobile Health Rules & Alerting Deep Dive", "MRUM Health Rules & Alerting", "Mobile Health Rules"],
                required_placeholders=("Table Placeholder 1",),
                exclude_ids=used_slide_ids
            )
        if hra_slide:
            used_slide_ids.add(id(hra_slide))
        else:
            logging.warning("[MRUM] Health Rules & Alerting Deep Dive slide not found.")
            hra_slide = None

        # Rectangles 10..12
        if hra_slide:
            col_violations    = first_present_col(df_health_rules, ["numberOfHealthRuleViolations", "NumberOfHealthRuleViolations"])
            col_actions_bound = first_present_col(df_health_rules, ["numberOfActionsBoundToEnabledPolicies", "NumberOfActionsBoundToEnabledPolicies"])
            col_custom_rules  = first_present_col(df_health_rules, ["numberOfCustomHealthRules", "NumberOfCustomHealthRules"])

            def count_metric_declines_num(df, metric_col):
                if not metric_col: return 0
                cnt = 0
                for _, r in df.iterrows():
                    val = r.get(metric_col, "")
                    if is_num_decline_cell(val):
                        cnt += 1
                return cnt

            set_shape_text(hra_slide, "Rectangle 10", str(count_metric_declines_num(df_health_rules, col_violations)))
            set_shape_text(hra_slide, "Rectangle 11", str(count_metric_declines_num(df_health_rules, col_actions_bound)))
            set_shape_text(hra_slide, "Rectangle 12", str(count_metric_declines_num(df_health_rules, col_custom_rules)))

        # Declined-only table for HRA
        if hra_slide and find_table_placeholder_by_name(hra_slide, "Table Placeholder 1"):
            all_grades = ['platinum', 'gold', 'silver', 'bronze']
            table_grades = ['Gold', 'Silver', 'Bronze']

            def norm_grade(s):
                s = str(s).strip().lower()
                for g in all_grades:
                    if g in s: return g
                return None

            def parse_transition(val):
                s = str(val or "").strip()
                if not s: return (None, None)
                if '→' in s:
                    parts = s.split('→', 1)
                    return norm_grade(parts[0]), norm_grade(parts[1])
                low = s.lower()
                import re
                m = re.search(r'(?:declined|downgraded)\s+(?:to\s+)?(platinum|gold|silver|bronze)', low)
                if m: return (None, m.group(1))
                m = re.search(r'(?:now|is\s+now|became)\s+(platinum|gold|silver|bronze)', low)
                if m: return (None, m.group(1))
                g = norm_grade(low)
                return (None, g) if g else (None, None)

            app_col_hr = first_present_col(df_health_rules, ["name", "Name", "applicationName", "Application Name", "Application"]) \
                         or next((c for c in df_health_rules.columns if str(c).lower() in ("app", "application", "application name")), None)

            def row_for_app_hr(app):
                if not app_col_hr: return None
                m = df_health_rules[df_health_rules[app_col_hr].astype(str).str.strip() == str(app)]
                return m.iloc[0] if not m.empty else None

            totals_by_grade_hr = {g: 0 for g in table_grades}
            grade_by_app_hr = {}
            if "HealthRulesAndAlertingMRUM" in df_analysis.columns:
                for _, r in df_analysis.iterrows():
                    app = str(r.get("name", "") or "").strip()
                    _, cg = parse_transition(r.get("HealthRulesAndAlertingMRUM", ""))
                    disp = cg.capitalize() if cg else None
                    if app and disp in totals_by_grade_hr:
                        totals_by_grade_hr[disp] += 1
                        grade_by_app_hr[app] = disp

            declined_from_analysis_hr = set()
            if "HealthRulesAndAlertingMRUM" in df_analysis.columns:
                for _, r in df_analysis.iterrows():
                    app = str(r.get("name", "") or "").strip()
                    val = r.get("HealthRulesAndAlertingMRUM", "")
                    prev_g, curr_g = parse_transition(val)
                    low = str(val).lower()
                    is_dg = ('declined' in low or 'downgraded' in low) or (
                        prev_g and curr_g and prev_g in all_grades and curr_g in all_grades and all_grades.index(prev_g) < all_grades.index(curr_g)
                    )
                    if is_dg and app:
                        declined_from_analysis_hr.add(app)

            metric_cols_hr = [
                first_present_col(df_health_rules, ["numberOfHealthRuleViolations", "NumberOfHealthRuleViolations"]),
                first_present_col(df_health_rules, ["numberOfActionsBoundToEnabledPolicies", "NumberOfActionsBoundToEnabledPolicies"]),
                first_present_col(df_health_rules, ["numberOfCustomHealthRules", "NumberOfCustomHealthRules"]),
            ]
            metric_declined_apps_hr = set()
            if app_col_hr:
                for col in metric_cols_hr:
                    if not col: continue
                    for _, r in df_health_rules.iterrows():
                        app = str(r.get(app_col_hr, "") or "").strip()
                        val = r.get(col, "")
                        if is_num_decline_cell(val):
                            metric_declined_apps_hr.add(app)

            declined_union_hr = declined_from_analysis_hr.union(metric_declined_apps_hr)

            def resolve_grade_for_app_hr(app):
                dest = grade_by_app_hr.get(app)
                if dest: return dest
                r = row_for_app_hr(app)
                if r is not None:
                    for v in r.values:
                        g = norm_grade(v)
                        if g: return g.capitalize()
                if "HealthRulesAndAlertingMRUM" in df_analysis.columns:
                    txt = df_analysis.loc[df_analysis["name"].astype(str).str.strip() == str(app), "HealthRulesAndAlertingMRUM"]
                    if not txt.empty:
                        _, cg = parse_transition(txt.iloc[0])
                        if cg: return cg.capitalize()
                return None

            declined_by_grade_hr = {g: [] for g in table_grades}
            missing_grade_hr = []
            for app in sorted(declined_union_hr):
                dest = resolve_grade_for_app_hr(app)
                if dest in declined_by_grade_hr:
                    declined_by_grade_hr[dest].append(app)
                else:
                    missing_grade_hr.append(app)
            if missing_grade_hr:
                logging.warning("[MRUM][Slide12] %d declined apps have no resolvable grade: %s", len(missing_grade_hr), missing_grade_hr[:20])

            headers = ["Grade", "Application Names", "Number of Applications", "Percentage Declined"]
            rows = []
            for g in table_grades:
                names = sorted(declined_by_grade_hr[g], key=str.lower)
                num_apps = len(names)
                denom = totals_by_grade_hr.get(g, 0)
                pct = round((num_apps / denom) * 100) if denom > 0 else 0
                rows.append([g, "\n".join(names) if names else "—", str(num_apps), f"{pct}%"])

            table = insert_table_at_placeholder(hra_slide, "Table Placeholder 1", len(rows) + 1, len(headers))
            if table:
                for i, h in enumerate(headers):
                    cell = table.cell(0, i); cell.text = h
                    p = cell.text_frame.paragraphs[0]; p.font.bold = True; p.font.size = Pt(12)
                for r_idx, row in enumerate(rows, start=1):
                    for c_idx, val in enumerate(row):
                        cell = table.cell(r_idx, c_idx); cell.text = val
                        p = cell.text_frame.paragraphs[0]; p.font.size = Pt(12)
                        if c_idx == 3:
                            run = p.runs[0] if p.runs else p.add_run()
                            pct_num = int(str(val).replace("%", "") or 0)
                            run.font.color.rgb = RGBColor(192, 0, 0) if pct_num > 0 else RGBColor(0, 176, 80)

        # ============================
        # Populate "TextBox 7" with number of MRUM applications (slide index 5)
        # ============================
        def set_textbox_value(prs_obj, shape_name, text, fallback_slide_index=5):
            for s in prs_obj.slides:
                for shp in s.shapes:
                    if getattr(shp, "name", "") == shape_name and hasattr(shp, "text_frame"):
                        shp.text_frame.clear()
                        shp.text_frame.text = str(text)
                        return True
            if len(prs_obj.slides) > fallback_slide_index:
                s = prs_obj.slides[fallback_slide_index]
                for shp in s.shapes:
                    if getattr(shp, "name", "") == shape_name and hasattr(shp, "text_frame"):
                        shp.text_frame.clear()
                        shp.text_frame.text = str(text)
                        return True
            return False

        if not set_textbox_value(prs, "TextBox 7", number_of_apps, fallback_slide_index=5):
            logging.warning("[MRUM] 'TextBox 7' not found; MRUM application count not written.")

        # Save once at the end.
        prs.save(powerpoint_output_path)
        logging.debug(f"[MRUM] PowerPoint saved to: {powerpoint_output_path}")

    except Exception as e:
        logging.error(f"[MRUM] Error generating PowerPoint: {e}", exc_info=True)
        raise


@app.route('/upload', methods=['POST'])
def upload():
    logging.debug("Request files: %s", request.files)
    
    if 'previous_file' not in request.files or 'current_file' not in request.files:
        logging.error("No file part")
        return render_template('index.html', message="Error: No file part was uploaded."), 400
    
    previous_file = request.files.get('previous_file')
    current_file = request.files.get('current_file')
    
    if not previous_file or previous_file.filename == '':
        logging.error("No selected file for previous")
        return render_template('index.html', message="Error: No previous file was selected."), 400
    
    if not current_file or current_file.filename == '':
        logging.error("No selected file for current")
        return render_template('index.html', message="Error: No current file was selected."), 400
    
    previous_file_path = os.path.join(app.config['UPLOAD_FOLDER'], 'previous.xlsx')
    current_file_path = os.path.join(app.config['UPLOAD_FOLDER'], 'current.xlsx')
    # output_file_path = os.path.join(app.config['RESULT_FOLDER'], 'comparison_result.xlsx')
    previous_sum_path = os.path.join(app.config['UPLOAD_FOLDER'], 'previous_sum.xlsx')
    current_sum_path = os.path.join(app.config['UPLOAD_FOLDER'], 'current_sum.xlsx')
    comparison_sum_path = os.path.join(app.config['RESULT_FOLDER'], 'comparison_sum.xlsx')
    # powerpoint_output_path = os.path.join(app.config['RESULT_FOLDER'], 'Analysis_Summary.pptx')

    ts = ts_now()
    domain = "APM"
    output_file_path = os.path.join(app.config['RESULT_FOLDER'], f"comparison_result_{domain.lower()}_{ts}.xlsx")
    powerpoint_output_path = os.path.join(app.config['RESULT_FOLDER'], f"analysis_summary_{domain.lower()}_{ts}.pptx")

    try:
        # Save the uploaded files
        previous_file.save(previous_file_path)
        current_file.save(current_file_path)

        # Automatically save the workbooks to recalculate formulas
        save_workbook(previous_file_path)
        save_workbook(current_file_path)

        # Check if controllers match
        if not check_controllers_match(previous_file_path, current_file_path):
            logging.error("Controllers do not match.")
            return render_template('index.html', message="Error: The controllers in the two files do not match. Please upload files from the same controller."), 400

        # Proceed with comparison...
        create_summary_workbooks(previous_file_path, current_file_path, previous_sum_path, current_sum_path)
        compare_files_summary(previous_sum_path, current_sum_path, comparison_sum_path)
        compare_files_other_sheets(previous_file_path, current_file_path, output_file_path)
        copy_summary_to_result(comparison_sum_path, output_file_path)

        # Build APM PowerPoint.
        generate_powerpoint_from_analysis(output_file_path, powerpoint_output_path, current_file_path, previous_file_path)

        # Build compact JSON summary.
        json_path, json_name, _ = build_comparison_json(
            domain="APM",
            comparison_result_path=output_file_path,
            current_file_path=current_file_path,
            previous_file_path=previous_file_path,
            result_folder=app.config['RESULT_FOLDER']
        )

        app.config['LAST_RESULT_APM'] = output_file_path
        app.config['LAST_JSON_APM'] = json_path

        # Optional internal workflow (guarded).
        CIRCUIT_ENABLED = os.getenv("CIRCUIT_ENABLED", "false").lower() == "true"
        if CIRCUIT_ENABLED:
            logging.info("[APM] CIRCUIT_ENABLED: prepared JSON at %s", json_path)
            # send_to_circuit(json_path)  # Implement internally.

        return render_template(
            'index.html', 
            message=(
                "APM comparison completed successfully.<br><br>"
                f"Download the APM results "
                f"<a href='/download/{os.path.basename(output_file_path)}' style='color: #32CD32;'>here</a>, "
                f"the APM PowerPoint "
                f"<a href='/download/{os.path.basename(powerpoint_output_path)}' style='color: #32CD32;'>here</a>, "
                f"and the APM JSON "
                f"<a href='/download/{json_name}' style='color: #32CD32;'>here</a>."
            )
        )
    except Exception as e:
        logging.error(f"Error during file upload or comparison: {e}", exc_info=True)
        return render_template('index.html', message="Error during file upload or comparison"), 500
    
@app.route('/upload_brum', methods=['POST'])
def upload_brum():
    logging.debug("[BRUM] Request files: %s", request.files)

    if 'previous_brum' not in request.files or 'current_brum' not in request.files:
        logging.error("[BRUM] No file part")
        return render_template('index.html', message="Error: No BRUM file was uploaded."), 400

    previous_file = request.files.get('previous_brum')
    current_file = request.files.get('current_brum')

    if not previous_file or previous_file.filename == '':
        logging.error("[BRUM] No selected file for previous")
        return render_template('index.html', message="Error: No previous BRUM file was selected."), 400

    if not current_file or current_file.filename == '':
        logging.error("[BRUM] No selected file for current")
        return render_template('index.html', message="Error: No current BRUM file was selected."), 400

    # Fixed upload paths (intentionally overwrite the working inputs).
    previous_file_path = os.path.join(app.config['UPLOAD_FOLDER'], 'previous_brum.xlsx')
    current_file_path  = os.path.join(app.config['UPLOAD_FOLDER'], 'current_brum.xlsx')

    try:
        # Save uploads.
        previous_file.save(previous_file_path)
        current_file.save(current_file_path)

        # Recalculate formulas for both workbooks.
        save_workbook(previous_file_path)
        save_workbook(current_file_path)

        # Build token based on controller + dates (controller.prevdate.currdate.comparedate).
        controller_slug, prev_date, curr_date, compare_date = build_compare_token(previous_file_path, current_file_path)
        domain = "BRUM"
        base = f"{controller_slug}.{prev_date}.{curr_date}.{compare_date}"

        # Result artifact paths (stamped).
        output_file_path    = os.path.join(app.config['RESULT_FOLDER'], f"comparison_result_{domain.lower()}_{base}.xlsx")
        previous_sum_path   = os.path.join(app.config['UPLOAD_FOLDER'],  f"previous_{domain.lower()}_sum.xlsx")
        current_sum_path    = os.path.join(app.config['UPLOAD_FOLDER'],  f"current_{domain.lower()}_sum.xlsx")
        comparison_sum_path = os.path.join(app.config['RESULT_FOLDER'], f"comparison_sum_{domain.lower()}_{base}.xlsx")
        pptx_out            = os.path.join(app.config['RESULT_FOLDER'], f"analysis_summary_{domain.lower()}_{base}.pptx")

        # Build Summary workbooks and comparison sheet.
        create_summary_workbooks(previous_file_path, current_file_path, previous_sum_path, current_sum_path)
        compare_files_summary(previous_sum_path, current_sum_path, comparison_sum_path)

        # BRUM compare for all sheets.
        compare_files_other_sheets_brum(previous_file_path, current_file_path, output_file_path)

        # Put Summary first in the output workbook.
        copy_summary_to_result(comparison_sum_path, output_file_path)

        # Build BRUM PowerPoint.
        generate_powerpoint_from_brum(
            comparison_result_path=output_file_path,
            powerpoint_output_path=pptx_out,
            current_file_path=current_file_path,
            previous_file_path=previous_file_path
        )

        # NEW: pass meta directly to the JSON builder.
        meta = {
            "domain": domain,
            "controller": controller_slug,
            "previousDate": prev_date,
            "currentDate": curr_date,
            "compareDate": compare_date
        }

        # Build compact JSON summary.
        json_path, json_name, _ = build_comparison_json(
            domain=domain,
            comparison_result_path=output_file_path,
            current_file_path=current_file_path,
            previous_file_path=previous_file_path,
            result_folder=app.config['RESULT_FOLDER']
        )

        # Rename JSON to stamped name if needed.
        stamped_json_name = f"summary_{domain.lower()}_{base}.json"
        stamped_json_path = os.path.join(app.config['RESULT_FOLDER'], stamped_json_name)
        if os.path.basename(json_path) != stamped_json_name:
            try:
                os.replace(json_path, stamped_json_path)
                json_path, json_name = stamped_json_path, stamped_json_name
            except Exception as e:
                logging.warning("[BRUM] Could not rename JSON to stamped name: %s", e)

        # Inject meta into JSON (domain/controller/dates).
        try:
            with open(json_path, "r+", encoding="utf-8") as f:
                data = json.load(f)
                data["meta"] = {
                    "domain": domain,
                    "controller": controller_slug,
                    "previousDate": prev_date,
                    "currentDate": curr_date,
                    "compareDate": compare_date
                }
                f.seek(0)
                json.dump(data, f, indent=2)
                f.truncate()
        except Exception as e:
            logging.warning("[BRUM] Failed to inject meta into JSON: %s", e)

        # Record latest artifacts for Insights.
        app.config['LAST_RESULT_BRUM'] = output_file_path
        app.config['LAST_JSON_BRUM']   = json_path

        # Optional internal workflow (guarded).
        CIRCUIT_ENABLED = os.getenv("CIRCUIT_ENABLED", "false").lower() == "true"
        if CIRCUIT_ENABLED:
            logging.info("[BRUM] CIRCUIT_ENABLED: prepared JSON at %s", json_path)

        # Use the final (possibly renamed) JSON basename in the message.
        json_name = os.path.basename(json_path)

        return render_template(
            'index.html',
            message=(
                "BRUM comparison completed successfully.<br><br>"
                f"Download the BRUM results "
                f"<a href='/download/{os.path.basename(output_file_path)}' style='color: #32CD32;'>here</a>, "
                f"the BRUM PowerPoint "
                f"<a href='/download/{os.path.basename(pptx_out)}' style='color: #32CD32;'>here</a>, "
                f"and the BRUM JSON "
                f"<a href='/download/{json_name}' style='color: #32CD32;'>here</a>."
            )
        )
    except Exception as e:
        logging.error(f"[BRUM] Error during upload or comparison: {e}", exc_info=True)
        return render_template('index.html', message="Error during BRUM upload or comparison"), 500

@app.route('/upload_mrum', methods=['POST'])
def upload_mrum():
    logging.debug("[MRUM] Request files: %s", request.files)

    if 'previous_mrum' not in request.files or 'current_mrum' not in request.files:
        logging.error("[MRUM] No file part")
        return render_template('index.html', message="Error: No MRUM file was uploaded."), 400

    previous_file = request.files.get('previous_mrum')
    current_file  = request.files.get('current_mrum')

    if not previous_file or previous_file.filename == '':
        logging.error("[MRUM] No selected file for previous")
        return render_template('index.html', message="Error: No previous MRUM file was selected."), 400

    if not current_file or current_file.filename == '':
        logging.error("[MRUM] No selected file for current")
        return render_template('index.html', message="Error: No current MRUM file was selected."), 400

    # Fixed upload paths (intentionally overwrite the working inputs).
    previous_file_path = os.path.join(app.config['UPLOAD_FOLDER'], 'previous_mrum.xlsx')
    current_file_path  = os.path.join(app.config['UPLOAD_FOLDER'], 'current_mrum.xlsx')

    try:
        # Save uploads.
        previous_file.save(previous_file_path)
        current_file.save(current_file_path)

        # Recalculate formulas for both workbooks.
        save_workbook(previous_file_path)
        save_workbook(current_file_path)

        # Build token based on controller + dates (controller.prevdate.currdate.comparedate).
        controller_slug, prev_date, curr_date, compare_date = build_compare_token(previous_file_path, current_file_path)
        domain = "MRUM"
        base = f"{controller_slug}.{prev_date}.{curr_date}.{compare_date}"

        # Result artifact paths (stamped).
        output_file_path    = os.path.join(app.config['RESULT_FOLDER'], f"comparison_result_{domain.lower()}_{base}.xlsx")
        previous_sum_path   = os.path.join(app.config['UPLOAD_FOLDER'],  f"previous_{domain.lower()}_sum.xlsx")
        current_sum_path    = os.path.join(app.config['UPLOAD_FOLDER'],  f"current_{domain.lower()}_sum.xlsx")
        comparison_sum_path = os.path.join(app.config['RESULT_FOLDER'], f"comparison_sum_{domain.lower()}_{base}.xlsx")
        pptx_out            = os.path.join(app.config['RESULT_FOLDER'], f"analysis_summary_{domain.lower()}_{base}.pptx")

        # Build Summary workbooks and comparison sheet.
        create_summary_workbooks(previous_file_path, current_file_path, previous_sum_path, current_sum_path)
        compare_files_summary(previous_sum_path, current_sum_path, comparison_sum_path)

        # MRUM compare for all sheets.
        compare_files_other_sheets_mrum(previous_file_path, current_file_path, output_file_path)

        # Put Summary first in the output workbook.
        copy_summary_to_result(comparison_sum_path, output_file_path)

        # Build MRUM PowerPoint.
        generate_powerpoint_from_mrum(
            comparison_result_path=output_file_path,
            powerpoint_output_path=pptx_out,
            current_file_path=current_file_path,
            previous_file_path=previous_file_path
        )

        # NEW: pass meta directly to the JSON builder.
        meta = {
            "domain": domain,
            "controller": controller_slug,
            "previousDate": prev_date,
            "currentDate": curr_date,
            "compareDate": compare_date
        }

        # Build compact JSON summary.
        json_path, json_name, _ = build_comparison_json(
            domain=domain,
            comparison_result_path=output_file_path,
            current_file_path=current_file_path,
            previous_file_path=previous_file_path,
            result_folder=app.config['RESULT_FOLDER']
        )

        # Rename JSON to stamped name if needed.
        stamped_json_name = f"summary_{domain.lower()}_{base}.json"
        stamped_json_path = os.path.join(app.config['RESULT_FOLDER'], stamped_json_name)
        if os.path.basename(json_path) != stamped_json_name:
            try:
                os.replace(json_path, stamped_json_path)
                json_path, json_name = stamped_json_path, stamped_json_name
            except Exception as e:
                logging.warning("[MRUM] Could not rename JSON to stamped name: %s", e)

        # Inject meta into JSON (domain/controller/dates).
        try:
            with open(json_path, "r+", encoding="utf-8") as f:
                data = json.load(f)
                data["meta"] = {
                    "domain": domain,
                    "controller": controller_slug,
                    "previousDate": prev_date,
                    "currentDate": curr_date,
                    "compareDate": compare_date
                }
                f.seek(0)
                json.dump(data, f, indent=2)
                f.truncate()
        except Exception as e:
            logging.warning("[MRUM] Failed to inject meta into JSON: %s", e)

        # Record latest artifacts for Insights.
        app.config['LAST_RESULT_MRUM'] = output_file_path
        app.config['LAST_JSON_MRUM']   = json_path

        # Optional internal workflow (guarded).
        CIRCUIT_ENABLED = os.getenv("CIRCUIT_ENABLED", "false").lower() == "true"
        if CIRCUIT_ENABLED:
            logging.info("[MRUM] CIRCUIT_ENABLED: prepared JSON at %s", json_path)

        # Use the final (possibly renamed) JSON basename in the message.
        json_name = os.path.basename(json_path)

        return render_template(
            'index.html',
            message=(
                "MRUM comparison completed successfully.<br><br>"
                f"Download the MRUM results "
                f"<a href='/download/{os.path.basename(output_file_path)}' style='color: #32CD32;'>here</a>, "
                f"the MRUM PowerPoint "
                f"<a href='/download/{os.path.basename(pptx_out)}' style='color: #32CD32;'>here</a>, "
                f"and the MRUM JSON "
                f"<a href='/download/{json_name}' style='color: #32CD32;'>here</a>."
            )
        )
    except Exception as e:
        logging.error(f"[MRUM] Error during upload or comparison: {e}", exc_info=True)
        return render_template('index.html', message="Error during MRUM upload or comparison"), 500

@app.route('/download/<filename>')
def download_file(filename):
    # Provide a download link for the output file
    return send_file(os.path.join(app.config['RESULT_FOLDER'], filename), as_attachment=True)


# Define color fills for Excel cells
red_fill = PatternFill(start_color='FF0000', end_color='FF0000', fill_type='solid')
green_fill = PatternFill(start_color='00FF00', end_color='00FF00', fill_type='solid')
added_fill = PatternFill(start_color='ADD8E6', end_color='ADD8E6', fill_type='solid')

# Helper function to find column index by name
def get_key_column(sheet, key_name):
    for i, cell in enumerate(sheet[1], 1):
        if cell.value == key_name:
            return i
    return None

# Function to create summary workbooks
def create_summary_workbooks(previous_file_path, current_file_path, previous_sum_path, current_sum_path):
    try:
        wb_previous = load_workbook(previous_file_path, data_only=True)
        wb_current = load_workbook(current_file_path, data_only=True)

        if 'Summary' not in wb_previous.sheetnames or 'Summary' not in wb_current.sheetnames:
            logging.error("'Summary' sheet is missing in one of the files.")
            return

        ws_previous = wb_previous['Summary']
        ws_current = wb_current['Summary']

        # Create new workbooks for the summaries
        wb_previous_sum = openpyxl.Workbook()
        wb_current_sum = openpyxl.Workbook()

        ws_previous_sum = wb_previous_sum.active
        ws_current_sum = wb_current_sum.active

        ws_previous_sum.title = 'Summary'
        ws_current_sum.title = 'Summary'

        # Copy data from original workbooks to summary workbooks as values only
        for row in ws_previous.iter_rows(values_only=True):
            ws_previous_sum.append(row)
        for row in ws_current.iter_rows(values_only=True):
            ws_current_sum.append(row)

        # Save the cleaned-up summary workbooks
        wb_previous_sum.save(previous_sum_path)
        wb_current_sum.save(current_sum_path)

    except Exception as e:
        logging.error(f"Error in create_summary_workbooks: {e}", exc_info=True)
        raise


# Function to compare 'Summary' sheet and save to a new workbook
def compare_files_summary(previous_sum_path, current_sum_path, comparison_sum_path):
    try:
        # Load the previous_sum and current_sum Excel files
        wb_previous = load_workbook(previous_sum_path, data_only=True)
        wb_current = load_workbook(current_sum_path, data_only=True)
        wb_output = openpyxl.Workbook()

        ws_previous = wb_previous['Summary']
        ws_current = wb_current['Summary']
        ws_output = wb_output.active
        ws_output.title = 'Summary'

        # Debugging: Print sheet names and some data from the 'Summary' sheet
        print("Previous Workbook Sheets:", wb_previous.sheetnames)
        print("Current Workbook Sheets:", wb_current.sheetnames)
        print("Example data from 'Summary' sheet in Previous Workbook (1,1):", ws_previous.cell(row=1, column=1).value)
        print("Example data from 'Summary' sheet in Current Workbook (1,1):", ws_current.cell(row=1, column=1).value)
        
        logging.debug(f"Processing sheet: 'Summary'")
        
        compare_summary(ws_previous, ws_current, ws_output)

        # Save the workbook after all modifications have been completed
        wb_output.save(comparison_sum_path)
        logging.debug(f"Summary comparison saved to: {comparison_sum_path}")

    except Exception as e:
        logging.error(f"Error in compare_files_summary: {e}", exc_info=True)
        raise


# Function to compare summaries of both sheets
def compare_summary(ws_previous, ws_current, ws_output):
    from openpyxl.styles import PatternFill

    # Define fill styles
    red_fill = PatternFill(start_color="FF0000", end_color="FF0000", fill_type="solid")
    green_fill = PatternFill(start_color="00FF00", end_color="00FF00", fill_type="solid")

    for row in ws_previous.iter_rows(min_row=1, min_col=1, max_col=ws_previous.max_column, max_row=ws_previous.max_row):
        for cell in row:
            prev_cell = ws_previous.cell(row=cell.row, column=cell.column)
            curr_cell = ws_current.cell(row=cell.row, column=cell.column)
            output_cell = ws_output.cell(row=cell.row, column=cell.column)

            prev_value = prev_cell.value
            curr_value = curr_cell.value

            if prev_value is None:
                prev_value = ''
            if curr_value is None:
                curr_value = ''

            logging.debug(f"Comparing cell ({cell.row},{cell.column}): Previous Value: {prev_value}, Current Value: {curr_value}")

            if prev_value != curr_value:
                if isinstance(prev_value, (int, float)) and isinstance(curr_value, (int, float)):
                    if curr_value > prev_value:
                        output_cell.fill = green_fill
                    else:
                        output_cell.fill = red_fill
                    output_cell.value = f"{prev_value} → {curr_value}"
                else:
                    output_cell.fill = red_fill
                    output_cell.value = f"{prev_value} → {curr_value}"
            else:
                output_cell.value = prev_value

            logging.debug(f"Cell ({cell.row},{cell.column}) updated to: {output_cell.value}")

# Function to copy the Summary sheet from comparison_sum to comparison_result
def copy_summary_to_result(comparison_sum_path, output_file_path):
    try:
        # Load the comparison_sum and output workbooks
        wb_comparison_sum = load_workbook(comparison_sum_path)
        wb_output = load_workbook(output_file_path)

        # Get the Summary sheet from comparison_sum
        ws_comparison_sum = wb_comparison_sum['Summary']

        # If the Summary sheet already exists in output, delete it
        if 'Summary' in wb_output.sheetnames:
            del wb_output['Summary']

        # Create a new Summary sheet in the output workbook
        ws_output = wb_output.create_sheet('Summary', 0)  # Insert Summary as the first sheet

        # Copy data and formatting from the comparison_sum Summary sheet to the output Summary sheet
        for row in ws_comparison_sum.iter_rows():
            for cell in row:
                new_cell = ws_output.cell(row=cell.row, column=cell.column, value=cell.value)
                
                # Copy individual style attributes
                if cell.has_style:
                    new_cell.font = copy(cell.font)
                    new_cell.border = copy(cell.border)
                    new_cell.fill = copy(cell.fill)
                    new_cell.number_format = cell.number_format
                    new_cell.protection = copy(cell.protection)
                    new_cell.alignment = copy(cell.alignment)

        # Define header colors for Bronze, Silver, Gold, and Platinum
        header_colors = {
            'B1': 'cd7f32',  # Bronze
            'C1': 'c0c0c0',  # Silver
            'D1': 'ffd700',  # Gold
            'E1': 'e5e4e2'   # Platinum
        }

        # Apply colors to the headers in the first row
        for col in header_colors:
            cell = ws_output[col]
            cell.fill = PatternFill(start_color=header_colors[col], end_color=header_colors[col], fill_type="solid")
            cell.font = Font(bold=True, color="000000")

        # Save the output workbook
        wb_output.save(output_file_path)
        logging.debug("Summary sheet copied to the final comparison result and placed as the first sheet with highlighted headers.")

    except Exception as e:
        logging.error(f"Error in copy_summary_to_result: {e}", exc_info=True)
        raise

# Function to compare files for all other sheets
def compare_files_other_sheets(previous_file_path, current_file_path, output_file_path):
    try:
        wb_previous = load_workbook(previous_file_path)
        wb_current = load_workbook(current_file_path)

        for sheet_name in wb_current.sheetnames:
            if sheet_name in wb_previous.sheetnames:
                ws_previous = wb_previous[sheet_name]
                ws_current = wb_current[sheet_name]

                logging.debug(f"Processing sheet: {sheet_name}")

                if sheet_name == 'Analysis':
                    compare_analysis(ws_previous, ws_current)
                elif  sheet_name == 'AppAgentsAPM':
                    compare_appagentsapm(ws_previous, ws_current) 
                elif  sheet_name == 'MachineAgentsAPM':
                    compare_machineagentsapm(ws_previous, ws_current)
                elif sheet_name == 'BusinessTransactionsAPM':
                    compare_businesstransactionsapm(ws_previous, ws_current)
                elif  sheet_name == 'BackendsAPM':
                    compare_backendsapm(ws_previous, ws_current) 
                elif  sheet_name == 'OverheadAPM':
                    compare_overheadapm(ws_previous, ws_current)
                elif sheet_name == 'ServiceEndpointsAPM':
                    compare_serviceendpointsapm(ws_previous, ws_current)
                elif sheet_name == 'ErrorConfigurationAPM':
                    compare_errorconfigurationapm(ws_previous, ws_current)
                elif  sheet_name == 'HealthRulesAndAlertingAPM':
                    compare_healthrulesandalertingapm(ws_previous, ws_current) 
                elif  sheet_name == 'DataCollectorsAPM':
                    compare_datacollectorsapm(ws_previous, ws_current)
                elif sheet_name == 'DashboardsAPM':
                    compare_dashboardsapm(ws_previous, ws_current)
                elif sheet_name == 'OverallAssessmentAPM':
                    compare_overallassessmentapm(ws_previous, ws_current)
                elif sheet_name == 'Summary':
                    continue
                else:
                    logging.warning(f"No comparison function defined for sheet: {sheet_name}")

        wb_current.save(output_file_path)
        logging.info(f"Comparison results saved to: {output_file_path}")

    except Exception as e:
        logging.error(f"Error in compare_files_other_sheets: {e}", exc_info=True)
        raise

# ==============================
# BRUM SHEET DISPATCHER (INLINE)
# ==============================
def compare_files_other_sheets_brum(previous_file_path, current_file_path, output_file_path):
    try:
        wb_previous = load_workbook(previous_file_path)
        wb_current = load_workbook(current_file_path)

        for sheet_name in wb_current.sheetnames:
            if sheet_name in wb_previous.sheetnames:
                ws_previous = wb_previous[sheet_name]
                ws_current = wb_current[sheet_name]
                logging.debug(f"[BRUM] Processing sheet: {sheet_name}")

                if sheet_name == 'Analysis':
                    compare_analysis_brum(ws_previous, ws_current)
                elif sheet_name == 'NetworkRequestsBRUM':
                    compare_networkrequestsbrum(ws_previous, ws_current)
                elif sheet_name == 'HealthRulesAndAlertingBRUM':
                    compare_healthrulesandalertingbrum(ws_previous, ws_current)
                elif sheet_name == 'OverallAssessmentBRUM':
                    compare_overallassessmentbrum(ws_previous, ws_current)
                elif sheet_name == 'Summary':
                    continue
                else:
                    logging.warning(f"[BRUM] No comparison defined for sheet: {sheet_name}")
            else:
                logging.warning(f"[BRUM] Sheet '{sheet_name}' missing in previous workbook.")
        wb_current.save(output_file_path)
        logging.info(f"[BRUM] Comparison results saved to: {output_file_path}")
    except Exception as e:
        logging.error(f"[BRUM] Error in compare_files_other_sheets_brum: {e}", exc_info=True)
        raise

# ==============================
# MRUM SHEET DISPATCHER (INLINE)
# ==============================
def compare_files_other_sheets_mrum(previous_file_path, current_file_path, output_file_path):
    try:
        wb_previous = load_workbook(previous_file_path)
        wb_current = load_workbook(current_file_path)

        for sheet_name in wb_current.sheetnames:
            if sheet_name in wb_previous.sheetnames:
                ws_previous = wb_previous[sheet_name]
                ws_current = wb_current[sheet_name]
                logging.debug(f"[MRUM] Processing sheet: {sheet_name}")

                if sheet_name == 'Analysis':
                    compare_analysis_mrum(ws_previous, ws_current)
                elif sheet_name == 'NetworkRequestsMRUM':
                    compare_networkrequestsmrum(ws_previous, ws_current)
                elif sheet_name == 'HealthRulesAndAlertingMRUM':
                    compare_healthrulesandalertingmrum(ws_previous, ws_current)
                elif sheet_name == 'OverallAssessmentMRUM':
                    compare_overallassessmentmrum(ws_previous, ws_current)
                elif sheet_name == 'Summary':
                    continue
                else:
                    logging.warning(f"[MRUM] No comparison defined for sheet: {sheet_name}")
            else:
                logging.warning(f"[MRUM] Sheet '{sheet_name}' missing in previous workbook.")
        wb_current.save(output_file_path)
        logging.info(f"[MRUM] Comparison results saved to: {output_file_path}")
    except Exception as e:
        logging.error(f"[MRUM] Error in compare_files_other_sheets_mrum: {e}", exc_info=True)
        raise

def compare_appagentsapm(ws_previous, ws_current):
    try:
        # Define column names and their specific comparison logic
        columns = {
            'metricLimitNotHit': None,
            'percentAgentsLessThan1YearOld': None,
            'percentAgentsLessThan2YearsOld': None,
            'percentAgentsReportingData': None,
            'percentAgentsRunningSameVersion': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                   # logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value; retain original formatting or clear output
                        continue

                    if column == 'metricLimitNotHit':
                        # Handle boolean logic for metricLimitNotHit (True/False)
                        prev_value_str = str(previous_value).strip().upper()  # Ensure it's in upper case and string
                        curr_value_str = str(current_value).strip().upper()  # Ensure it's in upper case and string

                        logging.info(f"Comparing {column}: Previous={prev_value_str}, Current={curr_value_str}")

                        # Compare "TRUE" vs "FALSE"
                        if prev_value_str == "FALSE" and curr_value_str == "TRUE":
                            cell_output.fill = green_fill  # Green for improvement (False → True)
                            cell_output.value = f"{previous_value} → {current_value} (Improved)"
                        elif prev_value_str == "TRUE" and curr_value_str == "FALSE":
                            cell_output.fill = red_fill  # Red for decline (True → False)
                            cell_output.value = f"{previous_value} → {current_value} (Declined)"
                        else:
                            cell_output.fill = red_fill
                            cell_output.value = f"{previous_value} → {current_value} (Changed)"

                    # Handle numeric logic for other columns
                    else:
                        try:
                            prev_value_num = float(previous_value)
                            curr_value_num = float(current_value)
                            formatted_prev_value = f"{prev_value_num:.2f}"
                            formatted_curr_value = f"{curr_value_num:.2f}"

                            if curr_value_num > prev_value_num:
                                cell_output.fill = green_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Improved)"
                            else:
                                cell_output.fill = red_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Declined)"
                        except ValueError:
                            logging.error(f"Non-numeric value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_appagentsapm: {e}", exc_info=True)
        raise
 

def compare_machineagentsapm(ws_previous, ws_current):
    try:
        # Define column names and their specific comparison logic
        columns = {
            'percentAgentsLessThan1YearOld': None,
            'percentAgentsLessThan2YearsOld': None,
            'percentAgentsReportingData': None,
            'percentAgentsRunningSameVersion': None,
            'percentAgentsInstalledAlongsideAppAgents': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                   # logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value
                        continue

                    # Handle numeric logic for percentage columns
                    try:
                        prev_value_num = float(previous_value)
                        curr_value_num = float(current_value)
                        formatted_prev_value = f"{prev_value_num:.2f}"
                        formatted_curr_value = f"{curr_value_num:.2f}"

                        if curr_value_num > prev_value_num:
                            cell_output.fill = green_fill
                            cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Improved)"
                        else:
                            cell_output.fill = red_fill
                            cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Declined)"
                    except ValueError:
                        logging.error(f"Non-numeric value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_machineagentsapm: {e}", exc_info=True)
        raise

def compare_datacollectorsapm(ws_previous, ws_current):
    try:
        # Define column names and their specific comparison logic
        columns = {
            'numberOfDataCollectorFieldsConfigured': None,
            'numberOfDataCollectorFieldsCollectedInSnapshots': None,
            'numberOfDataCollectorFieldsCollectedInAnalytics': None,
            'biqEnabled': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                    # logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value; retain original formatting or clear output
                        cell_output.value = previous_value  # Ensure the value is set to the previous value
                        continue  # Skip any further formatting or changes

                    if column == 'biqEnabled':
                        # Handle boolean logic for biqEnabled
                        prev_value_str = str(previous_value).strip().upper()  # Ensure it's in upper case and string
                        curr_value_str = str(current_value).strip().upper()  # Ensure it's in upper case and string

                        # Log the comparison for debugging
                        logging.info(f"Comparing {column}: Previous={prev_value_str}, Current={curr_value_str}")

                        # Compare "TRUE" vs "FALSE"
                        if prev_value_str == "FALSE" and curr_value_str == "TRUE":
                            cell_output.fill = green_fill  # Green for improvement (False → True)
                            cell_output.value = f"{previous_value} → {current_value} (Improved)"
                        elif prev_value_str == "TRUE" and curr_value_str == "FALSE":
                            cell_output.fill = red_fill  # Red for decline (True → False)
                            cell_output.value = f"{previous_value} → {current_value} (Declined)"
                        else:
                            # For other cases, we just mark it as changed
                            cell_output.fill = red_fill
                            cell_output.value = f"{previous_value} → {current_value} (Changed)"

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_datacollectorsapm: {e}", exc_info=True)
        raise


# Function to compare 'Backends' sheet 
def compare_backendsapm(ws_previous, ws_current):
    try:
        # Define column names and their specific comparison logic
        columns = {
            'percentBackendsWithLoad': None,
            'backendLimitNotHit': None,  # Column to compare
            'numberOfCustomBackendRules': None
        }

        # Retrieve column indices for the columns to be compared
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices (application and controller columns)
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                    # logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value; retain original formatting or clear output
                        cell_output.value = previous_value  # Ensure the value is set to the previous value
                        continue  # Skip any further formatting or changes

                    if column == 'backendLimitNotHit':
                        # Handle boolean logic for backendLimitNotHit
                        prev_value = str(previous_value).strip().upper()  # Convert to string and handle case insensitivity
                        curr_value = str(current_value).strip().upper()  # Same for the current value

                        # Log the comparison for debugging
                        logging.info(f"Comparing backendLimitNotHit: Previous={prev_value}, Current={curr_value}")

                        # Compare "FALSE" vs "TRUE"
                        if prev_value == "FALSE" and curr_value == "TRUE":
                            cell_output.fill = green_fill  # Correct color for FALSE → TRUE
                            cell_output.value = "FALSE → TRUE"
                        elif prev_value == "TRUE" and curr_value == "FALSE":
                            cell_output.fill = red_fill  # Correct color for TRUE → FALSE
                            cell_output.value = "TRUE → FALSE"
                        else:
                            cell_output.fill = red_fill  # Default for unexpected values
                            cell_output.value = f"{previous_value} → {current_value} (Changed)"
                    else:
                        # Handle numeric logic for other columns
                        try:
                            prev_value_num = float(previous_value)
                            curr_value_num = float(current_value)
                            formatted_prev_value = f"{prev_value_num:.2f}"
                            formatted_curr_value = f"{curr_value_num:.2f}"

                            if curr_value_num > prev_value_num:
                                cell_output.fill = green_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Increased)"
                            else:
                                cell_output.fill = red_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Decreased)"
                        except ValueError:
                            logging.error(f"Non-numeric value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_backendsapm: {e}", exc_info=True)
        raise

def compare_overheadapm(ws_previous, ws_current):
    try:
        # Define column names and their specific comparison logic
        columns = {
            'developerModeNotEnabledForAnyBT': None,
            'findEntryPointsNotEnabled': None,
            'aggressiveSnapshottingNotEnabled': None,
            'developerModeNotEnabledForApplication': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                    # logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value
                        continue

                    # Handle boolean logic for specified columns
                    if column == 'developerModeNotEnabledForAnyBT' or column == 'findEntryPointsNotEnabled' or column == 'aggressiveSnapshottingNotEnabled' or column == 'developerModeNotEnabledForApplication':
                        prev_value_str = str(previous_value).strip().upper()  # Ensure it's in upper case and string
                        curr_value_str = str(current_value).strip().upper()  # Ensure it's in upper case and string

                        # Log the comparison for debugging
                        logging.info(f"Comparing {column}: Previous={prev_value_str}, Current={curr_value_str}")

                        # Compare "TRUE" vs "FALSE"
                        if prev_value_str == "FALSE" and curr_value_str == "TRUE":
                            cell_output.fill = green_fill  # Green for True (Improvement)
                            cell_output.value = f"{previous_value} → {current_value} (Improved)"
                        elif prev_value_str == "TRUE" and curr_value_str == "FALSE":
                            cell_output.fill = red_fill  # Red for False (Decline)
                            cell_output.value = f"{previous_value} → {current_value} (Declined)"
                        else:
                            # For other cases, we just mark it as changed
                            cell_output.fill = red_fill
                            cell_output.value = f"{previous_value} → {current_value} (Changed)"

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_overheadapm: {e}", exc_info=True)
        raise


def compare_healthrulesandalertingapm(ws_previous, ws_current):
    try:
        # Define column names and their specific comparison logic
        columns = {
            'numberOfHealthRuleViolations': None,
            'numberOfDefaultHealthRulesModified': None,
            'numberOfActionsBoundToEnabledPolicies': None,
            'numberOfCustomHealthRules': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                   # logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value
                        continue

                    try:
                        prev_value_num = float(previous_value)
                        curr_value_num = float(current_value)
                        formatted_prev_value = f"{prev_value_num:.2f}"
                        formatted_curr_value = f"{curr_value_num:.2f}"

                        if column == 'numberOfHealthRuleViolations':
                            if curr_value_num > prev_value_num:
                                cell_output.fill = red_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Increased)"
                            else:
                                cell_output.fill = green_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Decreased)"
                        else:
                            # For other columns
                            if curr_value_num > prev_value_num:
                                cell_output.fill = green_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Increased)"
                            else:
                                cell_output.fill = red_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Decreased)"
                    except ValueError:
                        logging.error(f"Non-numeric value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_healthrulesandalertingapm: {e}", exc_info=True)
        raise

def compare_errorconfigurationapm(ws_previous, ws_current):
    try:
        # Define the columns for comparison
        columns = {
            'successPercentageOfWorstTransaction': None,
            'numberOfCustomRules': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                   # logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value
                        continue

                    try:
                        prev_value_num = float(previous_value)
                        curr_value_num = float(current_value)
                        formatted_prev_value = f"{prev_value_num:.2f}"
                        formatted_curr_value = f"{curr_value_num:.2f}"

                        if column == 'successPercentageOfWorstTransaction':
                            if curr_value_num > prev_value_num:
                                cell_output.fill = green_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Increased)"
                            else:
                                cell_output.fill = red_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Decreased)"
                        elif column == 'numberOfCustomRules':
                            if curr_value_num > prev_value_num:
                                cell_output.fill = green_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Increased)"
                            else:
                                cell_output.fill = red_fill
                                cell_output.value = f"{formatted_prev_value} → {formatted_curr_value} (Decreased)"
                    except ValueError:
                        logging.error(f"Non-numeric value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_errorconfigurationapm: {e}", exc_info=True)
        raise

# Function to compare 'Service Endpoints' sheet
def compare_serviceendpointsapm(ws_previous, ws_current):
    try:
        # Define the columns for comparison
        columns = {
            'numberOfCustomServiceEndpointRules': None,
            'serviceEndpointLimitNotHit': None,
            'percentServiceEndpointsWithLoadOrDisabled': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                    logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value
                        continue

                    try:
                        # Handle each column's specific comparison logic
                        if column == 'numberOfCustomServiceEndpointRules':
                            prev_value_num = float(previous_value)
                            curr_value_num = float(current_value)
                            if curr_value_num > prev_value_num:
                                cell_output.fill = green_fill
                                cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Increased)"
                            elif curr_value_num < prev_value_num:
                                cell_output.fill = red_fill
                                cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Decreased)"

                        elif column == 'serviceEndpointLimitNotHit':
                            # Explicitly handle TRUE/FALSE as strings and booleans
                            prev_value_str = str(previous_value).strip().upper()  # Ensure it's in upper case and string
                            curr_value_str = str(current_value).strip().upper()  # Ensure it's in upper case and string

                            logging.info(f"Comparing serviceEndpointLimitNotHit: Previous={prev_value_str}, Current={curr_value_str}")

                            # Compare "TRUE" vs "FALSE"
                            if prev_value_str == "FALSE" and curr_value_str == "TRUE":
                                cell_output.fill = green_fill
                                cell_output.value = "FALSE → TRUE"
                            elif prev_value_str == "TRUE" and curr_value_str == "FALSE":
                                cell_output.fill = red_fill
                                cell_output.value = "TRUE → FALSE"
                            else:
                                # Log if we encounter an unexpected value
                                logging.info(f"Unexpected values for serviceEndpointLimitNotHit: Previous={previous_value}, Current={current_value}")

                        elif column == 'percentServiceEndpointsWithLoadOrDisabled':
                            prev_value_num = float(previous_value)
                            curr_value_num = float(current_value)
                            if curr_value_num > prev_value_num:
                                cell_output.fill = green_fill
                                cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Increased)"
                            elif curr_value_num < prev_value_num:
                                cell_output.fill = red_fill
                                cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Decreased)"

                    except ValueError:
                        logging.error(f"Non-numeric or invalid value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_serviceendpointsapm: {e}", exc_info=True)
        raise


def compare_dashboardsapm(ws_previous, ws_current):
    try:
        # Define the columns for comparison
        columns = {
            'numberOfDashboards': None,
            'percentageOfDashboardsModifiedLast6Months': None,
            'numberOfDashboardsUsingBiQ': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                 #   logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value
                        continue

                    try:
                        # Handle each column's specific comparison logic
                        prev_value_num = float(previous_value)
                        curr_value_num = float(current_value)
                        if curr_value_num > prev_value_num:
                            cell_output.fill = green_fill
                            cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Increased)"
                        elif curr_value_num < prev_value_num:
                            cell_output.fill = red_fill
                            cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Decreased)"
                    except ValueError:
                        logging.error(f"Non-numeric or invalid value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_dashboardsapm: {e}", exc_info=True)
        raise

def compare_overallassessmentapm(ws_previous, ws_current):
    try:
        # Define the columns for comparison
        columns = {
            'percentageTotalPlatinum': None,
            'percentageTotalGoldOrBetter': None,
            'percentageTotalSilverOrBetter': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')

        if app_col_prev is None or app_col_curr is None:
            logging.error("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                 #   logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value
                        continue

                    try:
                        # Handle each column's specific comparison logic
                        prev_value_num = float(previous_value)
                        curr_value_num = float(current_value)
                        if curr_value_num > prev_value_num:
                            cell_output.fill = green_fill
                            cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Increased)"
                        elif curr_value_num < prev_value_num:
                            cell_output.fill = red_fill
                            cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Decreased)"
                    except ValueError:
                        logging.error(f"Non-numeric or invalid value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_overallassessmentapm: {e}", exc_info=True)
        raise

# Function to compare 'Business Transactions' sheet
def compare_businesstransactionsapm(ws_previous, ws_current):
    try:
        # Define the columns for comparison
        columns = {
            'numberOfBTs': None,
            'percentBTsWithLoad': None,
            'btLockdownEnabled': None,
            'numberCustomMatchRules': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                print(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')
        app_col_prev = get_key_column(ws_previous, 'application')
        app_col_curr = get_key_column(ws_current, 'application')

        if app_col_prev is None or app_col_curr is None:
            print("The 'application' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            print("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            app_value = row[app_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (app_value, ctrl_value)
            if app_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    if previous_value == current_value:
                        # No change in value
                        continue

                    try:
                        print(f"Column: {column}, Previous Value: {previous_value}, Current Value: {current_value}")

                        # Handle each column's specific comparison logic
                        if column == 'numberOfBTs':
                            prev_value_num = float(previous_value)
                            curr_value_num = float(current_value)
                            if 201 <= prev_value_num <= 600:
                                if curr_value_num < prev_value_num:
                                    cell_output.fill = green_fill
                                    cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Decreased)"
                                elif curr_value_num > prev_value_num:
                                    cell_output.fill = red_fill
                                    cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Increased)"

                        elif column == 'percentBTsWithLoad':
                            prev_value_num = float(previous_value)
                            curr_value_num = float(current_value)
                            if curr_value_num > prev_value_num:
                                cell_output.fill = green_fill
                                cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Increased)"
                            elif curr_value_num < prev_value_num:
                                cell_output.fill = red_fill
                                cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Decreased)"

                        elif column == 'btLockdownEnabled':
                            # Normalize TRUE/FALSE as strings
                            prev_value_str = str(previous_value).strip().upper()  # Ensure it's in upper case and string
                            curr_value_str = str(current_value).strip().upper()  # Ensure it's in upper case and string

                            # Log the comparison for debugging
                            print(f"Comparing btLockdownEnabled for app {key}: Previous={prev_value_str}, Current={curr_value_str}")

                            # Compare "TRUE" vs "FALSE"
                            if prev_value_str == "FALSE" and curr_value_str == "TRUE":
                                print(f"Update: FALSE → TRUE for app {key}")
                                cell_output.fill = green_fill
                                cell_output.value = "FALSE → TRUE"
                            elif prev_value_str == "TRUE" and curr_value_str == "FALSE":
                                print(f"Update: TRUE → FALSE for app {key}")
                                cell_output.fill = red_fill
                                cell_output.value = "TRUE → FALSE"
                            else:
                                print(f"Unexpected values for btLockdownEnabled: Previous={previous_value}, Current={current_value}")

                        elif column == 'numberCustomMatchRules':
                            prev_value_num = float(previous_value)
                            curr_value_num = float(current_value)
                            if curr_value_num > prev_value_num:
                                cell_output.fill = green_fill
                                cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Increased)"
                            elif curr_value_num < prev_value_num:
                                cell_output.fill = red_fill
                                cell_output.value = f"{prev_value_num:.2f} → {curr_value_num:.2f} (Decreased)"
                    except ValueError:
                        print(f"Non-numeric or invalid value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        print(f"Error in compare_businesstransactionsapm: {e}")
        raise


def compare_analysis(ws_previous, ws_current):
    try:
        # Define the columns for comparison
        columns = {
            'AppAgentsAPM': None,
            'MachineAgentsAPM': None,
            'BusinessTransactionsAPM': None,
            'BackendsAPM': None,
            'OverheadAPM': None,
            'ServiceEndpointsAPM': None,
            'ErrorConfigurationAPM': None,
            'HealthRulesAndAlertingAPM': None,
            'DataCollectorsAPM': None,
            'DashboardsAPM': None,
            'OverallAssessment': None
        }

        # Retrieve column indices
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"The '{column}' column is missing in one of the sheets. Cannot proceed with comparison.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        # Retrieve key column indices
        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')
        name_col_prev = get_key_column(ws_previous, 'name')
        name_col_curr = get_key_column(ws_current, 'name')

        if name_col_prev is None or name_col_curr is None:
            logging.error("The 'name' column is missing in one of the sheets. Cannot proceed with comparison.")
            return

        if ctrl_col_prev is None or ctrl_col_curr is None:
            logging.error("The 'controller' column is missing in one of the sheets. This might affect the comparison.")
            return

        previous_data = {}
        current_data = {}

        # Read previous data
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            name_value = row[name_col_prev - 1].value
            ctrl_value = row[ctrl_col_prev - 1].value
            key = (name_value, ctrl_value)
            if name_value and ctrl_value:
                previous_data[key] = row

        # Read current data
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            name_value = row[name_col_curr - 1].value
            ctrl_value = row[ctrl_col_curr - 1].value
            key = (name_value, ctrl_value)
            if name_value and ctrl_value:
                current_data[key] = row

        # Compare previous data with current data
        for key, previous_row in previous_data.items():
            current_row = current_data.get(key)
            if current_row:
                for column, (col_idx_prev, col_idx_curr) in columns.items():
                    previous_value = previous_row[col_idx_prev - 1].value
                    current_value = current_row[col_idx_curr - 1].value
                    cell_output = ws_current.cell(row=current_row[0].row, column=col_idx_curr)

                    # Log values being compared
                 #   logging.debug(f"Comparing '{column}' for key '{key}': Previous={previous_value}, Current={current_value}")

                    if previous_value == current_value:
                        # No change in value
                        continue

                    try:
                        # Comparison logic based on ranking
                        ranking = {'bronze': 1, 'silver': 2, 'gold': 3, 'platinum': 4}
                        prev_rank = ranking.get(previous_value.lower(), 0)
                        curr_rank = ranking.get(current_value.lower(), 0)

                        if curr_rank > prev_rank:
                            cell_output.fill = green_fill
                            cell_output.value = f"{previous_value} → {current_value} (Upgraded)"
                        elif curr_rank < prev_rank:
                            cell_output.fill = red_fill
                            cell_output.value = f"{previous_value} → {current_value} (Downgraded)"
                    except ValueError:
                        logging.error(f"Invalid ranking value encountered for column '{column}': Previous={previous_value}, Current={current_value}")

        # Add new entries in the current sheet
        for key, current_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(current_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill

    except Exception as e:
        logging.error(f"Error in compare_analysis: {e}", exc_info=True)
        raise

# =========================
# BRUM COMPARISON FUNCTIONS
# =========================

def compare_analysis_brum(ws_previous, ws_current):
    try:
        columns = {
            'NetworkRequestsBRUM': None,
            'HealthRulesAndAlertingBRUM': None,
            'OverallAssessment': None
        }
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"[BRUM] '{column}' missing in Analysis.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')
        name_col_prev = get_key_column(ws_previous, 'name')
        name_col_curr = get_key_column(ws_current, 'name')
        if name_col_prev is None or name_col_curr is None:
            logging.error("[BRUM] 'name' missing in Analysis.")
            return

        previous_data, current_data = {}, {}
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            key = (row[name_col_prev - 1].value, row[ctrl_col_prev - 1].value if ctrl_col_prev else None)
            if key[0]:
                previous_data[key] = row
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            key = (row[name_col_curr - 1].value, row[ctrl_col_curr - 1].value if ctrl_col_curr else None)
            if key[0]:
                current_data[key] = row

        ranking = {'bronze': 1, 'silver': 2, 'gold': 3, 'platinum': 4}
        for key, prev_row in previous_data.items():
            cur_row = current_data.get(key)
            if not cur_row:
                continue
            for column, (p_idx, c_idx) in columns.items():
                prev_val = prev_row[p_idx - 1].value
                cur_val = cur_row[c_idx - 1].value
                if prev_val == cur_val:
                    continue
                cell_out = ws_current.cell(row=cur_row[0].row, column=c_idx)
                prev_rank = ranking.get(str(prev_val).strip().lower(), 0)
                cur_rank = ranking.get(str(cur_val).strip().lower(), 0)
                if cur_rank > prev_rank:
                    cell_out.fill = green_fill
                    cell_out.value = f"{prev_val} → {cur_val} (Upgraded)"
                elif cur_rank < prev_rank:
                    cell_out.fill = red_fill
                    cell_out.value = f"{prev_val} → {cur_val} (Downgraded)"

        # Add new entries
        for key, cur_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(cur_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill
    except Exception as e:
        logging.error(f"[BRUM] Error in compare_analysis_brum: {e}", exc_info=True)
        raise

def compare_networkrequestsbrum(ws_previous, ws_current):
    try:
        columns = {
            'collectingDataPastOneDay': 'bool',
            'networkRequestLimitNotHit': 'bool',
            'numberCustomMatchRules': 'num',
            'hasBtCorrelation': 'bool',
            'hasCustomEventServiceIncludeRule': 'bool'
        }
        idx = {}
        for c in columns:
            p = get_key_column(ws_previous, c); n = get_key_column(ws_current, c)
            if p is None or n is None:
                logging.error(f"[BRUM] Missing '{c}' in NetworkRequestsBRUM.")
                return
            idx[c] = (p, n)

        app_prev = get_key_column(ws_previous, 'application')
        app_curr = get_key_column(ws_current, 'application')
        ctrl_prev = get_key_column(ws_previous, 'controller')
        ctrl_curr = get_key_column(ws_current, 'controller')

        prev_map, curr_map = {}, {}
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            key = (row[app_prev - 1].value, row[ctrl_prev - 1].value if ctrl_prev else None)
            if key[0]: prev_map[key] = row
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            key = (row[app_curr - 1].value, row[ctrl_curr - 1].value if ctrl_curr else None)
            if key[0]: curr_map[key] = row

        for key, prow in prev_map.items():
            crow = curr_map.get(key)
            if not crow: continue
            for col, kind in columns.items():
                p_idx, c_idx = idx[col]
                pv, cv = prow[p_idx - 1].value, crow[c_idx - 1].value
                if pv == cv: continue
                cell = ws_current.cell(row=crow[0].row, column=c_idx)
                if kind == 'bool':
                    p = str(pv).strip().upper(); c = str(cv).strip().upper()
                    if p == 'FALSE' and c == 'TRUE':
                        cell.fill = green_fill; cell.value = f"{pv} → {cv} (Improved)"
                    elif p == 'TRUE' and c == 'FALSE':
                        cell.fill = red_fill; cell.value = f"{pv} → {cv} (Declined)"
                    else:
                        cell.fill = red_fill; cell.value = f"{pv} → {cv} (Changed)"
                else:
                    try:
                        pnum = float(pv); cnum = float(cv)
                        if cnum > pnum:
                            cell.fill = green_fill; cell.value = f"{pnum:.2f} → {cnum:.2f} (Increased)"
                        else:
                            cell.fill = red_fill; cell.value = f"{pnum:.2f} → {cnum:.2f} (Decreased)"
                    except Exception:
                        logging.error(f"[BRUM] Non-numeric '{col}': {pv} vs {cv}")

        # New entries
        for key, crow in curr_map.items():
            if key not in prev_map:
                r = ws_current.max_row + 1
                for i, c in enumerate(crow, 1):
                    nc = ws_current.cell(row=r, column=i, value=c.value); nc.fill = added_fill
    except Exception as e:
        logging.error(f"[BRUM] Error in compare_networkrequestsbrum: {e}", exc_info=True)
        raise

def compare_healthrulesandalertingbrum(ws_previous, ws_current):
    try:
        columns = {
            'numberOfHealthRuleViolations': 'lower_better',
            'numberOfActionsBoundToEnabledPolicies': 'higher_better',
            'numberOfCustomHealthRules': 'higher_better'
        }
        idx = {}
        for c in columns:
            p = get_key_column(ws_previous, c); n = get_key_column(ws_current, c)
            if p is None or n is None:
                logging.error(f"[BRUM] Missing '{c}' in HealthRulesAndAlertingBRUM.")
                return
            idx[c] = (p, n)

        app_prev = get_key_column(ws_previous, 'application')
        app_curr = get_key_column(ws_current, 'application')
        ctrl_prev = get_key_column(ws_previous, 'controller')
        ctrl_curr = get_key_column(ws_current, 'controller')

        prev_map, curr_map = {}, {}
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            key = (row[app_prev - 1].value, row[ctrl_prev - 1].value if ctrl_prev else None)
            if key[0]: prev_map[key] = row
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            key = (row[app_curr - 1].value, row[ctrl_curr - 1].value if ctrl_curr else None)
            if key[0]: curr_map[key] = row

        for key, prow in prev_map.items():
            crow = curr_map.get(key)
            if not crow: continue
            for col, rule in columns.items():
                p_idx, c_idx = idx[col]
                pv, cv = prow[p_idx - 1].value, crow[c_idx - 1].value
                if pv == cv: continue
                cell = ws_current.cell(row=crow[0].row, column=c_idx)
                try:
                    pnum = float(pv); cnum = float(cv)
                    if rule == 'lower_better':
                        if cnum < pnum:
                            cell.fill = green_fill; cell.value = f"{pnum:.2f} → {cnum:.2f} (Improved)"
                        else:
                            cell.fill = red_fill; cell.value = f"{pnum:.2f} → {cnum:.2f} (Declined)"
                    else:
                        if cnum > pnum:
                            cell.fill = green_fill; cell.value = f"{pnum:.2f} → {cnum:.2f} (Increased)"
                        else:
                            cell.fill = red_fill; cell.value = f"{pnum:.2f} → {cnum:.2f} (Decreased)"
                except Exception:
                    logging.error(f"[BRUM] Non-numeric '{col}': {pv} vs {cv}")

        # New entries
        for key, crow in curr_map.items():
            if key not in prev_map:
                r = ws_current.max_row + 1
                for i, c in enumerate(crow, 1):
                    nc = ws_current.cell(row=r, column=i, value=c.value); nc.fill = added_fill
    except Exception as e:
        logging.error(f"[BRUM] Error in compare_healthrulesandalertingbrum: {e}", exc_info=True)
        raise

def compare_overallassessmentbrum(ws_previous, ws_current):
    try:
        columns = {
            'percentageTotalPlatinum': None,
            'percentageTotalGoldOrBetter': None,
            'percentageTotalSilverOrBetter': None
        }
        for c in columns.keys():
            p = get_key_column(ws_previous, c); n = get_key_column(ws_current, c)
            if p is None or n is None:
                logging.error(f"[BRUM] Missing '{c}' in OverallAssessmentBRUM.")
                return
            columns[c] = (p, n)

        ctrl_prev = get_key_column(ws_previous, 'controller')
        ctrl_curr = get_key_column(ws_current, 'controller')
        app_prev = get_key_column(ws_previous, 'application')
        app_curr = get_key_column(ws_current, 'application')

        prev_map, curr_map = {}, {}
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            key = (row[app_prev - 1].value, row[ctrl_prev - 1].value if ctrl_prev else None)
            if key[0]: prev_map[key] = row
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            key = (row[app_curr - 1].value, row[ctrl_curr - 1].value if ctrl_curr else None)
            if key[0]: curr_map[key] = row

        for key, prow in prev_map.items():
            crow = curr_map.get(key)
            if not crow: continue
            for col, (p_idx, c_idx) in columns.items():
                pv, cv = prow[p_idx - 1].value, crow[c_idx - 1].value
                if pv == cv: continue
                cell = ws_current.cell(row=crow[0].row, column=c_idx)
                try:
                    pnum = float(str(pv).replace('%', ''))
                    cnum = float(str(cv).replace('%', ''))
                    if cnum > pnum:
                        cell.fill = green_fill; cell.value = f"{pnum:.2f}% → {cnum:.2f}% (Increased)"
                    else:
                        cell.fill = red_fill; cell.value = f"{pnum:.2f}% → {cnum:.2f}% (Decreased)"
                except Exception:
                    logging.error(f"[BRUM] Non-numeric '{col}': {pv} vs {cv}")

        # New entries
        for key, crow in curr_map.items():
            if key not in prev_map:
                r = ws_current.max_row + 1
                for i, c in enumerate(crow, 1):
                    nc = ws_current.cell(row=r, column=i, value=c.value); nc.fill = added_fill
    except Exception as e:
        logging.error(f"[BRUM] Error in compare_overallassessmentbrum: {e}", exc_info=True)
        raise

# =========================
# MRUM COMPARISON FUNCTIONS
# =========================

def compare_analysis_mrum(ws_previous, ws_current):
    try:
        columns = {
            'NetworkRequestsMRUM': None,
            'HealthRulesAndAlertingMRUM': None,
            'OverallAssessment': None
        }
        for column in columns.keys():
            col_idx_prev = get_key_column(ws_previous, column)
            col_idx_curr = get_key_column(ws_current, column)
            if col_idx_prev is None or col_idx_curr is None:
                logging.error(f"[MRUM] '{column}' missing in Analysis.")
                return
            columns[column] = (col_idx_prev, col_idx_curr)

        ctrl_col_prev = get_key_column(ws_previous, 'controller')
        ctrl_col_curr = get_key_column(ws_current, 'controller')
        name_col_prev = get_key_column(ws_previous, 'name')
        name_col_curr = get_key_column(ws_current, 'name')
        if name_col_prev is None or name_col_curr is None:
            logging.error("[MRUM] 'name' missing in Analysis.")
            return

        previous_data, current_data = {}, {}
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            key = (row[name_col_prev - 1].value, row[ctrl_col_prev - 1].value if ctrl_col_prev else None)
            if key[0]:
                previous_data[key] = row
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            key = (row[name_col_curr - 1].value, row[ctrl_col_curr - 1].value if ctrl_col_curr else None)
            if key[0]:
                current_data[key] = row

        ranking = {'bronze': 1, 'silver': 2, 'gold': 3, 'platinum': 4}
        for key, prev_row in previous_data.items():
            cur_row = current_data.get(key)
            if not cur_row:
                continue
            for column, (p_idx, c_idx) in columns.items():
                prev_val = prev_row[p_idx - 1].value
                cur_val = cur_row[c_idx - 1].value
                if prev_val == cur_val:
                    continue
                cell_out = ws_current.cell(row=cur_row[0].row, column=c_idx)
                prev_rank = ranking.get(str(prev_val).strip().lower(), 0)
                cur_rank = ranking.get(str(cur_val).strip().lower(), 0)
                if cur_rank > prev_rank:
                    cell_out.fill = green_fill
                    cell_out.value = f"{prev_val} → {cur_val} (Upgraded)"
                elif cur_rank < prev_rank:
                    cell_out.fill = red_fill
                    cell_out.value = f"{prev_val} → {cur_val} (Downgraded)"

        # Add new entries
        for key, cur_row in current_data.items():
            if key not in previous_data:
                row_index = ws_current.max_row + 1
                for col_num, cell in enumerate(cur_row, 1):
                    new_cell = ws_current.cell(row=row_index, column=col_num, value=cell.value)
                    new_cell.fill = added_fill
    except Exception as e:
        logging.error(f"[MRUM] Error in compare_analysis_mrum: {e}", exc_info=True)
        raise


def compare_networkrequestsmrum(ws_previous, ws_current):
    try:
        columns = {
            'collectingDataPastOneDay': 'bool',
            'networkRequestLimitNotHit': 'bool',
            'numberCustomMatchRules': 'num',
            'hasBtCorrelation': 'bool',
            'hasCustomEventServiceIncludeRule': 'bool'
        }
        idx = {}
        for c in columns:
            p = get_key_column(ws_previous, c); n = get_key_column(ws_current, c)
            if p is None or n is None:
                logging.error(f"[MRUM] Missing '{c}' in NetworkRequestsMRUM.")
                return
            idx[c] = (p, n)

        app_prev = get_key_column(ws_previous, 'application')
        app_curr = get_key_column(ws_current, 'application')
        ctrl_prev = get_key_column(ws_previous, 'controller')
        ctrl_curr = get_key_column(ws_current, 'controller')

        previous_map, current_map = {}, {}
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            key = (row[app_prev - 1].value, row[ctrl_prev - 1].value if ctrl_prev else None)
            if key[0]:
                previous_map[key] = row
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            key = (row[app_curr - 1].value, row[ctrl_curr - 1].value if ctrl_curr else None)
            if key[0]:
                current_map[key] = row

        for key, prow in previous_map.items():
            crow = current_map.get(key)
            if not crow:
                continue
            for col, kind in columns.items():
                p_idx, c_idx = idx[col]
                pv, cv = prow[p_idx - 1].value, crow[c_idx - 1].value
                if pv == cv:
                    continue
                cell = ws_current.cell(row=crow[0].row, column=c_idx)
                if kind == 'bool':
                    p = str(pv).strip().upper()
                    c = str(cv).strip().upper()
                    if p == 'FALSE' and c == 'TRUE':
                        cell.fill = green_fill; cell.value = f"{pv} → {cv} (Improved)"
                    elif p == 'TRUE' and c == 'FALSE':
                        cell.fill = red_fill; cell.value = f"{pv} → {cv} (Declined)"
                    else:
                        cell.fill = red_fill; cell.value = f"{pv} → {cv} (Changed)"
                else:
                    try:
                        pnum = float(pv); cnum = float(cv)
                        if cnum > pnum:
                            cell.fill = green_fill; cell.value = f"{pnum:.2f} → {cnum:.2f} (Increased)"
                        else:
                            cell.fill = red_fill; cell.value = f"{pnum:.2f} → {cnum:.2f} (Decreased)"
                    except Exception:
                        logging.error(f"[MRUM] Non-numeric '{col}': {pv} vs {cv}")

        # New entries
        for key, crow in current_map.items():
            if key not in previous_map:
                r = ws_current.max_row + 1
                for i, c in enumerate(crow, 1):
                    nc = ws_current.cell(row=r, column=i, value=c.value); nc.fill = added_fill
    except Exception as e:
        logging.error(f"[MRUM] Error in compare_networkrequestsmrum: {e}", exc_info=True)
        raise


def compare_healthrulesandalertingmrum(ws_previous, ws_current):
    try:
        columns = {
            'numberOfHealthRuleViolations': 'lower_better',
            'numberOfActionsBoundToEnabledPolicies': 'higher_better',
            'numberOfCustomHealthRules': 'higher_better'
        }
        idx = {}
        for c in columns:
            p = get_key_column(ws_previous, c); n = get_key_column(ws_current, c)
            if p is None or n is None:
                logging.error(f"[MRUM] Missing '{c}' in HealthRulesAndAlertingMRUM.")
                return
            idx[c] = (p, n)

        app_prev = get_key_column(ws_previous, 'application')
        app_curr = get_key_column(ws_current, 'application')
        ctrl_prev = get_key_column(ws_previous, 'controller')
        ctrl_curr = get_key_column(ws_current, 'controller')

        previous_map, current_map = {}, {}
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            key = (row[app_prev - 1].value, row[ctrl_prev - 1].value if ctrl_prev else None)
            if key[0]:
                previous_map[key] = row
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            key = (row[app_curr - 1].value, row[ctrl_curr - 1].value if ctrl_curr else None)
            if key[0]:
                current_map[key] = row

        for key, prow in previous_map.items():
            crow = current_map.get(key)
            if not crow:
                continue
            for col, rule in columns.items():
                p_idx, c_idx = idx[col]
                pv, cv = prow[p_idx - 1].value, crow[c_idx - 1].value
                if pv == cv:
                    continue
                cell = ws_current.cell(row=crow[0].row, column=c_idx)
                try:
                    pnum = float(pv); cnum = float(cv)
                    if rule == 'lower_better':
                        if cnum < pnum:
                            cell.fill = green_fill; cell.value = f"{pnum:.2f} → {cnum:.2f} (Improved)"
                        else:
                            cell.fill = red_fill; cell.value = f"{pnum:.2f} → {cnum:.2f} (Declined)"
                    else:
                        if cnum > pnum:
                            cell.fill = green_fill; cell.value = f"{pnum:.2f} → {cnum:.2f} (Increased)"
                        else:
                            cell.fill = red_fill; cell.value = f"{pnum:.2f} → {cnum:.2f} (Decreased)"
                except Exception:
                    logging.error(f"[MRUM] Non-numeric '{col}': {pv} vs {cv}")

        # New entries
        for key, crow in current_map.items():
            if key not in previous_map:
                r = ws_current.max_row + 1
                for i, c in enumerate(crow, 1):
                    nc = ws_current.cell(row=r, column=i, value=c.value); nc.fill = added_fill
    except Exception as e:
        logging.error(f"[MRUM] Error in compare_healthrulesandalertingmrum: {e}", exc_info=True)
        raise


def compare_overallassessmentmrum(ws_previous, ws_current):
    try:
        columns = {
            'percentageTotalPlatinum': None,
            'percentageTotalGoldOrBetter': None,
            'percentageTotalSilverOrBetter': None
        }
        for c in columns.keys():
            p = get_key_column(ws_previous, c); n = get_key_column(ws_current, c)
            if p is None or n is None:
                logging.error(f"[MRUM] Missing '{c}' in OverallAssessmentMRUM.")
                return
            columns[c] = (p, n)

        ctrl_prev = get_key_column(ws_previous, 'controller')
        ctrl_curr = get_key_column(ws_current, 'controller')
        app_prev = get_key_column(ws_previous, 'application')
        app_curr = get_key_column(ws_current, 'application')

        previous_map, current_map = {}, {}
        for row in ws_previous.iter_rows(min_row=2, values_only=False):
            key = (row[app_prev - 1].value, row[ctrl_prev - 1].value if ctrl_prev else None)
            if key[0]:
                previous_map[key] = row
        for row in ws_current.iter_rows(min_row=2, values_only=False):
            key = (row[app_curr - 1].value, row[ctrl_curr - 1].value if ctrl_curr else None)
            if key[0]:
                current_map[key] = row

        for key, prow in previous_map.items():
            crow = current_map.get(key)
            if not crow:
                continue
            for col, (p_idx, c_idx) in columns.items():
                pv, cv = prow[p_idx - 1].value, crow[c_idx - 1].value
                if pv == cv:
                    continue
                cell = ws_current.cell(row=crow[0].row, column=c_idx)
                try:
                    pnum = float(str(pv).replace('%', ''))
                    cnum = float(str(cv).replace('%', ''))
                    if cnum > pnum:
                        cell.fill = green_fill; cell.value = f"{pnum:.2f}% → {cnum:.2f}% (Increased)"
                    else:
                        cell.fill = red_fill; cell.value = f"{pnum:.2f}% → {cnum:.2f}% (Decreased)"
                except Exception:
                    logging.error(f"[MRUM] Non-numeric '{col}': {pv} vs {cv}")

        # New entries
        for key, crow in current_map.items():
            if key not in previous_map:
                r = ws_current.max_row + 1
                for i, c in enumerate(crow, 1):
                    nc = ws_current.cell(row=r, column=i, value=c.value); nc.fill = added_fill
    except Exception as e:
        logging.error(f"[MRUM] Error in compare_overallassessmentmrum: {e}", exc_info=True)
        raise

    def _last_json_path_for_domain(domain):
        return app.config.get(f"LAST_JSON_{domain.upper()}")

"""     @app.route("/api/apps", methods=["GET"])
    def api_apps():
        domain = (request.args.get("domain") or "").upper()
        json_path = _last_json_path_for_domain(domain)
        if not domain or not json_path or not os.path.isfile(json_path):
            return jsonify({"error": "No JSON found for domain."}), 400
        with open(json_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        return jsonify({"domain": domain, "apps": data.get("apps", {}).get("names", [])})

    @app.route("/api/insights", methods=["GET"])
    def api_insights():
        domain = (request.args.get("domain") or "").upper()
        app_name = (request.args.get("app") or "").strip()
        json_path = _last_json_path_for_domain(domain)
        if not domain or not app_name or not json_path or not os.path.isfile(json_path):
            return jsonify({"error": "Missing or invalid parameters."}), 400
        with open(json_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        apps_index = data.get("appsIndex", {})
        if app_name not in apps_index:
            return jsonify({"error": "Application not found."}), 404
        payload = {
            "domain": domain,
            "app": app_name,
            "areas": apps_index[app_name].get("areas", []),
            "detail": apps_index[app_name].get("detail", {})
        }
        return jsonify(payload) """

# Auto-open the default web browser after a short delay.
import threading, time, webbrowser
port = 5000
url = f"http://127.0.0.1:{port}"

def _open_browser():
    time.sleep(1.0)  # small delay so the server is ready
    try:
        webbrowser.open(url)
    except Exception:
        pass

threading.Thread(target=_open_browser, daemon=True).start()

# Run the Flask app once (outside of any loop).
app.run(debug=True, host='127.0.0.1', port=port, use_reloader=False)
